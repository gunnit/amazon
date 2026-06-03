"""Export endpoints for Excel, PowerPoint, CSV, and PDF packages."""
from typing import Dict, List, Optional
from datetime import date, timedelta
from uuid import UUID
import io
import os
import logging
from fastapi import APIRouter, HTTPException, Query, status
from fastapi.responses import StreamingResponse
from pydantic import BaseModel, Field
from sqlalchemy import select

from app.api.deps import CurrentUser, CurrentOrganization, DbSession
from app.models.amazon_account import AmazonAccount
from app.models.forecast_export_job import ForecastExportJob
from app.models.sales_data import SalesData
from app.schemas.exports import ForecastExportCreate, ForecastExportJobResponse
from app.services.data_extraction import DAILY_TOTAL_ASIN
from app.services.export_service import ExportService
from app.services.forecast_export_service import (
    ForecastExportService,
    build_forecast_excel_filename,
    build_forecast_workbook_bytes,
)
from app.models.advertising import AdvertisingCampaign, AdvertisingMetrics
from app.services.strategic_recommendations_export import (
    build_recommendations_workbook_bytes,
)
from app.services.strategic_recommendations_service import (
    VALID_CATEGORIES as REC_VALID_CATEGORIES,
    VALID_STATUSES as REC_VALID_STATUSES,
    StrategicRecommendationsService,
)
from workers.tasks.forecast_exports import process_forecast_export

logger = logging.getLogger(__name__)

router = APIRouter()


def _forecast_export_job_to_response(job: ForecastExportJob) -> ForecastExportJobResponse:
    """Convert a forecast export job model to its response schema."""
    return ForecastExportJobResponse(
        id=str(job.id),
        forecast_id=str(job.forecast_id),
        status=job.status,
        progress_step=job.progress_step,
        progress_pct=job.progress_pct or 0,
        error_message=job.error_message,
        include_insights=job.include_insights,
        template=job.template,
        language=job.language,
        download_ready=job.status == "completed" and bool(job.artifact_data),
        created_at=job.created_at.isoformat() if job.created_at else "",
        completed_at=job.completed_at.isoformat() if job.completed_at else None,
    )


# Brand palette sampled from the Niuexa/Libera "Inthezon — Strategia Amazon" deck.
# The deck's dominant field is a near-black deep navy; the lockup carries a
# red→teal→blue→gold accent strip under "INSPIRE MORE.".
_PPT_NAVY = (7, 8, 12)            # cover / band background (deep navy-black)
_PPT_INK = (31, 41, 55)          # titles on light slides
_PPT_GREY = (90, 90, 90)
_PPT_WHITE = (255, 255, 255)
_PPT_BODY = (40, 40, 40)
_PPT_ACCENT = (0, 94, 132)       # brand blue from the accent strip (#005E84)
_PPT_ACCENT_STRIP = [
    (207, 41, 43),   # red
    (18, 176, 177),  # teal
    (0, 94, 132),    # blue
    (206, 157, 62),  # gold
]
_PPT_FONT = "Nunito"

_PPT_LOGO_PATH = os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(__file__))), "assets", "inthezon_logo.png")


class _PowerPointBuilder:
    """Builds a Niuexa/Libera-branded, Italian/European-formatted sales deck."""

    def __init__(self, is_it: bool) -> None:
        from pptx import Presentation

        self.is_it = is_it
        self.prs = Presentation()
        self.width = self.prs.slide_width
        self.height = self.prs.slide_height

    def _blank(self):
        return self.prs.slides.add_slide(self.prs.slide_layouts[6])

    def _rgb(self, rgb):
        from pptx.dml.color import RGBColor

        return RGBColor(*rgb)

    def _set_font(self, font):
        font.name = _PPT_FONT

    def _accent_strip(self, slide, left, top, width, height):
        """Draw the brand's multicolour accent strip (red→teal→blue→gold)."""
        from pptx.util import Inches

        seg = width / len(_PPT_ACCENT_STRIP)
        for i, col in enumerate(_PPT_ACCENT_STRIP):
            bar = slide.shapes.add_shape(
                1, Inches(left + seg * i), Inches(top), Inches(seg), Inches(height)
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = self._rgb(col)
            bar.line.fill.background()
            bar.shadow.inherit = False

    def _textbox(self, slide, left, top, width, height, text, size, *, bold=False,
                 color=_PPT_INK, align=None):
        from pptx.util import Inches, Pt

        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        para = box.text_frame.paragraphs[0]
        para.text = text
        para.font.size = Pt(size)
        para.font.bold = bold
        para.font.color.rgb = self._rgb(color)
        self._set_font(para.font)
        if align is not None:
            para.alignment = align
        box.text_frame.word_wrap = True
        return box

    def cover(self, title, period, scope_label, scope_caption, period_caption,
              source_caption, source_value, footer):
        from pptx.util import Inches, Pt

        slide = self._blank()
        # Full-bleed dark band so the white-on-navy logo lockup sits seamlessly.
        band = slide.shapes.add_shape(1, Inches(0), Inches(0), self.width, Inches(3.7))
        band.fill.solid()
        band.fill.fore_color.rgb = self._rgb(_PPT_NAVY)
        band.line.fill.background()
        band.shadow.inherit = False

        slide_w_in = self.width / 914400
        if os.path.exists(_PPT_LOGO_PATH):
            logo_w = 3.6
            logo_h = logo_w * 282 / 880
            slide.shapes.add_picture(
                _PPT_LOGO_PATH, Inches(0.6), Inches(0.55),
                Inches(logo_w), Inches(logo_h),
            )
        self._textbox(slide, 0.6, 2.15, slide_w_in - 1.2, 1.0, title, 32,
                      bold=True, color=_PPT_WHITE)
        self._textbox(slide, 0.6, 3.05, slide_w_in - 1.2, 0.5, period, 16,
                      color=_PPT_WHITE)

        self._accent_strip(slide, 0.6, 4.0, 3.2, 0.08)

        rows = [
            (scope_caption, scope_label),
            (period_caption, period),
            (source_caption, source_value),
        ]
        top = 4.4
        for caption, value in rows:
            self._textbox(slide, 0.6, top, 2.6, 0.4, caption, 12, bold=True, color=_PPT_GREY)
            self._textbox(slide, 3.2, top, 6.2, 0.4, value, 12, color=_PPT_BODY)
            top += 0.5
        self._textbox(slide, 0.6, 6.9, 9, 0.4, footer, 11, color=_PPT_GREY)

    def _slide_title(self, slide, title, subtitle=None):
        self._textbox(slide, 0.6, 0.45, 9, 0.8, title, 28, bold=True, color=_PPT_INK)
        self._accent_strip(slide, 0.6, 1.18, 1.4, 0.06)
        if subtitle is not None:
            self._textbox(slide, 0.6, 1.3, 9, 0.4, subtitle, 12, color=_PPT_GREY)

    def executive_summary(self, title, lines, scope_label, period):
        from pptx.util import Inches, Pt

        slide = self._blank()
        self._slide_title(slide, title, f"{period}  ·  {scope_label}")
        box = slide.shapes.add_textbox(
            Inches(0.6), Inches(2.0), Inches(8.8), Inches(4.0)
        )
        tf = box.text_frame
        tf.word_wrap = True
        for idx, line in enumerate(lines):
            para = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            para.text = f"•  {line}"
            para.font.size = Pt(16)
            para.font.color.rgb = self._rgb(_PPT_BODY)
            self._set_font(para.font)
            para.space_after = Pt(14)

    def kpi_slide(self, title, subtitle, kpis):
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN

        slide = self._blank()
        self._slide_title(slide, title, subtitle)

        box_w, box_h, gap, cols = Inches(2.9), Inches(1.4), Inches(0.1), 3
        for index, (label, value) in enumerate(kpis):
            row, col = divmod(index, cols)
            left = Inches(0.5) + (box_w + gap) * col
            top = Inches(2.0) + (box_h + gap) * row
            shape = slide.shapes.add_shape(1, left, top, box_w, box_h)
            shape.fill.solid()
            shape.fill.fore_color.rgb = self._rgb(_PPT_NAVY)
            shape.line.fill.background()
            shape.shadow.inherit = False
            tf = shape.text_frame
            tf.word_wrap = True
            p_label = tf.paragraphs[0]
            p_label.text = label
            p_label.font.size = Pt(12)
            p_label.font.color.rgb = self._rgb(_PPT_WHITE)
            self._set_font(p_label.font)
            p_label.alignment = PP_ALIGN.CENTER
            p_value = tf.add_paragraph()
            p_value.text = value
            p_value.font.size = Pt(22)
            p_value.font.bold = True
            p_value.font.color.rgb = self._rgb(_PPT_WHITE)
            self._set_font(p_value.font)
            p_value.alignment = PP_ALIGN.CENTER

    def trend_slide(self, title, trend_rows, value_caption):
        from pptx.util import Inches
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE

        slide = self._blank()
        self._slide_title(slide, title)

        chart_data = CategoryChartData()
        chart_data.categories = [
            r["report_date"].isoformat() if hasattr(r["report_date"], "isoformat")
            else str(r["report_date"])
            for r in trend_rows
        ]
        chart_data.add_series(value_caption, [float(r["revenue"]) for r in trend_rows])
        gframe = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(0.5), Inches(1.7), Inches(9), Inches(4.8),
            chart_data,
        )
        self._brand_chart(gframe.chart)

    def _brand_chart(self, chart):
        from pptx.dml.color import RGBColor

        chart.has_legend = False
        plot = chart.plots[0]
        for series in plot.series:
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = RGBColor(*_PPT_ACCENT)
        try:
            chart.font.name = _PPT_FONT
        except Exception:
            pass

    def top_products_slide(self, title, note, headers, rows):
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN

        slide = self._blank()
        self._slide_title(slide, title, note)

        n_rows = len(rows) + 1
        n_cols = len(headers)
        table_shape = slide.shapes.add_table(
            n_rows, n_cols, Inches(0.5), Inches(1.9), Inches(9), Inches(0.4 * n_rows)
        )
        table = table_shape.table
        table.columns[0].width = Inches(1.6)
        table.columns[1].width = Inches(4.4)
        table.columns[2].width = Inches(1.3)
        table.columns[3].width = Inches(1.7)

        for c, header in enumerate(headers):
            cell = table.cell(0, c)
            cell.text = header
            para = cell.text_frame.paragraphs[0]
            para.font.bold = True
            para.font.size = Pt(12)
            para.font.color.rgb = self._rgb(_PPT_WHITE)
            self._set_font(para.font)
            cell.fill.solid()
            cell.fill.fore_color.rgb = self._rgb(_PPT_NAVY)

        for r, row in enumerate(rows, start=1):
            for c, value in enumerate(row):
                cell = table.cell(r, c)
                cell.text = str(value)
                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(11)
                self._set_font(para.font)
                if c >= 2:
                    para.alignment = PP_ALIGN.RIGHT

    def agency_slide(self, title, description, services_caption, services,
                     contacts_caption, contacts, footer):
        from pptx.util import Inches, Pt

        slide = self._blank()
        band = slide.shapes.add_shape(1, Inches(0), Inches(0), self.width, Inches(1.5))
        band.fill.solid()
        band.fill.fore_color.rgb = self._rgb(_PPT_NAVY)
        band.line.fill.background()
        band.shadow.inherit = False

        slide_w_in = self.width / 914400
        if os.path.exists(_PPT_LOGO_PATH):
            logo_w = 2.8
            logo_h = logo_w * 282 / 880
            slide.shapes.add_picture(
                _PPT_LOGO_PATH, Inches(0.6), Inches(0.35),
                Inches(logo_w), Inches(logo_h),
            )
        self._accent_strip(slide, 0.6, 1.75, 2.0, 0.07)

        self._textbox(slide, 0.6, 2.0, slide_w_in - 1.2, 0.6, title, 24,
                      bold=True, color=_PPT_INK)

        box = slide.shapes.add_textbox(
            Inches(0.6), Inches(2.7), Inches(slide_w_in - 1.2), Inches(1.6)
        )
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = description
        p.font.size = Pt(13)
        p.font.color.rgb = self._rgb(_PPT_BODY)
        self._set_font(p.font)

        self._textbox(slide, 0.6, 4.5, 4.4, 0.4, services_caption, 12,
                      bold=True, color=_PPT_ACCENT)
        sbox = slide.shapes.add_textbox(Inches(0.6), Inches(4.9), Inches(4.4), Inches(2.0))
        stf = sbox.text_frame
        stf.word_wrap = True
        for idx, svc in enumerate(services):
            para = stf.paragraphs[0] if idx == 0 else stf.add_paragraph()
            para.text = f"•  {svc}"
            para.font.size = Pt(12)
            para.font.color.rgb = self._rgb(_PPT_BODY)
            self._set_font(para.font)
            para.space_after = Pt(6)

        self._textbox(slide, 5.4, 4.5, 4.0, 0.4, contacts_caption, 12,
                      bold=True, color=_PPT_ACCENT)
        cbox = slide.shapes.add_textbox(Inches(5.4), Inches(4.9), Inches(4.0), Inches(2.0))
        ctf = cbox.text_frame
        ctf.word_wrap = True
        for idx, line in enumerate(contacts):
            para = ctf.paragraphs[0] if idx == 0 else ctf.add_paragraph()
            para.text = line
            para.font.size = Pt(12)
            para.font.color.rgb = self._rgb(_PPT_BODY)
            self._set_font(para.font)
            para.space_after = Pt(6)

        self._textbox(slide, 0.6, 6.95, slide_w_in - 1.2, 0.4, footer, 10,
                      color=_PPT_GREY)

    def to_bytes(self) -> io.BytesIO:
        output = io.BytesIO()
        self.prs.save(output)
        output.seek(0)
        return output


@router.post("/csv")
async def export_to_csv_package(
    current_user: CurrentUser,
    organization: CurrentOrganization,
    db: DbSession,
    report_type: str = Query(..., regex="^(sales|inventory|advertising)$"),
    start_date: date = Query(default_factory=lambda: date.today() - timedelta(days=30)),
    end_date: date = Query(default_factory=lambda: date.today()),
    account_ids: Optional[List[UUID]] = Query(default=None),
    group_by: str = Query(default="day", regex="^(day|week|month)$"),
    low_stock_only: bool = Query(default=False),
    language: str = Query(default="en", regex="^(en|it)$"),
    include_comparison: bool = Query(default=True),
):
    """Generate a professional CSV ZIP package for the selected report."""
    export_service = ExportService(db)
    package_bytes, filename = await export_service.generate_csv_package(
        organization_id=organization.id,
        report_type=report_type,
        start_date=start_date,
        end_date=end_date,
        account_ids=account_ids,
        group_by=group_by,
        low_stock_only=low_stock_only,
        language=language,
        include_comparison=include_comparison,
    )

    return StreamingResponse(
        io.BytesIO(package_bytes),
        media_type="application/zip",
        headers={"Content-Disposition": f"attachment; filename={filename}"},
    )


@router.post("/bundle")
async def export_bundle(
    current_user: CurrentUser,
    organization: CurrentOrganization,
    db: DbSession,
    report_types: List[str] = Query(...),
    start_date: date = Query(default_factory=lambda: date.today() - timedelta(days=30)),
    end_date: date = Query(default_factory=lambda: date.today()),
    account_ids: Optional[List[UUID]] = Query(default=None),
    group_by: str = Query(default="day", regex="^(day|week|month)$"),
    low_stock_only: bool = Query(default=False),
    language: str = Query(default="en", regex="^(en|it)$"),
    include_comparison: bool = Query(default=True),
):
    """Generate a bundled CSV ZIP package with multiple report types."""
    valid_types = {"sales", "inventory", "advertising"}
    for rt in report_types:
        if rt not in valid_types:
            raise HTTPException(
                status_code=status.HTTP_422_UNPROCESSABLE_ENTITY,
                detail=f"Invalid report_type: {rt}. Must be one of {valid_types}",
            )

    export_service = ExportService(db)
    package_bytes, filename = await export_service.generate_bundle_package(
        organization_id=organization.id,
        report_types=report_types,
        start_date=start_date,
        end_date=end_date,
        account_ids=account_ids,
        group_by=group_by,
        low_stock_only=low_stock_only,
        language=language,
        include_comparison=include_comparison,
    )

    return StreamingResponse(
        io.BytesIO(package_bytes),
        media_type="application/zip",
        headers={"Content-Disposition": f"attachment; filename={filename}"},
    )


@router.post("/excel-bundle")
async def export_excel_bundle(
    current_user: CurrentUser,
    organization: CurrentOrganization,
    db: DbSession,
    report_types: List[str] = Query(...),
    start_date: date = Query(default_factory=lambda: date.today() - timedelta(days=30)),
    end_date: date = Query(default_factory=lambda: date.today()),
    account_ids: Optional[List[UUID]] = Query(default=None),
    group_by: str = Query(default="day", regex="^(day|week|month)$"),
    low_stock_only: bool = Query(default=False),
    language: str = Query(default="en", regex="^(en|it)$"),
    include_comparison: bool = Query(default=True),
    template: str = Query(default="corporate", regex="^(clean|corporate|executive)$"),
):
    """Generate a styled Excel workbook with multiple report types."""
    valid_types = {"sales", "inventory", "advertising"}
    for rt in report_types:
        if rt not in valid_types:
            raise HTTPException(
                status_code=status.HTTP_422_UNPROCESSABLE_ENTITY,
                detail=f"Invalid report_type: {rt}. Must be one of {valid_types}",
            )

    export_service = ExportService(db)
    package_bytes, filename = await export_service.generate_excel_bundle(
        organization_id=organization.id,
        report_types=report_types,
        start_date=start_date,
        end_date=end_date,
        template=template,
        account_ids=account_ids,
        group_by=group_by,
        low_stock_only=low_stock_only,
        language=language,
        include_comparison=include_comparison,
    )

    return StreamingResponse(
        io.BytesIO(package_bytes),
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={"Content-Disposition": f"attachment; filename={filename}"},
    )


@router.post("/excel")
async def export_to_excel(
    current_user: CurrentUser,
    organization: CurrentOrganization,
    db: DbSession,
    start_date: date = Query(default_factory=lambda: date.today() - timedelta(days=30)),
    end_date: date = Query(default_factory=lambda: date.today()),
    account_ids: Optional[List[UUID]] = Query(default=None),
    include_sales: bool = Query(default=True),
    include_advertising: bool = Query(default=True),
):
    """Generate Excel export of data."""
    from datetime import datetime
    from openpyxl import Workbook

    from app.services.excel_templates import TEMPLATES, ExcelTemplateRenderer

    accounts_query = select(AmazonAccount.id).where(
        AmazonAccount.organization_id == organization.id
    )
    if account_ids:
        accounts_query = accounts_query.where(AmazonAccount.id.in_(account_ids))

    sales_rows: List[Dict] = []
    ads_rows: List[Dict] = []

    if include_sales:
        sales_query = (
            select(SalesData)
            .where(
                SalesData.account_id.in_(accounts_query),
                SalesData.asin != DAILY_TOTAL_ASIN,
                SalesData.date >= start_date,
                SalesData.date <= end_date,
            )
            .order_by(SalesData.date.desc())
        )
        result = await db.execute(sales_query)
        for sale in result.scalars().all():
            sales_rows.append({
                "date": sale.date.isoformat(),
                "asin": sale.asin,
                "sku": sale.sku or "",
                "units_ordered": sale.units_ordered,
                "revenue": float(sale.ordered_product_sales),
                "orders": sale.total_order_items,
                "currency": sale.currency,
            })

    if include_advertising:
        campaigns_query = select(AdvertisingCampaign.id).where(
            AdvertisingCampaign.account_id.in_(accounts_query)
        )
        ads_query = (
            select(AdvertisingMetrics, AdvertisingCampaign.campaign_name)
            .join(AdvertisingCampaign)
            .where(
                AdvertisingMetrics.campaign_id.in_(campaigns_query),
                AdvertisingMetrics.date >= start_date,
                AdvertisingMetrics.date <= end_date,
            )
            .order_by(AdvertisingMetrics.date.desc())
        )
        result = await db.execute(ads_query)
        for metric, campaign_name in result.all():
            ads_rows.append({
                "date": metric.date.isoformat(),
                "campaign": campaign_name or "",
                "impressions": metric.impressions,
                "clicks": metric.clicks,
                "cost": float(metric.cost),
                "sales": float(metric.attributed_sales_7d),
                "roas": float(metric.roas or 0),
                "acos": float(metric.acos or 0),
            })

    renderer = ExcelTemplateRenderer(TEMPLATES["corporate"])
    wb = Workbook()
    wb.remove(wb.active)

    # Summary / KPI sheet
    summary_rows: List[Dict] = []
    if include_sales:
        summary_rows.append({"metric": "Units Ordered", "current_value": sum(r["units_ordered"] or 0 for r in sales_rows)})
        summary_rows.append({"metric": "Revenue", "current_value": round(sum(r["revenue"] for r in sales_rows), 2)})
        summary_rows.append({"metric": "Orders", "current_value": sum(r["orders"] or 0 for r in sales_rows)})
    if include_advertising:
        summary_rows.append({"metric": "Ad Impressions", "current_value": sum(r["impressions"] or 0 for r in ads_rows)})
        summary_rows.append({"metric": "Ad Clicks", "current_value": sum(r["clicks"] or 0 for r in ads_rows)})
        summary_rows.append({"metric": "Ad Spend", "current_value": round(sum(r["cost"] for r in ads_rows), 2)})
        summary_rows.append({"metric": "Ad Sales", "current_value": round(sum(r["sales"] for r in ads_rows), 2)})

    ws_summary = wb.create_sheet(title="Summary")
    next_row = renderer.write_title_banner(
        ws_summary,
        title="Inthezon Report",
        subtitle=f"{start_date.isoformat()} — {end_date.isoformat()}",
        metadata_rows=[
            ("Organization", organization.name),
            ("Period", f"{start_date} to {end_date}"),
            ("Generated at", datetime.utcnow().strftime("%Y-%m-%d %H:%M")),
        ],
    )
    renderer.write_summary_sheet(
        ws_summary,
        summary_rows,
        ["metric", "current_value"],
        ["Metric", "Value"],
        start_row=next_row,
    )

    if include_sales:
        ws_sales = wb.create_sheet(title="Sales Data")
        renderer.write_data_sheet(
            ws_sales,
            sales_rows,
            ["date", "asin", "sku", "units_ordered", "revenue", "orders", "currency"],
            ["Date", "ASIN", "SKU", "Units Ordered", "Revenue", "Orders", "Currency"],
        )

    if include_advertising:
        ws_ads = wb.create_sheet(title="Advertising")
        renderer.write_data_sheet(
            ws_ads,
            ads_rows,
            ["date", "campaign", "impressions", "clicks", "cost", "sales", "roas", "acos"],
            ["Date", "Campaign", "Impressions", "Clicks", "Cost", "Sales", "ROAS", "ACoS"],
        )

    output = io.BytesIO()
    wb.save(output)
    output.seek(0)

    filename = f"inthezon_report_{start_date}_{end_date}.xlsx"

    return StreamingResponse(
        output,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={"Content-Disposition": f"attachment; filename={filename}"}
    )


@router.post("/powerpoint")
async def export_to_powerpoint(
    current_user: CurrentUser,
    organization: CurrentOrganization,
    db: DbSession,
    start_date: date = Query(default_factory=lambda: date.today() - timedelta(days=30)),
    end_date: date = Query(default_factory=lambda: date.today()),
    account_ids: Optional[List[UUID]] = Query(default=None),
    group_by: str = Query(default="month", regex="^(day|week|month)$"),
    template: str = Query(default="default"),
    language: str = Query(default="it", regex="^(en|it)$"),
):
    """Generate an Italian, European-formatted PowerPoint deck."""
    from app.services.export_service import format_money_eu, format_int_eu

    export_service = ExportService(db)
    scoped_accounts = await export_service._get_accounts(organization.id, account_ids)
    scoped_account_ids = [a.id for a in scoped_accounts]

    summary = await export_service._sales_summary(scoped_account_ids, start_date, end_date)
    trend_rows = await export_service._sales_trend_rows(
        scoped_account_ids, start_date, end_date, group_by
    )
    product_rows = await export_service._sales_product_rows(
        scoped_account_ids, start_date, end_date
    )

    currency = "EUR"
    for row in trend_rows:
        if row.get("currency"):
            currency = row["currency"]
            break

    is_it = language == "it"

    def L(en: str, it: str) -> str:
        return it if is_it else en

    def money(value: float) -> str:
        return format_money_eu(value, currency)

    scope_label = (
        ", ".join(a.account_name for a in scoped_accounts)
        if scoped_accounts
        else L("All accounts", "Tutti gli account")
    )
    if len(scope_label) > 80:
        scope_label = scope_label[:77] + "…"

    builder = _PowerPointBuilder(is_it=is_it)
    builder.cover(
        title=L("Amazon Performance Report", "Report Prestazioni Amazon"),
        period=f"{start_date.isoformat()} — {end_date.isoformat()}",
        scope_label=scope_label,
        scope_caption=L("Account scope", "Ambito account"),
        period_caption=L("Period", "Periodo"),
        source_caption=L("Data source", "Fonte dati"),
        source_value=L(
            "Amazon SP-API — confirmed daily sales",
            "Amazon SP-API — vendite giornaliere confermate",
        ),
        footer=L("Generated by Inthezon", "Generato da Inthezon"),
    )

    builder.executive_summary(
        title=L("Executive Summary", "Sintesi Esecutiva"),
        lines=[
            L(
                f"Confirmed revenue of {money(summary['revenue'])} over the selected period.",
                f"Fatturato confermato di {money(summary['revenue'])} nel periodo selezionato.",
            ),
            L(
                f"{format_int_eu(summary['units'])} units sold across {format_int_eu(summary['orders'])} orders.",
                f"{format_int_eu(summary['units'])} unità vendute su {format_int_eu(summary['orders'])} ordini.",
            ),
            L(
                f"Average order value of {money(summary['average_order_value'])} "
                f"and {format_int_eu(summary['active_asins'])} active ASINs.",
                f"Valore medio ordine di {money(summary['average_order_value'])} "
                f"e {format_int_eu(summary['active_asins'])} ASIN attivi.",
            ),
        ],
        scope_label=scope_label,
        period=f"{start_date.isoformat()} — {end_date.isoformat()}",
    )

    builder.kpi_slide(
        title=L("Key Performance Indicators", "Indicatori di Performance"),
        subtitle=f"{L('Period', 'Periodo')}: {start_date.isoformat()} — {end_date.isoformat()}  ·  "
        f"{L('Scope', 'Ambito')}: {scope_label}",
        kpis=[
            (L("Total Revenue", "Fatturato Totale"), money(summary["revenue"])),
            (L("Units Sold", "Unità Vendute"), format_int_eu(summary["units"])),
            (L("Total Orders", "Ordini Totali"), format_int_eu(summary["orders"])),
            (L("Avg. Order Value", "Valore Medio Ordine"), money(summary["average_order_value"])),
            (L("Avg. Selling Price", "Prezzo Medio Vendita"), money(summary["average_selling_price"])),
            (L("Active ASINs", "ASIN Attivi"), format_int_eu(summary["active_asins"])),
        ],
    )

    if trend_rows:
        builder.trend_slide(
            title=L("Revenue Trend", "Andamento Fatturato"),
            trend_rows=trend_rows,
            value_caption=L("Revenue", "Fatturato"),
        )

    if product_rows:
        builder.top_products_slide(
            title=L("Top Products", "Prodotti Principali"),
            note=L(
                "Estimated revenue from Amazon's by-ASIN report.",
                "Fatturato stimato dal report per ASIN di Amazon.",
            ),
            headers=[
                L("ASIN", "ASIN"),
                L("Product", "Prodotto"),
                L("Units", "Unità"),
                L("Revenue", "Fatturato"),
            ],
            rows=[
                (
                    r["asin"],
                    (r["title"] or "")[:42],
                    format_int_eu(r["units"]),
                    money(r["revenue"]),
                )
                for r in product_rows[:10]
            ],
        )

    builder.agency_slide(
        title=L("About the agency", "Chi siamo"),
        description=L(
            "Inthezon is the Amazon consulting practice of Libera Brand Building "
            "Group, an independent Italian communication group. We design "
            "omnichannel strategies that grow brand value on the marketplaces, "
            "combining creativity, technology, trend analysis and media planning.",
            "Inthezon è la practice di consulenza Amazon di Libera Brand Building "
            "Group, gruppo di comunicazione italiano e indipendente. Realizziamo "
            "strategie omnicanale per accrescere il valore del brand sui "
            "marketplace, unendo creatività, tecnologia, trend analysis e media "
            "planning.",
        ),
        services_caption=L("Amazon consulting", "Consulenza Amazon"),
        services=[
            L("Profile & account management", "Gestione profilo e account"),
            L("Media planning & advertising", "Media planning e advertising"),
            L("Catalogue & content optimization", "Ottimizzazione catalogo e contenuti"),
            L("Performance analytics & reporting", "Analisi performance e reporting"),
        ],
        contacts_caption=L("Contacts", "Contatti"),
        contacts=[
            "Libera Brand Building Group",
            "Via Andreis 18 — Torino",
            "Via Rutilia 10/8 — Milano",
            "liberabrandbuilding.group",
        ],
        footer=L("Generated by Inthezon", "Generato da Inthezon"),
    )

    output = builder.to_bytes()
    filename = f"inthezon_presentation_{start_date}_{end_date}_{language}.pptx"

    return StreamingResponse(
        output,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f"attachment; filename={filename}"}
    )


@router.post("/forecast-excel")
async def export_forecast_excel(
    current_user: CurrentUser,
    organization: CurrentOrganization,
    db: DbSession,
    forecast_id: UUID = Query(...),
    template: str = Query(default="corporate", regex="^(clean|corporate|executive)$"),
    language: str = Query(default="en", regex="^(en|it)$"),
):
    """Generate a styled Excel workbook for a single forecast."""
    service = ForecastExportService(db)
    try:
        context = await service._get_forecast_context(
            org_id=organization.id,
            forecast_id=forecast_id,
            template=template,
            language=language,
        )
    except ValueError:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Forecast not found",
        )

    workbook_bytes = build_forecast_workbook_bytes(
        context.forecast,
        context.account_name,
        template=template,
        language=language,
    )
    filename = build_forecast_excel_filename(
        context.forecast,
        context.account_name,
        template=template,
        language=language,
    )

    return StreamingResponse(
        io.BytesIO(workbook_bytes),
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


@router.post("/forecast-csv")
async def export_forecast_csv(
    current_user: CurrentUser,
    organization: CurrentOrganization,
    db: DbSession,
    forecast_id: UUID = Query(...),
):
    """Generate a CSV export for a single forecast (history + predictions)."""
    import csv

    service = ForecastExportService(db)
    try:
        context = await service._get_forecast_context(
            org_id=organization.id,
            forecast_id=forecast_id,
            template="corporate",
            language="en",
        )
    except ValueError:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Forecast not found",
        )

    forecast = context.forecast
    buffer = io.StringIO()
    writer = csv.writer(buffer)
    writer.writerow(["segment", "date", "predicted_value", "lower_bound", "upper_bound", "currency"])

    for point in context.historical:
        writer.writerow([
            "historical",
            point.date.isoformat(),
            round(point.value, 2),
            "",
            "",
            "EUR",
        ])

    for p in (forecast.predictions or []):
        pred_date = p["date"]
        value = p["value"]
        writer.writerow([
            "forecast",
            pred_date if isinstance(pred_date, str) else pred_date.isoformat(),
            round(value, 2),
            round(p.get("lower", value * 0.8), 2),
            round(p.get("upper", value * 1.2), 2),
            "EUR",
        ])

    safe_name = context.account_name.replace(" ", "_")[:30]
    horizon = forecast.forecast_horizon_days or 30
    filename = f"inthezon_forecast_{safe_name}_{forecast.forecast_type or 'sales'}_{horizon}d_{date.today().isoformat()}.csv"

    return StreamingResponse(
        io.BytesIO(buffer.getvalue().encode("utf-8")),
        media_type="text/csv",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


@router.post("/forecast-package", response_model=ForecastExportJobResponse)
async def create_forecast_export_package(
    request: ForecastExportCreate,
    current_user: CurrentUser,
    organization: CurrentOrganization,
    db: DbSession,
):
    """Create an async forecast export package job."""
    service = ForecastExportService(db)
    try:
        job = await service.create_job(
            org_id=organization.id,
            user_id=current_user.id,
            forecast_id=UUID(request.forecast_id),
            template=request.template,
            language=request.language,
            include_insights=request.include_insights,
        )
        await db.commit()
    except ValueError as exc:
        message = str(exc)
        if "not configured" in message:
            raise HTTPException(status_code=status.HTTP_422_UNPROCESSABLE_ENTITY, detail=message)
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail=message)

    try:
        process_forecast_export.delay(str(job.id))
    except Exception:
        logger.exception(
            "Failed to enqueue forecast export %s on Celery; falling back to in-process thread",
            job.id,
        )
        import threading

        from app.services.forecast_export_service import process_forecast_export_job

        thread = threading.Thread(
            target=process_forecast_export_job,
            args=(str(job.id),),
            daemon=True,
        )
        thread.start()

    return _forecast_export_job_to_response(job)


@router.get("/forecast-package/{job_id}", response_model=ForecastExportJobResponse)
async def get_forecast_export_package(
    job_id: UUID,
    current_user: CurrentUser,
    organization: CurrentOrganization,
    db: DbSession,
):
    """Get forecast export package job status."""
    service = ForecastExportService(db)
    job = await service.get_job(job_id, organization.id)
    if not job:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Forecast export job not found",
        )
    return _forecast_export_job_to_response(job)


@router.get("/forecast-package/{job_id}/download")
async def download_forecast_export_package(
    job_id: UUID,
    current_user: CurrentUser,
    organization: CurrentOrganization,
    db: DbSession,
):
    """Download a completed forecast export package."""
    service = ForecastExportService(db)
    job = await service.get_job(job_id, organization.id)
    if not job:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Forecast export job not found",
        )
    if job.status != "completed" or not job.artifact_data:
        raise HTTPException(
            status_code=status.HTTP_422_UNPROCESSABLE_ENTITY,
            detail=f"Forecast export is not ready (status: {job.status})",
        )

    filename = job.artifact_filename or f"forecast_export_{job.id}.zip"
    return StreamingResponse(
        io.BytesIO(job.artifact_data),
        media_type=job.artifact_content_type or "application/zip",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


class MarketResearchPdfRequest(BaseModel):
    """Request to generate a PDF for a market research report."""
    report_id: str
    language: str = Field(default="en", pattern="^(en|it)$")
    chart_images: Optional[Dict[str, str]] = None  # key -> base64 PNG


@router.post("/market-research-pdf")
async def export_market_research_pdf(
    request: MarketResearchPdfRequest,
    current_user: CurrentUser,
    organization: CurrentOrganization,
    db: DbSession,
):
    """Generate a professional PDF for a completed market research report."""
    from app.services.market_research_service import MarketResearchService
    from app.services.pdf_service import MarketResearchPdfBuilder

    service = MarketResearchService(db)
    report = await service.get_report(UUID(request.report_id), organization.id)

    if not report:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Report not found",
        )
    if report.status != "completed":
        raise HTTPException(
            status_code=status.HTTP_422_UNPROCESSABLE_ENTITY,
            detail=f"Report is not completed (status: {report.status})",
        )

    try:
        builder = MarketResearchPdfBuilder(
            report=report,
            chart_images=request.chart_images,
            language=request.language,
        )
        pdf_bytes = builder.build()
    except Exception as e:
        logger.exception("PDF generation failed for report %s", request.report_id)
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"PDF generation failed: {str(e)[:200]}",
        )

    # Build filename. The title may contain non-ASCII characters (e.g. the
    # narrow no-break space U+202F from Amazon listings) which break the
    # latin-1-only HTTP Content-Disposition header, so reduce it to safe ASCII.
    raw_title = (report.title or "report").replace(":", "")
    safe_title = "".join(
        ch if (ch.isascii() and (ch.isalnum() or ch in "-_")) else "_"
        for ch in raw_title
    )[:60].strip("_") or "report"
    filename = f"inthezon_{safe_title}_{date.today()}.pdf"

    return StreamingResponse(
        io.BytesIO(pdf_bytes),
        media_type="application/pdf",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


@router.post("/recommendations-xlsx")
async def export_recommendations_xlsx(
    db: DbSession,
    organization: CurrentOrganization,
    current_user: CurrentUser,
    status_: Optional[str] = Query(default=None, alias="status"),
    category: Optional[str] = Query(default=None),
    account_id: Optional[UUID] = Query(default=None),
    asin: Optional[str] = Query(default=None),
    ids: Optional[List[UUID]] = Query(default=None),
    language: str = Query(default="en", regex="^(en|it)$"),
    limit: int = Query(default=200, ge=1, le=500),
):
    """Export the matching strategic recommendations as a client-facing XLSX."""
    if status_ and status_ not in REC_VALID_STATUSES:
        raise HTTPException(status_code=400, detail=f"Invalid status: {status_}")
    if category and category not in REC_VALID_CATEGORIES:
        raise HTTPException(status_code=400, detail=f"Invalid category: {category}")

    service = StrategicRecommendationsService(db)
    recommendations = await service.list_recommendations(
        organization.id,
        status=status_,
        category=category,
        account_id=account_id,
        limit=limit,
    )

    if ids:
        id_set = {str(rec_id) for rec_id in ids}
        recommendations = [rec for rec in recommendations if str(rec.id) in id_set]

    if asin:
        normalized = asin.strip().upper()

        def _matches_asin(rec) -> bool:
            ctx = rec.context if isinstance(rec.context, dict) else None
            if not ctx:
                return False
            asins = ctx.get("asins")
            if isinstance(asins, list) and normalized in {
                str(value).upper() for value in asins if isinstance(value, str)
            }:
                return True
            filters = ctx.get("generation_filters")
            if isinstance(filters, dict):
                raw = filters.get("asin")
                if isinstance(raw, str) and raw.upper() == normalized:
                    return True
            return False

        recommendations = [rec for rec in recommendations if _matches_asin(rec)]

    if not recommendations:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="No recommendations match the requested filters",
        )

    accounts_result = await db.execute(
        select(AmazonAccount.id, AmazonAccount.account_name).where(
            AmazonAccount.organization_id == organization.id
        )
    )
    account_names = {str(row.id): row.account_name for row in accounts_result.all()}

    workbook_bytes = build_recommendations_workbook_bytes(
        recommendations,
        account_names=account_names,
        language=language,
        scope_account_id=str(account_id) if account_id else None,
        scope_asin=asin.strip().upper() if asin else None,
    )

    today = date.today().isoformat()
    suffix = (
        asin.strip().upper()
        if asin
        else (str(account_id)[:8] if account_id else "all")
    )
    filename = f"inthezon_recommendations_{suffix}_{today}_{language}.xlsx"

    return StreamingResponse(
        io.BytesIO(workbook_bytes),
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


@router.get("/{export_id}/download")
async def download_export(
    export_id: UUID,
    current_user: CurrentUser,
    organization: CurrentOrganization,
    db: DbSession,
):
    """Download a previously generated export by job id."""
    service = ForecastExportService(db)
    job = await service.get_job(export_id, organization.id)
    if not job or not job.artifact_data:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Export not found",
        )

    filename = job.artifact_filename or f"export_{job.id}.zip"
    return StreamingResponse(
        io.BytesIO(job.artifact_data),
        media_type=job.artifact_content_type or "application/octet-stream",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )
