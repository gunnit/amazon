"""Thin python-pptx drawing layer driven by :class:`DeckTheme`.

Blocks call these instead of touching python-pptx directly, so geometry, colour
and type are consistent. All coordinates are inches on a 13.333 x 7.5 slide.
"""
from __future__ import annotations

import io
from typing import Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from app.services.brand_analysis.deck.theme import RGB, DeckTheme

_ALIGN = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
_ANCHOR = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}


class DeckBuilder:
    """Owns the Presentation and exposes themed drawing primitives."""

    def __init__(self, brand: str) -> None:
        self.theme = DeckTheme
        self.brand = brand
        self.prs = Presentation()
        self.prs.slide_width = Inches(DeckTheme.SLIDE_W)
        self.prs.slide_height = Inches(DeckTheme.SLIDE_H)

    # Slide lifecycle ---------------------------------------------------------
    def blank_slide(self):
        return self.prs.slides.add_slide(self.prs.slide_layouts[6])

    def footer_strip(self, slide) -> None:
        """Seven-segment colour strip along the bottom edge (reference chrome)."""
        x = 0.0
        for color, frac in DeckTheme.FOOTER_STRIP:
            w = DeckTheme.SLIDE_W * frac
            self.rect(slide, x, DeckTheme.SLIDE_H - DeckTheme.STRIP_H, w, DeckTheme.STRIP_H, color)
            x += w

    def chrome(self, slide, page: int) -> None:
        """Reference chrome: red bar + brand eyebrow top-left, page number
        bottom-right, colour strip along the bottom edge."""
        self.rect(slide, DeckTheme.MARGIN, 0, 3.47, 0.15, DeckTheme.BRAND_PRIMARY)
        self.text(slide, DeckTheme.MARGIN, 0.21, 5.0, 0.27, self.brand.upper(),
                  size=DeckTheme.TYPE["eyebrow"], bold=True, color=DeckTheme.INK)
        self.text(slide, DeckTheme.SLIDE_W - 1.0, DeckTheme.SLIDE_H - 0.44, 0.65, 0.25,
                  str(page), size=10, color=DeckTheme.MUTED, align="right")
        self.footer_strip(slide)

    def heading(self, slide, title: str, subtitle: str = "") -> None:
        self.text(slide, DeckTheme.MARGIN, 0.50, DeckTheme.content_w(), 0.62, title,
                  size=DeckTheme.TYPE["h1"], bold=True, color=DeckTheme.INK)
        if subtitle:
            self.text(slide, DeckTheme.MARGIN, 1.16, DeckTheme.content_w(), 0.28, subtitle,
                      size=DeckTheme.TYPE["subtitle"], color=DeckTheme.SUBTLE_INK)

    # Drawing primitives ------------------------------------------------------
    @staticmethod
    def _flat(shape):
        """Kill theme-inherited effects. The empty <a:effectLst/> from
        ``shadow.inherit = False`` is not enough: the shape's <p:style> carries
        an effectRef that LibreOffice still renders as a drop shadow."""
        shape.shadow.inherit = False
        sp = shape._element
        style = sp.find(qn("p:style"))
        if style is not None:
            sp.remove(style)
        return shape

    def rect(self, slide, x, y, w, h, fill: RGB, *, line: Optional[RGB] = None, radius: bool = False):
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
        shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
        self._flat(shape)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*fill)
        if line is None:
            shape.line.fill.background()
        else:
            shape.line.color.rgb = RGBColor(*line)
            shape.line.width = Pt(0.75)
        return shape

    def oval(self, slide, x, y, w, h, fill: RGB):
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h))
        self._flat(shape)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*fill)
        shape.line.fill.background()
        return shape

    def text(self, slide, x, y, w, h, value, *, size, bold=False,
             color: RGB = DeckTheme.INK, align="left", anchor="top"):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        frame = box.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.margin_left = frame.margin_right = 0
        frame.margin_top = frame.margin_bottom = 0
        frame.vertical_anchor = _ANCHOR.get(anchor, MSO_ANCHOR.TOP)
        paragraph = frame.paragraphs[0]
        paragraph.alignment = _ALIGN.get(align, PP_ALIGN.LEFT)
        run = paragraph.add_run()
        run.text = str(value or "")
        run.font.name = DeckTheme.FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = RGBColor(*color)
        return box

    def bullets(self, slide, x, y, w, h, items, *, size, color: RGB = DeckTheme.SUBTLE_INK):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        frame = box.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.margin_left = frame.margin_right = 0
        frame.margin_top = frame.margin_bottom = 0
        for idx, item in enumerate(items):
            paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
            paragraph.space_after = Pt(3)
            run = paragraph.add_run()
            run.text = f"•  {item}"
            run.font.name = DeckTheme.FONT
            run.font.size = Pt(size)
            run.font.color.rgb = RGBColor(*color)
        return box

    def kpi(self, slide, x, y, w, h, label: str, value: str, *,
            accent: Optional[RGB] = None, fill: Optional[RGB] = None,
            chip: str = "") -> None:
        """Reference-style KPI tile: solid colour block, big centred value,
        small centred label below. ``accent``/``fill`` set the block colour;
        light fills get ink text, saturated fills get white."""
        block_fill = fill or accent or DeckTheme.SURFACE
        is_light = sum(block_fill) > 540
        text_color = DeckTheme.INK if is_light else DeckTheme.WHITE
        label_color = DeckTheme.SUBTLE_INK if is_light else DeckTheme.WHITE
        line = DeckTheme.HAIRLINE if is_light else None
        self.rect(slide, x, y, w, h, block_fill, line=line)
        value_size = DeckTheme.TYPE["kpi_value"] if h >= 1.35 else DeckTheme.TYPE["kpi_value_small"]
        label_text = f"{label}  ·  {chip}" if chip else label
        value_h = h * 0.62
        self.text(slide, x + 0.1, y + 0.06, w - 0.2, value_h, value,
                  size=value_size, bold=True, color=text_color,
                  align="center", anchor="middle")
        self.text(slide, x + 0.1, y + value_h + 0.1, w - 0.2, h - value_h - 0.16,
                  label_text, size=DeckTheme.TYPE["kpi_label"], color=label_color,
                  align="center")

    def callout(self, slide, x, y, w, h, title: str, body, *,
                accent: Optional[RGB] = None) -> None:
        """A titled body card. ``body`` may be a string or a list of bullets."""
        self.rect(slide, x, y + 0.07, w, h - 0.07, DeckTheme.SURFACE, line=DeckTheme.HAIRLINE)
        self.rect(slide, x, y, w, 0.07, accent or DeckTheme.BRAND_PRIMARY)
        self.text(slide, x + 0.25, y + 0.22, w - 0.5, 0.3, title,
                  size=DeckTheme.TYPE["body"] + 1, bold=True, color=DeckTheme.INK)
        if isinstance(body, (list, tuple)):
            self.bullets(slide, x + 0.25, y + 0.62, w - 0.5, h - 0.76,
                         [b for b in body if str(b).strip()], size=DeckTheme.TYPE["bullet"])
        elif body:
            self.text(slide, x + 0.25, y + 0.62, w - 0.5, h - 0.76, body,
                      size=DeckTheme.TYPE["bullet"], color=DeckTheme.SUBTLE_INK)

    def table(self, slide, x, y, headers, rows, widths, *,
              accent_columns=None, max_rows: int = 12) -> None:
        rows = rows[:max_rows]
        row_h = 0.38
        total_w = sum(widths)
        table_shape = slide.shapes.add_table(
            len(rows) + 1, len(headers), Inches(x), Inches(y),
            Inches(total_w), Inches(row_h * (len(rows) + 1)),
        )
        table = table_shape.table
        table.first_row = False
        table.horz_banding = False
        for idx, width in enumerate(widths):
            table.columns[idx].width = Inches(width)
        for col, header in enumerate(headers):
            cell = table.cell(0, col)
            cell.text = str(header)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(*DeckTheme.TABLE_HEADER)
            self._cell_font(cell, bold=True, size=DeckTheme.TYPE["table_header"],
                            color=DeckTheme.WHITE)
        for r, row in enumerate(rows, start=1):
            band = DeckTheme.WHITE if r % 2 else DeckTheme.SURFACE
            for c, value in enumerate(row):
                cell = table.cell(r, c)
                cell.text = str(value)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(*band)
                color = None
                if accent_columns and c in accent_columns:
                    color = _signed_color(value)
                self._cell_font(cell, size=DeckTheme.TYPE["table_cell"],
                                bold=(c == 0), color=color)

    def picture(self, slide, png: bytes, x, y, w, h) -> None:
        slide.shapes.add_picture(io.BytesIO(png), Inches(x), Inches(y), Inches(w), Inches(h))

    def _cell_font(self, cell, *, bold=False, size=8, color: Optional[RGB] = None) -> None:
        cell.margin_left = Inches(0.07)
        cell.margin_right = Inches(0.05)
        cell.margin_top = Inches(0.03)
        cell.margin_bottom = Inches(0.03)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = DeckTheme.FONT
                run.font.size = Pt(size)
                run.font.bold = bold
                if color is not None:
                    run.font.color.rgb = RGBColor(*color)

    def to_bytes(self) -> bytes:
        output = io.BytesIO()
        self.prs.save(output)
        return output.getvalue()


def _signed_color(value) -> Optional[RGB]:
    text = str(value).strip().replace("–", "-").replace("−", "-")
    if not text or text in {"—", "N/A"}:
        return None
    if text.startswith("+"):
        return DeckTheme.POSITIVE
    if text.startswith("-") and len(text) > 1 and text[1].isdigit():
        return DeckTheme.NEGATIVE
    return None
