"""Export service for generating Excel, PowerPoint, and CSV packages."""
from __future__ import annotations

import csv
import io
import zipfile
from datetime import date, datetime, timedelta
from decimal import Decimal
from typing import Any, Literal, Optional
from uuid import UUID

from sqlalchemy import Date, and_, cast, func, select
from sqlalchemy.ext.asyncio import AsyncSession

from app.models.advertising import AdvertisingCampaign, AdvertisingMetrics
from app.models.amazon_account import AmazonAccount
from app.models.inventory import InventoryData
from app.models.product import Product
from app.models.sales_data import SalesData
from app.services.data_extraction import DAILY_TOTAL_ASIN

Language = Literal["en", "it"]
ReportType = Literal["sales", "inventory", "advertising"]

LOW_STOCK_THRESHOLD = 10
AT_RISK_THRESHOLD = 25

TRANSLATIONS: dict[str, dict[str, str]] = {
    "en": {
        "section": "Section",
        "metric": "Metric",
        "current_value": "Current Value",
        "previous_value": "Previous Value",
        "change_percent": "Change %",
        "unit": "Unit",
        "notes": "Notes",
        "report_date": "Date",
        "snapshot_date": "Snapshot Date",
        "period_start": "Period Start",
        "period_end": "Period End",
        "generated_at": "Generated At",
        "account_name": "Account",
        "report_type": "Report Type",
        "report_scope": "Report Scope",
        "comparison_basis": "Comparison Basis",
        "group_by": "Group By",
        "asin": "ASIN",
        "sku": "SKU",
        "fnsku": "FNSKU",
        "title": "Title",
        "brand": "Brand",
        "category": "Category",
        "units": "Units",
        "revenue": "Revenue",
        "orders": "Orders",
        "currency": "Currency",
        "average_order_value": "Average Order Value",
        "average_selling_price": "Average Selling Price",
        "units_per_order": "Units per Order",
        "revenue_share": "Revenue Share %",
        "total_skus": "Total SKUs",
        "exported_rows": "Exported Rows",
        "total_available": "Total Available Qty",
        "total_inbound": "Total Inbound Qty",
        "reserved_quantity": "Reserved Qty",
        "low_stock_skus": "Low Stock SKUs",
        "fba_fulfillable": "FBA Fulfillable",
        "mfn_fulfillable": "MFN Fulfillable",
        "inbound_working": "Inbound Working",
        "inbound_shipped": "Inbound Shipped",
        "inbound_total": "Inbound Total",
        "total_fba_quantity": "Total FBA Qty",
        "stock_status": "Stock Status",
        "stock_severity": "Severity",
        "campaign_id": "Campaign ID",
        "campaign_name": "Campaign",
        "campaign_type": "Campaign Type",
        "campaign_state": "State",
        "daily_budget": "Daily Budget",
        "impressions": "Impressions",
        "clicks": "Clicks",
        "spend": "Spend",
        "attributed_sales_7d": "Attributed Sales 7d",
        "attributed_units_7d": "Attributed Units 7d",
        "ctr": "CTR",
        "cpc": "CPC",
        "acos": "ACoS",
        "roas": "ROAS",
        "active_asins": "Active ASINs",
        "active_campaigns": "Active Campaigns",
        "metadata": "Metadata",
        "performance": "Performance",
        "catalog": "Catalog",
        "inventory": "Inventory",
        "advertising": "Advertising",
        "sales_report": "Sales Report",
        "inventory_report": "Inventory Report",
        "advertising_report": "Advertising Report",
        "all_accounts": "All selected accounts",
        "low_stock_only_on": "Low-stock filter enabled",
        "low_stock_only_off": "Full inventory snapshot",
        "day": "Day",
        "week": "Week",
        "month": "Month",
        "currency_unit": "currency",
        "units_unit": "units",
        "orders_unit": "orders",
        "percent_unit": "%",
        "count_unit": "count",
        "status_out_of_stock": "Out of Stock",
        "status_low_stock": "Low Stock",
        "status_at_risk": "At Risk",
        "status_healthy": "Healthy",
        "severity_critical": "Critical",
        "severity_high": "High",
        "severity_medium": "Medium",
        "severity_low": "Low",
        "summary_notes_sales": "Compared with the immediately preceding period.",
        "summary_notes_inventory": "Compared with the previous available inventory snapshot.",
        "summary_notes_advertising": "Compared with the immediately preceding period.",
        "low_stock_note": f"Threshold: available quantity below {LOW_STOCK_THRESHOLD}.",
    },
    "it": {
        "section": "Sezione",
        "metric": "Metrica",
        "current_value": "Valore Attuale",
        "previous_value": "Valore Precedente",
        "change_percent": "Variazione %",
        "unit": "Unità",
        "notes": "Note",
        "report_date": "Data",
        "snapshot_date": "Data Snapshot",
        "period_start": "Inizio Periodo",
        "period_end": "Fine Periodo",
        "generated_at": "Generato Il",
        "account_name": "Account",
        "report_type": "Tipo Report",
        "report_scope": "Ambito Report",
        "comparison_basis": "Base Confronto",
        "group_by": "Raggruppamento",
        "asin": "ASIN",
        "sku": "SKU",
        "fnsku": "FNSKU",
        "title": "Titolo",
        "brand": "Brand",
        "category": "Categoria",
        "units": "Unità",
        "revenue": "Fatturato",
        "orders": "Ordini",
        "currency": "Valuta",
        "average_order_value": "Valore Medio Ordine",
        "average_selling_price": "Prezzo Medio Vendita",
        "units_per_order": "Unità per Ordine",
        "revenue_share": "Quota Fatturato %",
        "total_skus": "SKU Totali",
        "exported_rows": "Righe Esportate",
        "total_available": "Quantità Disponibile Totale",
        "total_inbound": "Quantità In Arrivo Totale",
        "reserved_quantity": "Quantità Riservata",
        "low_stock_skus": "SKU a Bassa Scorta",
        "fba_fulfillable": "FBA Disponibile",
        "mfn_fulfillable": "MFN Disponibile",
        "inbound_working": "Inbound In Lavorazione",
        "inbound_shipped": "Inbound Spedito",
        "inbound_total": "Inbound Totale",
        "total_fba_quantity": "Quantità FBA Totale",
        "stock_status": "Stato Stock",
        "stock_severity": "Severità",
        "campaign_id": "ID Campagna",
        "campaign_name": "Campagna",
        "campaign_type": "Tipo Campagna",
        "campaign_state": "Stato",
        "daily_budget": "Budget Giornaliero",
        "impressions": "Impression",
        "clicks": "Clic",
        "spend": "Spesa",
        "attributed_sales_7d": "Vendite Attribuite 7g",
        "attributed_units_7d": "Unità Attribuite 7g",
        "ctr": "CTR",
        "cpc": "CPC",
        "acos": "ACoS",
        "roas": "ROAS",
        "active_asins": "ASIN Attivi",
        "active_campaigns": "Campagne Attive",
        "metadata": "Metadati",
        "performance": "Prestazioni",
        "catalog": "Catalogo",
        "inventory": "Inventario",
        "advertising": "Pubblicità",
        "sales_report": "Report Vendite",
        "inventory_report": "Report Inventario",
        "advertising_report": "Report Pubblicità",
        "all_accounts": "Tutti gli account selezionati",
        "low_stock_only_on": "Filtro bassa scorta attivo",
        "low_stock_only_off": "Snapshot inventario completo",
        "day": "Giorno",
        "week": "Settimana",
        "month": "Mese",
        "currency_unit": "valuta",
        "units_unit": "unità",
        "orders_unit": "ordini",
        "percent_unit": "%",
        "count_unit": "conteggio",
        "status_out_of_stock": "Esaurito",
        "status_low_stock": "Bassa Scorta",
        "status_at_risk": "A Rischio",
        "status_healthy": "Sano",
        "severity_critical": "Critica",
        "severity_high": "Alta",
        "severity_medium": "Media",
        "severity_low": "Bassa",
        "summary_notes_sales": "Confrontato con il periodo immediatamente precedente.",
        "summary_notes_inventory": "Confrontato con lo snapshot inventario precedente disponibile.",
        "summary_notes_advertising": "Confrontato con il periodo immediatamente precedente.",
        "low_stock_note": f"Soglia: quantità disponibile inferiore a {LOW_STOCK_THRESHOLD}.",
    },
}


def _as_float(value: Any) -> float:
    """Normalize Decimal-like values to floats."""
    if value is None:
        return 0.0
    if isinstance(value, Decimal):
        return float(value)
    return float(value)


def _as_int(value: Any) -> int:
    """Normalize nullable numeric values to ints."""
    if value is None:
        return 0
    return int(value)


def _round(value: float, digits: int = 2) -> float:
    """Round floats consistently for CSV output."""
    return round(value, digits)


def _percent_change(current: Optional[float], previous: Optional[float]) -> Optional[float]:
    """Compute percentage change from previous to current."""
    if previous is None:
        return None
    if previous == 0:
        return 100.0 if current and current > 0 else 0.0
    return ((current or 0.0) - previous) / previous * 100


def _previous_period(start_date: date, end_date: date) -> tuple[date, date]:
    """Return the immediately preceding period with matching inclusive length."""
    period_days = (end_date - start_date).days + 1
    prev_end = start_date - timedelta(days=1)
    prev_start = prev_end - timedelta(days=period_days - 1)
    return prev_start, prev_end


def _safe_divide(numerator: float, denominator: float) -> float:
    """Divide safely, returning zero when the denominator is missing."""
    if denominator == 0:
        return 0.0
    return numerator / denominator


class ExportService:
    """Service for exporting data to various formats."""

    def __init__(self, db: AsyncSession):
        self.db = db

    async def generate_excel_report(
        self,
        account_ids: list[UUID],
        start_date: date,
        end_date: date,
        include_sales: bool = True,
        include_inventory: bool = True,
        include_advertising: bool = True,
    ) -> bytes:
        """Generate comprehensive Excel report."""
        from openpyxl import Workbook
        from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
        from openpyxl.utils.dataframe import dataframe_to_rows
        import pandas as pd

        wb = Workbook()

        header_font = Font(bold=True, color="FFFFFF")
        header_fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
        thin_border = Border(
            left=Side(style="thin"),
            right=Side(style="thin"),
            top=Side(style="thin"),
            bottom=Side(style="thin"),
        )

        ws = wb.active
        ws.title = "Summary"
        ws["A1"] = "Inthezon Performance Report"
        ws["A1"].font = Font(bold=True, size=16)
        ws["A3"] = f"Report Period: {start_date} to {end_date}"

        output = io.BytesIO()
        wb.save(output)
        output.seek(0)

        return output.getvalue()

    async def generate_powerpoint_report(
        self,
        account_ids: list[UUID],
        start_date: date,
        end_date: date,
        template: str = "default",
    ) -> bytes:
        """Generate PowerPoint presentation."""
        from pptx import Presentation
        from pptx.util import Inches, Pt

        prs = Presentation()

        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = "Amazon Performance Report"
        subtitle.text = f"{start_date} to {end_date}"

        output = io.BytesIO()
        prs.save(output)
        output.seek(0)

        return output.getvalue()

    async def generate_csv_package(
        self,
        organization_id: UUID,
        report_type: ReportType,
        start_date: date,
        end_date: date,
        account_ids: Optional[list[UUID]] = None,
        group_by: str = "day",
        low_stock_only: bool = False,
        language: Language = "en",
        include_comparison: bool = True,
    ) -> tuple[bytes, str]:
        """Generate a professional CSV package as a ZIP archive."""
        lang: Language = "it" if language == "it" else "en"
        accounts = await self._get_accounts(organization_id, account_ids)

        if report_type == "sales":
            files = await self._build_sales_package(accounts, start_date, end_date, group_by, lang, include_comparison)
        elif report_type == "inventory":
            files = await self._build_inventory_package(accounts, low_stock_only, lang, include_comparison)
        else:
            files = await self._build_advertising_package(accounts, start_date, end_date, lang, include_comparison)

        archive = io.BytesIO()
        with zipfile.ZipFile(archive, "w", compression=zipfile.ZIP_DEFLATED) as zip_file:
            for filename, data in files:
                zip_file.writestr(filename, data)

        archive.seek(0)
        filename = f"inthezon_{report_type}_report_{start_date}_{end_date}_professional_{lang}.zip"
        return archive.getvalue(), filename

    async def generate_bundle_package(
        self,
        organization_id: UUID,
        report_types: list[ReportType],
        start_date: date,
        end_date: date,
        account_ids: Optional[list[UUID]] = None,
        group_by: str = "day",
        low_stock_only: bool = False,
        language: Language = "en",
        include_comparison: bool = True,
    ) -> tuple[bytes, str]:
        """Generate a bundled ZIP containing CSVs for multiple report types."""
        lang: Language = "it" if language == "it" else "en"
        accounts = await self._get_accounts(organization_id, account_ids)

        all_files: list[tuple[str, str]] = []
        for report_type in report_types:
            if report_type == "sales":
                files = await self._build_sales_package(accounts, start_date, end_date, group_by, lang, include_comparison)
            elif report_type == "inventory":
                files = await self._build_inventory_package(accounts, low_stock_only, lang, include_comparison)
            else:
                files = await self._build_advertising_package(accounts, start_date, end_date, lang, include_comparison)
            all_files.extend(files)

        archive = io.BytesIO()
        with zipfile.ZipFile(archive, "w", compression=zipfile.ZIP_DEFLATED) as zip_file:
            for fname, data in all_files:
                zip_file.writestr(fname, data)

        archive.seek(0)
        filename = f"inthezon_export_{start_date}_{end_date}_{lang}.zip"
        return archive.getvalue(), filename

    async def generate_excel_bundle(
        self,
        organization_id: UUID,
        report_types: list[ReportType],
        start_date: date,
        end_date: date,
        template: str = "corporate",
        account_ids: Optional[list[UUID]] = None,
        group_by: str = "day",
        low_stock_only: bool = False,
        language: Language = "en",
        include_comparison: bool = True,
    ) -> tuple[bytes, str]:
        """Generate a styled Excel workbook with multiple report sheets."""
        from openpyxl import Workbook

        from app.services.excel_templates import TEMPLATES, ExcelTemplateRenderer

        lang: Language = "it" if language == "it" else "en"
        style = TEMPLATES.get(template, TEMPLATES["corporate"])
        renderer = ExcelTemplateRenderer(style)
        accounts = await self._get_accounts(organization_id, account_ids)

        wb = Workbook()
        # Remove default sheet — we'll create named ones
        wb.remove(wb.active)

        for report_type in report_types:
            if report_type == "sales":
                collected = await self._collect_sales_data(accounts, start_date, end_date, group_by, lang, include_comparison)
            elif report_type == "inventory":
                collected = await self._collect_inventory_data(accounts, low_stock_only, lang, include_comparison)
            else:
                collected = await self._collect_advertising_data(accounts, start_date, end_date, lang, include_comparison)

            # Summary sheet
            summary_rows = collected["summary_rows"]
            summary_columns = collected["summary_columns"]
            summary_headers = [self._text(lang, c) for c in summary_columns]
            ws_summary = wb.create_sheet(title=self._text(lang, f"{report_type}_report")[:31])
            next_row = renderer.write_title_banner(
                ws_summary,
                title=self._text(lang, f"{report_type}_report"),
                subtitle=f"{start_date.isoformat()} — {end_date.isoformat()}",
                metadata_rows=[
                    (self._text(lang, "account_name"), self._accounts_label(accounts, lang)),
                    (self._text(lang, "generated_at"), datetime.utcnow().strftime("%Y-%m-%d %H:%M")),
                ],
            )
            renderer.write_summary_sheet(ws_summary, summary_rows, summary_columns, summary_headers, start_row=next_row)

            # Data sheets
            for sheet_info in collected["sheets"]:
                ws = wb.create_sheet(title=sheet_info["name"][:31])
                cols = sheet_info["columns"]
                hdrs = [self._text(lang, c) for c in cols]
                renderer.write_data_sheet(ws, sheet_info["rows"], cols, hdrs)

        output = io.BytesIO()
        wb.save(output)
        output.seek(0)
        filename = f"inthezon_export_{start_date}_{end_date}_{lang}_{style.name}.xlsx"
        return output.getvalue(), filename

    # ── Data collectors (shared by CSV and Excel) ─────────────────────

    async def _collect_sales_data(
        self,
        accounts: list[AmazonAccount],
        start_date: date,
        end_date: date,
        group_by: str,
        lang: Language,
        include_comparison: bool,
    ) -> dict[str, Any]:
        """Collect all sales data for export rendering."""
        account_ids = self._account_ids(accounts)
        prev_start, prev_end = _previous_period(start_date, end_date)
        current_summary = await self._sales_summary(account_ids, start_date, end_date)
        previous_summary = await self._sales_summary(account_ids, prev_start, prev_end) if include_comparison else None
        trend_rows = await self._sales_trend_rows(account_ids, start_date, end_date, group_by)
        product_rows = await self._sales_product_rows(account_ids, start_date, end_date)
        category_rows = await self._sales_category_rows(account_ids, start_date, end_date)

        summary_rows = [
            self._metric_row(lang, "metadata", "report_type", self._text(lang, "sales_report")),
            self._metric_row(lang, "metadata", "period_start", start_date),
            self._metric_row(lang, "metadata", "period_end", end_date),
            self._metric_row(lang, "metadata", "generated_at", datetime.utcnow()),
            self._metric_row(lang, "metadata", "account_name", self._accounts_label(accounts, lang)),
            self._metric_row(lang, "metadata", "group_by", self._text(lang, group_by)),
        ]

        if include_comparison:
            summary_rows.append(
                self._metric_row(
                    lang, "metadata", "comparison_basis",
                    f"{prev_start.isoformat()} -> {prev_end.isoformat()}",
                    notes_key="summary_notes_sales",
                )
            )

        summary_rows.extend([
            self._metric_row(lang, "performance", "revenue", current_summary["revenue"],
                             previous_summary["revenue"] if previous_summary else "",
                             _percent_change(current_summary["revenue"], previous_summary["revenue"]) if previous_summary else None, "currency_unit"),
            self._metric_row(lang, "performance", "units", current_summary["units"],
                             previous_summary["units"] if previous_summary else "",
                             _percent_change(float(current_summary["units"]), float(previous_summary["units"])) if previous_summary else None, "units_unit"),
            self._metric_row(lang, "performance", "orders", current_summary["orders"],
                             previous_summary["orders"] if previous_summary else "",
                             _percent_change(float(current_summary["orders"]), float(previous_summary["orders"])) if previous_summary else None, "orders_unit"),
            self._metric_row(lang, "performance", "average_order_value", current_summary["average_order_value"],
                             previous_summary["average_order_value"] if previous_summary else "",
                             _percent_change(current_summary["average_order_value"], previous_summary["average_order_value"]) if previous_summary else None, "currency_unit"),
            self._metric_row(lang, "performance", "average_selling_price", current_summary["average_selling_price"],
                             previous_summary["average_selling_price"] if previous_summary else "",
                             _percent_change(current_summary["average_selling_price"], previous_summary["average_selling_price"]) if previous_summary else None, "currency_unit"),
            self._metric_row(lang, "performance", "units_per_order", current_summary["units_per_order"],
                             previous_summary["units_per_order"] if previous_summary else "",
                             _percent_change(current_summary["units_per_order"], previous_summary["units_per_order"]) if previous_summary else None, "units_unit"),
            self._metric_row(lang, "catalog", "active_asins", current_summary["active_asins"],
                             previous_summary["active_asins"] if previous_summary else "",
                             _percent_change(float(current_summary["active_asins"]), float(previous_summary["active_asins"])) if previous_summary else None, "count_unit"),
        ])

        summary_columns = ["section", "metric", "current_value", "previous_value", "change_percent", "unit", "notes"]
        trend_columns = ["report_date", "revenue", "units", "orders", "average_order_value", "average_selling_price", "units_per_order", "currency"]
        product_columns = ["account_name", "asin", "sku", "title", "brand", "category", "units", "revenue", "orders", "average_selling_price", "revenue_share", "currency"]
        category_columns = ["category", "revenue", "units", "orders", "average_order_value", "average_selling_price", "revenue_share", "currency"]

        return {
            "summary_rows": summary_rows,
            "summary_columns": summary_columns,
            "sheets": [
                {"name": self._text(lang, "report_date") + " Trend", "rows": trend_rows, "columns": trend_columns},
                {"name": "Product Performance", "rows": product_rows, "columns": product_columns},
                {"name": "Category Breakdown", "rows": category_rows, "columns": category_columns},
            ],
        }

    async def _collect_inventory_data(
        self,
        accounts: list[AmazonAccount],
        low_stock_only: bool,
        lang: Language,
        include_comparison: bool,
    ) -> dict[str, Any]:
        """Collect all inventory data for export rendering."""
        account_ids = self._account_ids(accounts)
        current_snapshot = await self._latest_snapshot_date(account_ids)
        previous_snapshot = await self._latest_snapshot_date(account_ids, before_date=current_snapshot) if include_comparison and current_snapshot else None
        current_rows = await self._inventory_snapshot_rows(account_ids, current_snapshot, lang) if current_snapshot else []
        previous_rows = await self._inventory_snapshot_rows(account_ids, previous_snapshot, lang) if previous_snapshot else []

        current_metrics = self._inventory_metrics(current_rows, low_stock_only)
        previous_metrics = self._inventory_metrics(previous_rows, low_stock_only) if previous_rows else None
        low_stock_rows = [row for row in current_rows if row["total_available"] < LOW_STOCK_THRESHOLD]
        snapshot_rows = low_stock_rows if low_stock_only else current_rows

        summary_rows = [
            self._metric_row(lang, "metadata", "report_type", self._text(lang, "inventory_report")),
            self._metric_row(lang, "metadata", "generated_at", datetime.utcnow()),
            self._metric_row(lang, "metadata", "account_name", self._accounts_label(accounts, lang)),
            self._metric_row(lang, "metadata", "snapshot_date", current_snapshot or ""),
            self._metric_row(lang, "metadata", "report_scope", self._text(lang, "low_stock_only_on" if low_stock_only else "low_stock_only_off")),
        ]

        if include_comparison:
            summary_rows.append(
                self._metric_row(lang, "metadata", "comparison_basis", previous_snapshot or "", notes_key="summary_notes_inventory")
            )

        summary_rows.extend([
            self._metric_row(lang, "inventory", "total_skus", current_metrics["total_skus"],
                             previous_metrics["total_skus"] if previous_metrics else "",
                             _percent_change(float(current_metrics["total_skus"]), float(previous_metrics["total_skus"])) if previous_metrics else None, "count_unit"),
            self._metric_row(lang, "inventory", "exported_rows", current_metrics["exported_rows"],
                             previous_metrics["exported_rows"] if previous_metrics else "",
                             _percent_change(float(current_metrics["exported_rows"]), float(previous_metrics["exported_rows"])) if previous_metrics else None, "count_unit",
                             notes_key="low_stock_note" if low_stock_only else ""),
            self._metric_row(lang, "inventory", "total_available", current_metrics["total_available"],
                             previous_metrics["total_available"] if previous_metrics else "",
                             _percent_change(current_metrics["total_available"], previous_metrics["total_available"]) if previous_metrics else None, "units_unit"),
            self._metric_row(lang, "inventory", "total_inbound", current_metrics["total_inbound"],
                             previous_metrics["total_inbound"] if previous_metrics else "",
                             _percent_change(current_metrics["total_inbound"], previous_metrics["total_inbound"]) if previous_metrics else None, "units_unit"),
            self._metric_row(lang, "inventory", "reserved_quantity", current_metrics["reserved_quantity"],
                             previous_metrics["reserved_quantity"] if previous_metrics else "",
                             _percent_change(current_metrics["reserved_quantity"], previous_metrics["reserved_quantity"]) if previous_metrics else None, "units_unit"),
            self._metric_row(lang, "inventory", "low_stock_skus", current_metrics["low_stock_skus"],
                             previous_metrics["low_stock_skus"] if previous_metrics else "",
                             _percent_change(float(current_metrics["low_stock_skus"]), float(previous_metrics["low_stock_skus"])) if previous_metrics else None, "count_unit",
                             notes_key="low_stock_note"),
        ])

        summary_columns = ["section", "metric", "current_value", "previous_value", "change_percent", "unit", "notes"]
        inv_columns = [
            "snapshot_date", "account_name", "asin", "sku", "fnsku", "title", "brand", "category",
            "fba_fulfillable", "mfn_fulfillable", "total_available", "inbound_working", "inbound_shipped",
            "inbound_total", "reserved_quantity", "total_fba_quantity", "stock_status", "stock_severity",
        ]

        return {
            "summary_rows": summary_rows,
            "summary_columns": summary_columns,
            "sheets": [
                {"name": "Inventory Snapshot", "rows": snapshot_rows, "columns": inv_columns},
                {"name": "Low Stock", "rows": low_stock_rows, "columns": inv_columns},
            ],
        }

    async def _collect_advertising_data(
        self,
        accounts: list[AmazonAccount],
        start_date: date,
        end_date: date,
        lang: Language,
        include_comparison: bool,
    ) -> dict[str, Any]:
        """Collect all advertising data for export rendering."""
        account_ids = self._account_ids(accounts)
        prev_start, prev_end = _previous_period(start_date, end_date)
        current_summary = await self._advertising_summary(account_ids, start_date, end_date)
        previous_summary = await self._advertising_summary(account_ids, prev_start, prev_end) if include_comparison else None
        campaign_rows = await self._advertising_daily_rows(account_ids, start_date, end_date)
        rollup_rows = await self._advertising_rollup_rows(account_ids, start_date, end_date)

        summary_rows = [
            self._metric_row(lang, "metadata", "report_type", self._text(lang, "advertising_report")),
            self._metric_row(lang, "metadata", "period_start", start_date),
            self._metric_row(lang, "metadata", "period_end", end_date),
            self._metric_row(lang, "metadata", "generated_at", datetime.utcnow()),
            self._metric_row(lang, "metadata", "account_name", self._accounts_label(accounts, lang)),
        ]

        if include_comparison:
            summary_rows.append(
                self._metric_row(lang, "metadata", "comparison_basis",
                                 f"{prev_start.isoformat()} -> {prev_end.isoformat()}",
                                 notes_key="summary_notes_advertising")
            )

        summary_rows.extend([
            self._metric_row(lang, "advertising", "spend", current_summary["spend"],
                             previous_summary["spend"] if previous_summary else "",
                             _percent_change(current_summary["spend"], previous_summary["spend"]) if previous_summary else None, "currency_unit"),
            self._metric_row(lang, "advertising", "attributed_sales_7d", current_summary["attributed_sales_7d"],
                             previous_summary["attributed_sales_7d"] if previous_summary else "",
                             _percent_change(current_summary["attributed_sales_7d"], previous_summary["attributed_sales_7d"]) if previous_summary else None, "currency_unit"),
            self._metric_row(lang, "advertising", "attributed_units_7d", current_summary["attributed_units_7d"],
                             previous_summary["attributed_units_7d"] if previous_summary else "",
                             _percent_change(float(current_summary["attributed_units_7d"]), float(previous_summary["attributed_units_7d"])) if previous_summary else None, "units_unit"),
            self._metric_row(lang, "advertising", "impressions", current_summary["impressions"],
                             previous_summary["impressions"] if previous_summary else "",
                             _percent_change(float(current_summary["impressions"]), float(previous_summary["impressions"])) if previous_summary else None, "count_unit"),
            self._metric_row(lang, "advertising", "clicks", current_summary["clicks"],
                             previous_summary["clicks"] if previous_summary else "",
                             _percent_change(float(current_summary["clicks"]), float(previous_summary["clicks"])) if previous_summary else None, "count_unit"),
            self._metric_row(lang, "advertising", "ctr", current_summary["ctr"],
                             previous_summary["ctr"] if previous_summary else "",
                             _percent_change(current_summary["ctr"], previous_summary["ctr"]) if previous_summary else None, "percent_unit"),
            self._metric_row(lang, "advertising", "cpc", current_summary["cpc"],
                             previous_summary["cpc"] if previous_summary else "",
                             _percent_change(current_summary["cpc"], previous_summary["cpc"]) if previous_summary else None, "currency_unit"),
            self._metric_row(lang, "advertising", "acos", current_summary["acos"],
                             previous_summary["acos"] if previous_summary else "",
                             _percent_change(current_summary["acos"], previous_summary["acos"]) if previous_summary else None, "percent_unit"),
            self._metric_row(lang, "advertising", "roas", current_summary["roas"],
                             previous_summary["roas"] if previous_summary else "",
                             _percent_change(current_summary["roas"], previous_summary["roas"]) if previous_summary else None),
            self._metric_row(lang, "advertising", "active_campaigns", current_summary["active_campaigns"],
                             previous_summary["active_campaigns"] if previous_summary else "",
                             _percent_change(float(current_summary["active_campaigns"]), float(previous_summary["active_campaigns"])) if previous_summary else None, "count_unit"),
        ])

        summary_columns = ["section", "metric", "current_value", "previous_value", "change_percent", "unit", "notes"]
        daily_columns = [
            "report_date", "account_name", "campaign_id", "campaign_name", "campaign_type",
            "campaign_state", "daily_budget", "impressions", "clicks", "spend",
            "attributed_sales_7d", "attributed_units_7d", "ctr", "cpc", "acos", "roas",
        ]
        rollup_columns = [c for c in daily_columns if c != "report_date"]

        return {
            "summary_rows": summary_rows,
            "summary_columns": summary_columns,
            "sheets": [
                {"name": "Campaign Performance", "rows": campaign_rows, "columns": daily_columns},
                {"name": "Campaign Rollup", "rows": rollup_rows, "columns": rollup_columns},
            ],
        }

    async def _get_accounts(
        self,
        organization_id: UUID,
        account_ids: Optional[list[UUID]],
    ) -> list[AmazonAccount]:
        """Fetch organization accounts constrained by the selected filters."""
        query = select(AmazonAccount).where(AmazonAccount.organization_id == organization_id)
        if account_ids:
            query = query.where(AmazonAccount.id.in_(account_ids))
        query = query.order_by(AmazonAccount.account_name)
        result = await self.db.execute(query)
        return result.scalars().all()

    def _account_ids(self, accounts: list[AmazonAccount]) -> list[UUID]:
        """Extract account ids from model objects."""
        return [account.id for account in accounts]

    def _accounts_label(self, accounts: list[AmazonAccount], lang: Language) -> str:
        """Create a compact account label for metadata rows."""
        if not accounts:
            return self._text(lang, "all_accounts")
        return ", ".join(account.account_name for account in accounts)

    def _text(self, lang: Language, key: str) -> str:
        """Translate a stable export label."""
        return TRANSLATIONS.get(lang, TRANSLATIONS["en"]).get(key, TRANSLATIONS["en"].get(key, key))

    def _csv_bytes(self, rows: list[dict[str, Any]], columns: list[str], lang: Language) -> bytes:
        """Write CSV rows using localized headers and Excel-friendly BOM encoding."""
        output = io.StringIO(newline="")
        localized_headers = [self._text(lang, column) for column in columns]
        writer = csv.DictWriter(output, fieldnames=localized_headers, lineterminator="\n")
        writer.writeheader()

        for row in rows:
            serialized = {
                self._text(lang, column): self._serialize_value(row.get(column))
                for column in columns
            }
            writer.writerow(serialized)

        return output.getvalue().encode("utf-8-sig")

    def _serialize_value(self, value: Any) -> Any:
        """Normalize values for CSV output."""
        if isinstance(value, Decimal):
            return _round(float(value))
        if isinstance(value, float):
            return _round(value, 4 if abs(value) < 10 else 2)
        if isinstance(value, datetime):
            return value.isoformat(timespec="seconds")
        if isinstance(value, date):
            return value.isoformat()
        return value if value is not None else ""

    def _metric_row(
        self,
        lang: Language,
        section_key: str,
        metric_key: str,
        current_value: Any,
        previous_value: Any = "",
        change_percent: Optional[float] = None,
        unit_key: str = "",
        notes_key: str = "",
        notes_value: str = "",
    ) -> dict[str, Any]:
        """Create a summary row with consistent structure."""
        notes = notes_value or (self._text(lang, notes_key) if notes_key else "")
        return {
            "section": self._text(lang, section_key),
            "metric": self._text(lang, metric_key),
            "current_value": current_value,
            "previous_value": previous_value,
            "change_percent": "" if change_percent is None else _round(change_percent),
            "unit": self._text(lang, unit_key) if unit_key else "",
            "notes": notes,
        }

    async def _build_sales_package(
        self,
        accounts: list[AmazonAccount],
        start_date: date,
        end_date: date,
        group_by: str,
        lang: Language,
        include_comparison: bool,
    ) -> list[tuple[str, bytes]]:
        """Build the sales report ZIP contents."""
        account_ids = self._account_ids(accounts)
        prev_start, prev_end = _previous_period(start_date, end_date)
        current_summary = await self._sales_summary(account_ids, start_date, end_date)
        previous_summary = await self._sales_summary(account_ids, prev_start, prev_end) if include_comparison else None
        trend_rows = await self._sales_trend_rows(account_ids, start_date, end_date, group_by)
        product_rows = await self._sales_product_rows(account_ids, start_date, end_date)
        category_rows = await self._sales_category_rows(account_ids, start_date, end_date)

        summary_rows = [
            self._metric_row(lang, "metadata", "report_type", self._text(lang, "sales_report")),
            self._metric_row(lang, "metadata", "period_start", start_date),
            self._metric_row(lang, "metadata", "period_end", end_date),
            self._metric_row(lang, "metadata", "generated_at", datetime.utcnow()),
            self._metric_row(lang, "metadata", "account_name", self._accounts_label(accounts, lang)),
            self._metric_row(lang, "metadata", "group_by", self._text(lang, group_by)),
        ]

        if include_comparison:
            summary_rows.append(
                self._metric_row(
                    lang,
                    "metadata",
                    "comparison_basis",
                    f"{prev_start.isoformat()} -> {prev_end.isoformat()}",
                    notes_key="summary_notes_sales",
                )
            )

        summary_rows.extend(
            [
                self._metric_row(
                    lang,
                    "performance",
                    "revenue",
                    current_summary["revenue"],
                    previous_summary["revenue"] if previous_summary else "",
                    _percent_change(current_summary["revenue"], previous_summary["revenue"]) if previous_summary else None,
                    "currency_unit",
                ),
                self._metric_row(
                    lang,
                    "performance",
                    "units",
                    current_summary["units"],
                    previous_summary["units"] if previous_summary else "",
                    _percent_change(float(current_summary["units"]), float(previous_summary["units"])) if previous_summary else None,
                    "units_unit",
                ),
                self._metric_row(
                    lang,
                    "performance",
                    "orders",
                    current_summary["orders"],
                    previous_summary["orders"] if previous_summary else "",
                    _percent_change(float(current_summary["orders"]), float(previous_summary["orders"])) if previous_summary else None,
                    "orders_unit",
                ),
                self._metric_row(
                    lang,
                    "performance",
                    "average_order_value",
                    current_summary["average_order_value"],
                    previous_summary["average_order_value"] if previous_summary else "",
                    _percent_change(current_summary["average_order_value"], previous_summary["average_order_value"]) if previous_summary else None,
                    "currency_unit",
                ),
                self._metric_row(
                    lang,
                    "performance",
                    "average_selling_price",
                    current_summary["average_selling_price"],
                    previous_summary["average_selling_price"] if previous_summary else "",
                    _percent_change(current_summary["average_selling_price"], previous_summary["average_selling_price"]) if previous_summary else None,
                    "currency_unit",
                ),
                self._metric_row(
                    lang,
                    "performance",
                    "units_per_order",
                    current_summary["units_per_order"],
                    previous_summary["units_per_order"] if previous_summary else "",
                    _percent_change(current_summary["units_per_order"], previous_summary["units_per_order"]) if previous_summary else None,
                    "units_unit",
                ),
                self._metric_row(
                    lang,
                    "catalog",
                    "active_asins",
                    current_summary["active_asins"],
                    previous_summary["active_asins"] if previous_summary else "",
                    _percent_change(float(current_summary["active_asins"]), float(previous_summary["active_asins"])) if previous_summary else None,
                    "count_unit",
                ),
            ]
        )

        return [
            ("00_summary.csv", self._csv_bytes(summary_rows, ["section", "metric", "current_value", "previous_value", "change_percent", "unit", "notes"], lang)),
            ("01_daily_trend.csv", self._csv_bytes(trend_rows, ["report_date", "revenue", "units", "orders", "average_order_value", "average_selling_price", "units_per_order", "currency"], lang)),
            ("02_product_performance.csv", self._csv_bytes(product_rows, ["account_name", "asin", "sku", "title", "brand", "category", "units", "revenue", "orders", "average_selling_price", "revenue_share", "currency"], lang)),
            ("03_category_breakdown.csv", self._csv_bytes(category_rows, ["category", "revenue", "units", "orders", "average_order_value", "average_selling_price", "revenue_share", "currency"], lang)),
        ]

    async def _build_inventory_package(
        self,
        accounts: list[AmazonAccount],
        low_stock_only: bool,
        lang: Language,
        include_comparison: bool,
    ) -> list[tuple[str, bytes]]:
        """Build the inventory report ZIP contents."""
        account_ids = self._account_ids(accounts)
        current_snapshot = await self._latest_snapshot_date(account_ids)
        previous_snapshot = await self._latest_snapshot_date(account_ids, before_date=current_snapshot) if include_comparison and current_snapshot else None
        current_rows = await self._inventory_snapshot_rows(account_ids, current_snapshot, lang) if current_snapshot else []
        previous_rows = await self._inventory_snapshot_rows(account_ids, previous_snapshot, lang) if previous_snapshot else []

        current_metrics = self._inventory_metrics(current_rows, low_stock_only)
        previous_metrics = self._inventory_metrics(previous_rows, low_stock_only) if previous_rows else None
        low_stock_rows = [row for row in current_rows if row["total_available"] < LOW_STOCK_THRESHOLD]
        snapshot_rows = low_stock_rows if low_stock_only else current_rows

        summary_rows = [
            self._metric_row(lang, "metadata", "report_type", self._text(lang, "inventory_report")),
            self._metric_row(lang, "metadata", "generated_at", datetime.utcnow()),
            self._metric_row(lang, "metadata", "account_name", self._accounts_label(accounts, lang)),
            self._metric_row(lang, "metadata", "snapshot_date", current_snapshot or ""),
            self._metric_row(lang, "metadata", "report_scope", self._text(lang, "low_stock_only_on" if low_stock_only else "low_stock_only_off")),
        ]

        if include_comparison:
            summary_rows.append(
                self._metric_row(
                    lang,
                    "metadata",
                    "comparison_basis",
                    previous_snapshot or "",
                    notes_key="summary_notes_inventory",
                )
            )

        summary_rows.extend(
            [
                self._metric_row(
                    lang,
                    "inventory",
                    "total_skus",
                    current_metrics["total_skus"],
                    previous_metrics["total_skus"] if previous_metrics else "",
                    _percent_change(float(current_metrics["total_skus"]), float(previous_metrics["total_skus"])) if previous_metrics else None,
                    "count_unit",
                ),
                self._metric_row(
                    lang,
                    "inventory",
                    "exported_rows",
                    current_metrics["exported_rows"],
                    previous_metrics["exported_rows"] if previous_metrics else "",
                    _percent_change(float(current_metrics["exported_rows"]), float(previous_metrics["exported_rows"])) if previous_metrics else None,
                    "count_unit",
                    notes_key="low_stock_note" if low_stock_only else "",
                ),
                self._metric_row(
                    lang,
                    "inventory",
                    "total_available",
                    current_metrics["total_available"],
                    previous_metrics["total_available"] if previous_metrics else "",
                    _percent_change(current_metrics["total_available"], previous_metrics["total_available"]) if previous_metrics else None,
                    "units_unit",
                ),
                self._metric_row(
                    lang,
                    "inventory",
                    "total_inbound",
                    current_metrics["total_inbound"],
                    previous_metrics["total_inbound"] if previous_metrics else "",
                    _percent_change(current_metrics["total_inbound"], previous_metrics["total_inbound"]) if previous_metrics else None,
                    "units_unit",
                ),
                self._metric_row(
                    lang,
                    "inventory",
                    "reserved_quantity",
                    current_metrics["reserved_quantity"],
                    previous_metrics["reserved_quantity"] if previous_metrics else "",
                    _percent_change(current_metrics["reserved_quantity"], previous_metrics["reserved_quantity"]) if previous_metrics else None,
                    "units_unit",
                ),
                self._metric_row(
                    lang,
                    "inventory",
                    "low_stock_skus",
                    current_metrics["low_stock_skus"],
                    previous_metrics["low_stock_skus"] if previous_metrics else "",
                    _percent_change(float(current_metrics["low_stock_skus"]), float(previous_metrics["low_stock_skus"])) if previous_metrics else None,
                    "count_unit",
                    notes_key="low_stock_note",
                ),
            ]
        )

        columns = [
            "snapshot_date",
            "account_name",
            "asin",
            "sku",
            "fnsku",
            "title",
            "brand",
            "category",
            "fba_fulfillable",
            "mfn_fulfillable",
            "total_available",
            "inbound_working",
            "inbound_shipped",
            "inbound_total",
            "reserved_quantity",
            "total_fba_quantity",
            "stock_status",
            "stock_severity",
        ]

        return [
            ("00_summary.csv", self._csv_bytes(summary_rows, ["section", "metric", "current_value", "previous_value", "change_percent", "unit", "notes"], lang)),
            ("01_inventory_snapshot.csv", self._csv_bytes(snapshot_rows, columns, lang)),
            ("02_low_stock.csv", self._csv_bytes(low_stock_rows, columns, lang)),
        ]

    async def _build_advertising_package(
        self,
        accounts: list[AmazonAccount],
        start_date: date,
        end_date: date,
        lang: Language,
        include_comparison: bool,
    ) -> list[tuple[str, bytes]]:
        """Build the advertising report ZIP contents."""
        account_ids = self._account_ids(accounts)
        prev_start, prev_end = _previous_period(start_date, end_date)
        current_summary = await self._advertising_summary(account_ids, start_date, end_date)
        previous_summary = await self._advertising_summary(account_ids, prev_start, prev_end) if include_comparison else None
        campaign_rows = await self._advertising_daily_rows(account_ids, start_date, end_date)
        rollup_rows = await self._advertising_rollup_rows(account_ids, start_date, end_date)

        summary_rows = [
            self._metric_row(lang, "metadata", "report_type", self._text(lang, "advertising_report")),
            self._metric_row(lang, "metadata", "period_start", start_date),
            self._metric_row(lang, "metadata", "period_end", end_date),
            self._metric_row(lang, "metadata", "generated_at", datetime.utcnow()),
            self._metric_row(lang, "metadata", "account_name", self._accounts_label(accounts, lang)),
        ]

        if include_comparison:
            summary_rows.append(
                self._metric_row(
                    lang,
                    "metadata",
                    "comparison_basis",
                    f"{prev_start.isoformat()} -> {prev_end.isoformat()}",
                    notes_key="summary_notes_advertising",
                )
            )

        summary_rows.extend(
            [
                self._metric_row(
                    lang,
                    "advertising",
                    "spend",
                    current_summary["spend"],
                    previous_summary["spend"] if previous_summary else "",
                    _percent_change(current_summary["spend"], previous_summary["spend"]) if previous_summary else None,
                    "currency_unit",
                ),
                self._metric_row(
                    lang,
                    "advertising",
                    "attributed_sales_7d",
                    current_summary["attributed_sales_7d"],
                    previous_summary["attributed_sales_7d"] if previous_summary else "",
                    _percent_change(current_summary["attributed_sales_7d"], previous_summary["attributed_sales_7d"]) if previous_summary else None,
                    "currency_unit",
                ),
                self._metric_row(
                    lang,
                    "advertising",
                    "attributed_units_7d",
                    current_summary["attributed_units_7d"],
                    previous_summary["attributed_units_7d"] if previous_summary else "",
                    _percent_change(float(current_summary["attributed_units_7d"]), float(previous_summary["attributed_units_7d"])) if previous_summary else None,
                    "units_unit",
                ),
                self._metric_row(
                    lang,
                    "advertising",
                    "impressions",
                    current_summary["impressions"],
                    previous_summary["impressions"] if previous_summary else "",
                    _percent_change(float(current_summary["impressions"]), float(previous_summary["impressions"])) if previous_summary else None,
                    "count_unit",
                ),
                self._metric_row(
                    lang,
                    "advertising",
                    "clicks",
                    current_summary["clicks"],
                    previous_summary["clicks"] if previous_summary else "",
                    _percent_change(float(current_summary["clicks"]), float(previous_summary["clicks"])) if previous_summary else None,
                    "count_unit",
                ),
                self._metric_row(
                    lang,
                    "advertising",
                    "ctr",
                    current_summary["ctr"],
                    previous_summary["ctr"] if previous_summary else "",
                    _percent_change(current_summary["ctr"], previous_summary["ctr"]) if previous_summary else None,
                    "percent_unit",
                ),
                self._metric_row(
                    lang,
                    "advertising",
                    "cpc",
                    current_summary["cpc"],
                    previous_summary["cpc"] if previous_summary else "",
                    _percent_change(current_summary["cpc"], previous_summary["cpc"]) if previous_summary else None,
                    "currency_unit",
                ),
                self._metric_row(
                    lang,
                    "advertising",
                    "acos",
                    current_summary["acos"],
                    previous_summary["acos"] if previous_summary else "",
                    _percent_change(current_summary["acos"], previous_summary["acos"]) if previous_summary else None,
                    "percent_unit",
                ),
                self._metric_row(
                    lang,
                    "advertising",
                    "roas",
                    current_summary["roas"],
                    previous_summary["roas"] if previous_summary else "",
                    _percent_change(current_summary["roas"], previous_summary["roas"]) if previous_summary else None,
                ),
                self._metric_row(
                    lang,
                    "advertising",
                    "active_campaigns",
                    current_summary["active_campaigns"],
                    previous_summary["active_campaigns"] if previous_summary else "",
                    _percent_change(float(current_summary["active_campaigns"]), float(previous_summary["active_campaigns"])) if previous_summary else None,
                    "count_unit",
                ),
            ]
        )

        daily_columns = [
            "report_date",
            "account_name",
            "campaign_id",
            "campaign_name",
            "campaign_type",
            "campaign_state",
            "daily_budget",
            "impressions",
            "clicks",
            "spend",
            "attributed_sales_7d",
            "attributed_units_7d",
            "ctr",
            "cpc",
            "acos",
            "roas",
        ]
        rollup_columns = [column for column in daily_columns if column != "report_date"]

        return [
            ("00_summary.csv", self._csv_bytes(summary_rows, ["section", "metric", "current_value", "previous_value", "change_percent", "unit", "notes"], lang)),
            ("01_campaign_performance.csv", self._csv_bytes(campaign_rows, daily_columns, lang)),
            ("02_campaign_rollup.csv", self._csv_bytes(rollup_rows, rollup_columns, lang)),
        ]

    async def _sales_summary(self, account_ids: list[UUID], start_date: date, end_date: date) -> dict[str, Any]:
        """Aggregate summary sales metrics from daily total rows."""
        if not account_ids:
            return {
                "revenue": 0.0,
                "units": 0,
                "orders": 0,
                "average_order_value": 0.0,
                "average_selling_price": 0.0,
                "units_per_order": 0.0,
                "active_asins": 0,
            }

        totals_query = (
            select(
                func.sum(SalesData.ordered_product_sales).label("revenue"),
                func.sum(SalesData.units_ordered).label("units"),
                func.sum(SalesData.total_order_items).label("orders"),
            )
            .where(
                SalesData.account_id.in_(account_ids),
                SalesData.asin == DAILY_TOTAL_ASIN,
                SalesData.date >= start_date,
                SalesData.date <= end_date,
            )
        )
        totals_row = (await self.db.execute(totals_query)).one()

        active_asins_query = (
            select(func.count(func.distinct(SalesData.asin)))
            .where(
                SalesData.account_id.in_(account_ids),
                SalesData.asin != DAILY_TOTAL_ASIN,
                SalesData.date >= start_date,
                SalesData.date <= end_date,
            )
        )
        active_asins = (await self.db.execute(active_asins_query)).scalar() or 0

        revenue = _round(_as_float(totals_row.revenue))
        units = _as_int(totals_row.units)
        orders = _as_int(totals_row.orders)

        return {
            "revenue": revenue,
            "units": units,
            "orders": orders,
            "average_order_value": _round(_safe_divide(revenue, orders)),
            "average_selling_price": _round(_safe_divide(revenue, units)),
            "units_per_order": _round(_safe_divide(units, orders)),
            "active_asins": int(active_asins),
        }

    async def _sales_trend_rows(
        self,
        account_ids: list[UUID],
        start_date: date,
        end_date: date,
        group_by: str,
    ) -> list[dict[str, Any]]:
        """Build the daily, weekly, or monthly sales trend table."""
        if not account_ids:
            return []

        if group_by == "week":
            period_expr = cast(func.date_trunc("week", SalesData.date), Date)
        elif group_by == "month":
            period_expr = cast(func.date_trunc("month", SalesData.date), Date)
        else:
            period_expr = SalesData.date

        query = (
            select(
                period_expr.label("report_date"),
                func.sum(SalesData.ordered_product_sales).label("revenue"),
                func.sum(SalesData.units_ordered).label("units"),
                func.sum(SalesData.total_order_items).label("orders"),
                func.max(SalesData.currency).label("currency"),
            )
            .where(
                SalesData.account_id.in_(account_ids),
                SalesData.asin == DAILY_TOTAL_ASIN,
                SalesData.date >= start_date,
                SalesData.date <= end_date,
            )
            .group_by(period_expr)
            .order_by(period_expr)
        )

        rows = (await self.db.execute(query)).all()
        output: list[dict[str, Any]] = []
        for row in rows:
            revenue = _as_float(row.revenue)
            units = _as_int(row.units)
            orders = _as_int(row.orders)
            output.append(
                {
                    "report_date": row.report_date,
                    "revenue": _round(revenue),
                    "units": units,
                    "orders": orders,
                    "average_order_value": _round(_safe_divide(revenue, orders)),
                    "average_selling_price": _round(_safe_divide(revenue, units)),
                    "units_per_order": _round(_safe_divide(units, orders)),
                    "currency": row.currency or "EUR",
                }
            )
        return output

    async def _sales_product_rows(
        self,
        account_ids: list[UUID],
        start_date: date,
        end_date: date,
    ) -> list[dict[str, Any]]:
        """Build the product performance table for sales."""
        if not account_ids:
            return []

        sku_expr = func.coalesce(SalesData.sku, Product.sku)
        query = (
            select(
                AmazonAccount.account_name.label("account_name"),
                SalesData.asin.label("asin"),
                sku_expr.label("sku"),
                Product.title.label("title"),
                Product.brand.label("brand"),
                func.coalesce(Product.category, "Uncategorized").label("category"),
                func.sum(SalesData.units_ordered).label("units"),
                func.sum(SalesData.ordered_product_sales).label("revenue"),
                func.sum(SalesData.total_order_items).label("orders"),
                func.max(SalesData.currency).label("currency"),
            )
            .select_from(SalesData)
            .join(AmazonAccount, AmazonAccount.id == SalesData.account_id)
            .outerjoin(
                Product,
                and_(
                    Product.account_id == SalesData.account_id,
                    Product.asin == SalesData.asin,
                ),
            )
            .where(
                SalesData.account_id.in_(account_ids),
                SalesData.asin != DAILY_TOTAL_ASIN,
                SalesData.date >= start_date,
                SalesData.date <= end_date,
            )
            .group_by(
                AmazonAccount.account_name,
                SalesData.asin,
                sku_expr,
                Product.title,
                Product.brand,
                Product.category,
            )
            .order_by(func.sum(SalesData.ordered_product_sales).desc(), SalesData.asin)
        )

        rows = (await self.db.execute(query)).all()
        total_revenue = sum(_as_float(row.revenue) for row in rows)
        output: list[dict[str, Any]] = []
        for row in rows:
            revenue = _as_float(row.revenue)
            units = _as_int(row.units)
            output.append(
                {
                    "account_name": row.account_name,
                    "asin": row.asin,
                    "sku": row.sku or "",
                    "title": row.title or "",
                    "brand": row.brand or "",
                    "category": row.category or "Uncategorized",
                    "units": units,
                    "revenue": _round(revenue),
                    "orders": _as_int(row.orders),
                    "average_selling_price": _round(_safe_divide(revenue, units)),
                    "revenue_share": _round(_safe_divide(revenue * 100, total_revenue)),
                    "currency": row.currency or "EUR",
                }
            )
        return output

    async def _sales_category_rows(
        self,
        account_ids: list[UUID],
        start_date: date,
        end_date: date,
    ) -> list[dict[str, Any]]:
        """Build the category breakdown table for sales."""
        if not account_ids:
            return []

        category_expr = func.coalesce(Product.category, "Uncategorized")
        query = (
            select(
                category_expr.label("category"),
                func.sum(SalesData.units_ordered).label("units"),
                func.sum(SalesData.ordered_product_sales).label("revenue"),
                func.sum(SalesData.total_order_items).label("orders"),
                func.max(SalesData.currency).label("currency"),
            )
            .select_from(SalesData)
            .outerjoin(
                Product,
                and_(
                    Product.account_id == SalesData.account_id,
                    Product.asin == SalesData.asin,
                ),
            )
            .where(
                SalesData.account_id.in_(account_ids),
                SalesData.asin != DAILY_TOTAL_ASIN,
                SalesData.date >= start_date,
                SalesData.date <= end_date,
            )
            .group_by(category_expr)
            .order_by(func.sum(SalesData.ordered_product_sales).desc())
        )

        rows = (await self.db.execute(query)).all()
        total_revenue = sum(_as_float(row.revenue) for row in rows)
        output: list[dict[str, Any]] = []
        for row in rows:
            revenue = _as_float(row.revenue)
            units = _as_int(row.units)
            orders = _as_int(row.orders)
            output.append(
                {
                    "category": row.category or "Uncategorized",
                    "revenue": _round(revenue),
                    "units": units,
                    "orders": orders,
                    "average_order_value": _round(_safe_divide(revenue, orders)),
                    "average_selling_price": _round(_safe_divide(revenue, units)),
                    "revenue_share": _round(_safe_divide(revenue * 100, total_revenue)),
                    "currency": row.currency or "EUR",
                }
            )
        return output

    async def _latest_snapshot_date(
        self,
        account_ids: list[UUID],
        before_date: Optional[date] = None,
    ) -> Optional[date]:
        """Fetch the latest inventory snapshot date for the selected accounts."""
        if not account_ids:
            return None

        query = select(func.max(InventoryData.snapshot_date)).where(InventoryData.account_id.in_(account_ids))
        if before_date:
            query = query.where(InventoryData.snapshot_date < before_date)
        return (await self.db.execute(query)).scalar()

    async def _inventory_snapshot_rows(
        self,
        account_ids: list[UUID],
        snapshot_date: Optional[date],
        lang: Language,
    ) -> list[dict[str, Any]]:
        """Build enriched inventory rows for a specific snapshot."""
        if not account_ids or snapshot_date is None:
            return []

        sku_expr = func.coalesce(InventoryData.sku, Product.sku)
        query = (
            select(
                InventoryData.snapshot_date.label("snapshot_date"),
                AmazonAccount.account_name.label("account_name"),
                InventoryData.asin.label("asin"),
                sku_expr.label("sku"),
                InventoryData.fnsku.label("fnsku"),
                Product.title.label("title"),
                Product.brand.label("brand"),
                func.coalesce(Product.category, "Uncategorized").label("category"),
                InventoryData.afn_fulfillable_quantity.label("fba_fulfillable"),
                InventoryData.mfn_fulfillable_quantity.label("mfn_fulfillable"),
                InventoryData.afn_inbound_working_quantity.label("inbound_working"),
                InventoryData.afn_inbound_shipped_quantity.label("inbound_shipped"),
                InventoryData.afn_reserved_quantity.label("reserved_quantity"),
                InventoryData.afn_total_quantity.label("total_fba_quantity"),
            )
            .select_from(InventoryData)
            .join(AmazonAccount, AmazonAccount.id == InventoryData.account_id)
            .outerjoin(
                Product,
                and_(
                    Product.account_id == InventoryData.account_id,
                    Product.asin == InventoryData.asin,
                ),
            )
            .where(
                InventoryData.account_id.in_(account_ids),
                InventoryData.snapshot_date == snapshot_date,
            )
            .order_by(AmazonAccount.account_name, InventoryData.asin)
        )

        rows = (await self.db.execute(query)).all()
        output: list[dict[str, Any]] = []
        for row in rows:
            total_available = _as_int(row.fba_fulfillable) + _as_int(row.mfn_fulfillable)
            inbound_total = _as_int(row.inbound_working) + _as_int(row.inbound_shipped)
            status_key, severity_key = self._inventory_status_keys(total_available)
            output.append(
                {
                    "snapshot_date": row.snapshot_date,
                    "account_name": row.account_name,
                    "asin": row.asin,
                    "sku": row.sku or "",
                    "fnsku": row.fnsku or "",
                    "title": row.title or "",
                    "brand": row.brand or "",
                    "category": row.category or "Uncategorized",
                    "fba_fulfillable": _as_int(row.fba_fulfillable),
                    "mfn_fulfillable": _as_int(row.mfn_fulfillable),
                    "total_available": total_available,
                    "inbound_working": _as_int(row.inbound_working),
                    "inbound_shipped": _as_int(row.inbound_shipped),
                    "inbound_total": inbound_total,
                    "reserved_quantity": _as_int(row.reserved_quantity),
                    "total_fba_quantity": _as_int(row.total_fba_quantity),
                    "stock_status": self._text(lang, status_key),
                    "stock_severity": self._text(lang, severity_key),
                }
            )
        return output

    def _inventory_status_keys(self, total_available: int) -> tuple[str, str]:
        """Classify stock health with simple severity bands."""
        if total_available <= 0:
            return "status_out_of_stock", "severity_critical"
        if total_available < LOW_STOCK_THRESHOLD:
            return "status_low_stock", "severity_high"
        if total_available < AT_RISK_THRESHOLD:
            return "status_at_risk", "severity_medium"
        return "status_healthy", "severity_low"

    def _inventory_metrics(self, rows: list[dict[str, Any]], low_stock_only: bool) -> dict[str, Any]:
        """Summarize inventory rows for the summary CSV."""
        low_stock_rows = [row for row in rows if int(row["total_available"]) < LOW_STOCK_THRESHOLD]
        exported_rows = len(low_stock_rows) if low_stock_only else len(rows)
        return {
            "total_skus": len(rows),
            "exported_rows": exported_rows,
            "total_available": _round(sum(_as_float(row["total_available"]) for row in rows)),
            "total_inbound": _round(sum(_as_float(row["inbound_total"]) for row in rows)),
            "reserved_quantity": _round(sum(_as_float(row["reserved_quantity"]) for row in rows)),
            "low_stock_skus": len(low_stock_rows),
        }

    async def _advertising_summary(
        self,
        account_ids: list[UUID],
        start_date: date,
        end_date: date,
    ) -> dict[str, Any]:
        """Aggregate advertising metrics over the selected period."""
        if not account_ids:
            return {
                "impressions": 0,
                "clicks": 0,
                "spend": 0.0,
                "attributed_sales_7d": 0.0,
                "attributed_units_7d": 0,
                "ctr": 0.0,
                "cpc": 0.0,
                "acos": 0.0,
                "roas": 0.0,
                "active_campaigns": 0,
            }

        campaign_ids_query = select(AdvertisingCampaign.id).where(AdvertisingCampaign.account_id.in_(account_ids))
        totals_query = (
            select(
                func.sum(AdvertisingMetrics.impressions).label("impressions"),
                func.sum(AdvertisingMetrics.clicks).label("clicks"),
                func.sum(AdvertisingMetrics.cost).label("spend"),
                func.sum(AdvertisingMetrics.attributed_sales_7d).label("sales"),
                func.sum(AdvertisingMetrics.attributed_units_ordered_7d).label("units"),
                func.count(func.distinct(AdvertisingMetrics.campaign_id)).label("campaigns"),
            )
            .where(
                AdvertisingMetrics.campaign_id.in_(campaign_ids_query),
                AdvertisingMetrics.date >= start_date,
                AdvertisingMetrics.date <= end_date,
            )
        )
        row = (await self.db.execute(totals_query)).one()

        impressions = _as_int(row.impressions)
        clicks = _as_int(row.clicks)
        spend = _round(_as_float(row.spend))
        sales = _round(_as_float(row.sales))

        return {
            "impressions": impressions,
            "clicks": clicks,
            "spend": spend,
            "attributed_sales_7d": sales,
            "attributed_units_7d": _as_int(row.units),
            "ctr": _round(_safe_divide(clicks * 100, impressions)),
            "cpc": _round(_safe_divide(spend, clicks)),
            "acos": _round(_safe_divide(spend * 100, sales)),
            "roas": _round(_safe_divide(sales, spend)),
            "active_campaigns": _as_int(row.campaigns),
        }

    async def _advertising_daily_rows(
        self,
        account_ids: list[UUID],
        start_date: date,
        end_date: date,
    ) -> list[dict[str, Any]]:
        """Build daily advertising performance rows per campaign."""
        if not account_ids:
            return []

        query = (
            select(
                AdvertisingMetrics.date.label("report_date"),
                AmazonAccount.account_name.label("account_name"),
                AdvertisingCampaign.campaign_id.label("campaign_id"),
                AdvertisingCampaign.campaign_name.label("campaign_name"),
                AdvertisingCampaign.campaign_type.label("campaign_type"),
                AdvertisingCampaign.state.label("campaign_state"),
                AdvertisingCampaign.daily_budget.label("daily_budget"),
                AdvertisingMetrics.impressions.label("impressions"),
                AdvertisingMetrics.clicks.label("clicks"),
                AdvertisingMetrics.cost.label("spend"),
                AdvertisingMetrics.attributed_sales_7d.label("attributed_sales_7d"),
                AdvertisingMetrics.attributed_units_ordered_7d.label("attributed_units_7d"),
            )
            .select_from(AdvertisingMetrics)
            .join(AdvertisingCampaign, AdvertisingCampaign.id == AdvertisingMetrics.campaign_id)
            .join(AmazonAccount, AmazonAccount.id == AdvertisingCampaign.account_id)
            .where(
                AdvertisingCampaign.account_id.in_(account_ids),
                AdvertisingMetrics.date >= start_date,
                AdvertisingMetrics.date <= end_date,
            )
            .order_by(AdvertisingMetrics.date, AmazonAccount.account_name, AdvertisingCampaign.campaign_name)
        )

        rows = (await self.db.execute(query)).all()
        return [self._advertising_row_dict(row) for row in rows]

    async def _advertising_rollup_rows(
        self,
        account_ids: list[UUID],
        start_date: date,
        end_date: date,
    ) -> list[dict[str, Any]]:
        """Build campaign rollup rows aggregated across the period."""
        if not account_ids:
            return []

        query = (
            select(
                AmazonAccount.account_name.label("account_name"),
                AdvertisingCampaign.campaign_id.label("campaign_id"),
                AdvertisingCampaign.campaign_name.label("campaign_name"),
                AdvertisingCampaign.campaign_type.label("campaign_type"),
                AdvertisingCampaign.state.label("campaign_state"),
                AdvertisingCampaign.daily_budget.label("daily_budget"),
                func.sum(AdvertisingMetrics.impressions).label("impressions"),
                func.sum(AdvertisingMetrics.clicks).label("clicks"),
                func.sum(AdvertisingMetrics.cost).label("spend"),
                func.sum(AdvertisingMetrics.attributed_sales_7d).label("attributed_sales_7d"),
                func.sum(AdvertisingMetrics.attributed_units_ordered_7d).label("attributed_units_7d"),
            )
            .select_from(AdvertisingMetrics)
            .join(AdvertisingCampaign, AdvertisingCampaign.id == AdvertisingMetrics.campaign_id)
            .join(AmazonAccount, AmazonAccount.id == AdvertisingCampaign.account_id)
            .where(
                AdvertisingCampaign.account_id.in_(account_ids),
                AdvertisingMetrics.date >= start_date,
                AdvertisingMetrics.date <= end_date,
            )
            .group_by(
                AmazonAccount.account_name,
                AdvertisingCampaign.campaign_id,
                AdvertisingCampaign.campaign_name,
                AdvertisingCampaign.campaign_type,
                AdvertisingCampaign.state,
                AdvertisingCampaign.daily_budget,
            )
            .order_by(func.sum(AdvertisingMetrics.cost).desc(), func.sum(AdvertisingMetrics.attributed_sales_7d).desc())
        )

        rows = (await self.db.execute(query)).all()
        return [self._advertising_row_dict(row, include_date=False) for row in rows]

    def _advertising_row_dict(self, row: Any, include_date: bool = True) -> dict[str, Any]:
        """Normalize advertising rows and compute ratio metrics from totals."""
        clicks = _as_int(row.clicks)
        impressions = _as_int(row.impressions)
        spend = _as_float(row.spend)
        sales = _as_float(row.attributed_sales_7d)

        data = {
            "account_name": row.account_name,
            "campaign_id": row.campaign_id,
            "campaign_name": row.campaign_name or "",
            "campaign_type": row.campaign_type or "",
            "campaign_state": row.campaign_state or "",
            "daily_budget": _round(_as_float(row.daily_budget)) if row.daily_budget is not None else "",
            "impressions": impressions,
            "clicks": clicks,
            "spend": _round(spend),
            "attributed_sales_7d": _round(sales),
            "attributed_units_7d": _as_int(row.attributed_units_7d),
            "ctr": _round(_safe_divide(clicks * 100, impressions)),
            "cpc": _round(_safe_divide(spend, clicks)),
            "acos": _round(_safe_divide(spend * 100, sales)),
            "roas": _round(_safe_divide(sales, spend)),
        }
        if include_date:
            data["report_date"] = row.report_date
        return data
