"""Brand Analysis Automation service."""
from __future__ import annotations

import io
import json
import logging
import math
import re
from dataclasses import dataclass
from datetime import date, datetime, timedelta
from pathlib import Path
from typing import Any, Optional
from uuid import UUID

import pandas as pd
from sqlalchemy import and_, delete, func, select, text as sa_text
from sqlalchemy.ext.asyncio import AsyncSession, async_sessionmaker, create_async_engine
from sqlalchemy.orm import selectinload

from app.config import settings
from app.models.amazon_account import AmazonAccount
from app.models.brand_analysis import BrandAnalysisJob, BrandAnalysisSourceFile
from app.models.sales_data import SalesData
from app.services.data_extraction import DAILY_TOTAL_ASIN

logger = logging.getLogger(__name__)


NARRATIVE_TEMPLATE_VERSION = "brand-analysis-narrative-v2"
PPTX_TEMPLATE_VERSION = "brand-analysis-pptx-v2"


def _canonical_mode(mode: Optional[str]) -> str:
    """Normalize legacy mode aliases to the current canonical names.

    ``amazon_sp_api`` is treated as ``internal`` (the Inthezon-native
    SP-API + Market Research path). Deprecated external-provider modes
    already present in older rows fall back to manual upload.
    """
    if mode == "amazon_sp_api":
        return "internal"
    if mode in {"helium10_api", "helium10", "helium10_browser"}:
        return "manual"
    return mode or "internal"


STATUS_PROGRESS = {
    "pending": 0,
    "capability_checking": 8,
    "preflight_checking": 14,
    "internal_sync_requested": 20,
    "syncing_internal_data": 28,
    "internal_sync_completed": 34,
    "internal_sync_failed": 34,
    "collecting_source_data": 30,
    "enriching_catalog": 55,
    "generating_metrics": 70,
    "generating_narrative": 82,
    "analyzing": 70,
    "generating_pptx": 90,
    "cancelling": 95,
    "completed": 100,
    "completed_with_limitations": 100,
    "failed": 100,
    "cancelled": 100,
    "waiting_for_user_action": 50,
    # Legacy statuses kept for backwards compatibility with existing rows:
    "configuring_market": 10,
    "waiting_for_ready": 25,
    "exporting_2025": 40,
    "exporting_2024": 50,
}

# Statuses a job moves through while work is in flight. The start guard and
# the cancel endpoint both consult this set; keep it the single source of truth.
RUNNING_STATUSES = frozenset(
    {
        "capability_checking",
        "preflight_checking",
        "internal_sync_requested",
        "syncing_internal_data",
        "internal_sync_completed",
        "internal_sync_failed",
        "analyzing",
        "generating_metrics",
        "generating_narrative",
        "generating_pptx",
        "collecting_source_data",
        "enriching_catalog",
        "cancelling",
        # Legacy statuses kept for older rows:
        "configuring_market",
        "waiting_for_ready",
        "exporting_2025",
        "exporting_2024",
    }
)

# Statuses from which a job will never advance on its own.
TERMINAL_STATUSES = frozenset(
    {
        "completed",
        "completed_with_limitations",
        "failed",
        "cancelled",
    }
)

BRAND_ANALYSIS_YEARS = (2024, 2025)
SALES_AND_TRAFFIC_REPORT_TYPE = "GET_SALES_AND_TRAFFIC_REPORT"


REQUIRED_CANONICAL_COLUMNS = ("asin", "revenue")
OPTIONAL_CANONICAL_COLUMNS = (
    "product_name",
    "brand",
    "category",
    "subcategory",
    "units",
    "price",
    "rating",
    "reviews",
    "images",
    "sellers",
    "offer_count",
    "buy_box_owner",
    "bsr",
    "status",
    "bullets",
    "bullet_count",
    "description",
    "aplus_content",
    "has_aplus_content",
    "aplus_module_count",
    "text_module_count",
    "image_module_count",
    "aplus_source",
    "aplus_limitation",
    "fulfillment",
    "fba_fees",
    "actual_fba_fees",
    "estimated_fba_fees",
    "fee_source",
    "fee_confidence",
    "fee_limitation",
)

COMPLETENESS_OPTIONAL_FIELDS = (
    "price",
    "rating",
    "reviews",
    "images",
    "sellers",
    "offer_count",
    "buy_box_owner",
)


class BrandAnalysisDataError(ValueError):
    """Raised for invalid or unsupported source exports."""


class BrandAnalysisJobRunningError(RuntimeError):
    """Raised when an operation requires the job not to be actively running."""


def _revoke_celery_task(task_id: str) -> None:
    """Best-effort Celery revoke; never raises into the caller."""
    try:
        from workers.celery_app import celery_app

        celery_app.control.revoke(task_id, terminate=True, signal="SIGTERM")
    except Exception as exc:  # pragma: no cover - broker may be unavailable
        logger.warning("Failed to revoke brand analysis task %s: %s", task_id, exc)


class MissingColumnError(BrandAnalysisDataError):
    """Raised when a source export is missing required columns."""

    def __init__(self, missing: list[str], available: list[str], year: Optional[int] = None):
        self.missing = missing
        self.available = available
        self.year = year
        parts = []
        year_label = f" in {year} export" if year else ""
        for canonical in missing:
            aliases = COLUMN_ALIASES.get(canonical, [])
            parts.append(
                f"missing required {canonical} column{year_label}; "
                f"expected one of: {', '.join(aliases)}"
            )
        super().__init__(
            "; ".join(parts)
            + ". Available columns: "
            + ", ".join(available[:40])
        )


class InsufficientDataError(BrandAnalysisDataError):
    """Raised when a data source can't deliver enough rows for a year.

    Routed to the ``waiting_for_user_action`` status by the processor so
    the user can fall back to manual upload.
    """

    def __init__(self, year: int, source_name: str, detail: str = ""):
        self.year = year
        self.source_name = source_name
        message = f"Data source '{source_name}' did not return any data for {year}"
        if detail:
            message += f": {detail}"
        super().__init__(message)


@dataclass
class ColumnValidationReport:
    """Result of validating an export's columns against canonical aliases."""

    required_found: list[str]
    required_missing: list[str]
    optional_found: list[str]
    optional_missing: list[str]
    detected_mapping: dict[str, str]
    available_columns: list[str]

    def to_dict(self) -> dict[str, Any]:
        return {
            "required_found": list(self.required_found),
            "required_missing": list(self.required_missing),
            "optional_found": list(self.optional_found),
            "optional_missing": list(self.optional_missing),
            "detected_mapping": dict(self.detected_mapping),
            "available_columns": list(self.available_columns),
        }


@dataclass
class ParsedBrandExport:
    """Normalized ASIN-level export data."""

    rows: pd.DataFrame
    columns: list[str]
    row_count: int
    source_name: str = "manual_upload"
    year: Optional[int] = None
    validation: Optional[ColumnValidationReport] = None


COLUMN_ALIASES = {
    "asin": [
        "asin",
        "product asin",
        "child asin",
        "parent asin",
        "sku asin",
    ],
    "product_name": [
        "product",
        "product name",
        "title",
        "product title",
        "item name",
        "name",
    ],
    "brand": ["brand", "brand name"],
    "category": [
        "category",
        "product category",
        "category path",
        "browse node",
    ],
    "subcategory": [
        "subcategory",
        "sub category",
        "niche",
        "market",
        "sub-category",
    ],
    "revenue": [
        "revenue",
        "total revenue",
        "sales",
        "total sales",
        "ordered product sales",
        "sales revenue",
        "gross revenue",
        "estimated revenue",
    ],
    "units": [
        "units",
        "units sold",
        "unit sales",
        "quantity",
        "quantity sold",
        "sales quantity",
        "total units",
    ],
    "price": [
        "price",
        "avg price",
        "average price",
        "selling price",
        "list price",
    ],
    "rating": [
        "rating",
        "average rating",
        "avg rating",
        "product rating",
        "stars",
    ],
    "reviews": [
        "reviews",
        "review count",
        "reviews count",
        "total reviews",
        "ratings count",
        "number of reviews",
    ],
    "images": [
        "images",
        "image count",
        "number of images",
        "images count",
        "product images",
        "image_count",
        "images_count",
    ],
    "sellers": [
        "sellers",
        "seller count",
        "reseller count",
        "number of sellers",
        "sellers count",
        "seller_count",
        "reseller_count",
        "offer count",
        "offers",
    ],
    "offer_count": [
        "offer count",
        "offers count",
        "offers",
        "total offers",
        "number of offers",
    ],
    "buy_box_owner": [
        "buy box",
        "buy box owner",
        "buy box seller",
        "buybox seller",
        "seller",
        "reseller",
        "merchant",
    ],
    "bsr": [
        "bsr",
        "best sellers rank",
        "sales rank",
        "rank",
    ],
    "status": [
        "status",
        "active status",
        "listing status",
        "availability",
        "active/inactive status",
    ],
    "bullets": [
        "bullets",
        "bullet points",
        "bullet_point",
        "bullet point",
        "key features",
    ],
    "bullet_count": [
        "bullet count",
        "bullets count",
        "number of bullets",
        "bullet_count",
    ],
    "description": [
        "description",
        "product description",
        "listing description",
    ],
    "aplus_content": [
        "a+ content",
        "a plus content",
        "aplus content",
        "aplus_content",
        "enhanced brand content",
    ],
    "has_aplus_content": [
        "has a+ content",
        "has a plus content",
        "has_aplus_content",
        "a+ present",
    ],
    "aplus_module_count": [
        "a+ module count",
        "aplus module count",
        "aplus_module_count",
    ],
    "text_module_count": [
        "a+ text module count",
        "text module count",
        "text_module_count",
    ],
    "image_module_count": [
        "a+ image module count",
        "image module count",
        "image_module_count",
    ],
    "aplus_source": [
        "a+ source",
        "aplus source",
        "aplus_source",
    ],
    "aplus_limitation": [
        "a+ limitation",
        "aplus limitation",
        "aplus_limitation",
    ],
    "fulfillment": [
        "fulfillment",
        "fulfillment channel",
        "fulfillment type",
        "shipping mode",
    ],
    "fba_fees": [
        "fba fees",
        "fba fee",
        "fulfillment fee",
        "fulfillment fees",
        "amazon fees",
    ],
    "actual_fba_fees": [
        "actual fba fees",
        "actual fulfillment fees",
        "settlement fba fees",
    ],
    "estimated_fba_fees": [
        "estimated fba fees",
        "estimated fulfillment fees",
        "product fees estimate",
    ],
    "fee_source": [
        "fee source",
        "fba fee source",
    ],
    "fee_confidence": [
        "fee confidence",
        "fba fee confidence",
    ],
    "fee_limitation": [
        "fee limitation",
        "fba fee limitation",
    ],
}

NUMERIC_COLUMNS = {
    "revenue",
    "units",
    "price",
    "rating",
    "reviews",
    "images",
    "sellers",
    "offer_count",
    "bsr",
    "bullet_count",
    "fba_fees",
    "actual_fba_fees",
    "estimated_fba_fees",
    "aplus_module_count",
    "text_module_count",
    "image_module_count",
}


def _normalize_header(value: Any) -> str:
    text = str(value or "").strip().lower()
    text = text.replace("\ufeff", "")
    return re.sub(r"[^a-z0-9]+", "", text)


def normalize_brand_text(value: Any) -> str:
    """Normalize a brand label for exact/fuzzy matching."""
    text = str(value or "").strip().lower()
    text = text.replace("&", "and")
    return re.sub(r"[^a-z0-9]+", "", text)


def brand_matches(candidate: Any, target: Any) -> bool:
    """Return true for exact or conservative fuzzy brand matches.

    This is intentionally stricter than general search relevance: the
    pipeline may include zero-revenue discovered ASINs, so false positives
    would pollute revenue and catalog-health metrics.
    """
    candidate_norm = normalize_brand_text(candidate)
    target_norm = normalize_brand_text(target)
    if not candidate_norm or not target_norm:
        return False
    if candidate_norm == target_norm:
        return True
    if len(target_norm) >= 4 and (target_norm in candidate_norm or candidate_norm in target_norm):
        return True

    from difflib import SequenceMatcher

    return SequenceMatcher(None, candidate_norm, target_norm).ratio() >= 0.88


NORMALIZED_ALIASES = {
    canonical: {_normalize_header(alias) for alias in aliases}
    for canonical, aliases in COLUMN_ALIASES.items()
}


def _parse_number(value: Any) -> Optional[float]:
    if value is None:
        return None
    if isinstance(value, float) and math.isnan(value):
        return None
    if isinstance(value, (int, float)):
        return float(value)

    text = str(value).strip()
    if not text or text.lower() in {"nan", "none", "n/a", "-", "--"}:
        return None
    negative = text.startswith("(") and text.endswith(")")
    text = text.strip("()")
    text = text.replace("\u00a0", "").replace(" ", "")
    text = re.sub(r"[^\d,.\-']", "", text)
    text = text.replace("'", "")
    if not text or text in {"-", ".", ","}:
        return None

    if "," in text and "." in text:
        if text.rfind(",") > text.rfind("."):
            text = text.replace(".", "").replace(",", ".")
        else:
            text = text.replace(",", "")
    elif "," in text:
        parts = text.split(",")
        if len(parts[-1]) <= 2:
            text = "".join(parts[:-1]) + "." + parts[-1]
        else:
            text = text.replace(",", "")

    try:
        number = float(text)
    except ValueError:
        return None
    return -number if negative else number


def _first_present(values: pd.Series) -> Optional[str]:
    for value in values:
        if value is None:
            continue
        if isinstance(value, float) and math.isnan(value):
            continue
        text = str(value).strip()
        if text:
            return text
    return None


def _mean_present(values: pd.Series) -> Optional[float]:
    present = [float(value) for value in values if value is not None and not pd.isna(value)]
    return sum(present) / len(present) if present else None


def _max_present(values: pd.Series) -> Optional[float]:
    present = [float(value) for value in values if value is not None and not pd.isna(value)]
    return max(present) if present else None


def _safe_round(value: Optional[float], digits: int = 2) -> Optional[float]:
    if value is None:
        return None
    return round(float(value), digits)


def yoy_percent(current: Optional[float], previous: Optional[float]) -> Optional[float]:
    """Calculate YoY percentage without inventing a value for zero baselines."""
    current_value = float(current or 0)
    if previous is None:
        return 0.0 if current_value == 0 else None
    previous_value = float(previous or 0)
    if previous_value == 0:
        return 0.0 if current_value == 0 else None
    return round(((current_value - previous_value) / previous_value) * 100, 1)


def _pct(part: float, total: float) -> float:
    return round((part / total) * 100, 1) if total else 0.0


def _read_export_dataframe(file_bytes: bytes, filename: str) -> pd.DataFrame:
    suffix = Path(filename).suffix.lower()
    data = io.BytesIO(file_bytes)
    if suffix in {".xlsx", ".xlsm", ".xls"}:
        return pd.read_excel(data)

    encodings = ["utf-8-sig", "utf-8", "latin1"]
    last_exc: Optional[Exception] = None
    for encoding in encodings:
        data.seek(0)
        try:
            return pd.read_csv(data, sep=None, engine="python", encoding=encoding)
        except Exception as exc:
            last_exc = exc
    raise BrandAnalysisDataError(f"Could not read source file {filename}: {last_exc}")


def parse_brand_export(
    file_bytes: bytes,
    filename: str,
    *,
    year: Optional[int] = None,
) -> ParsedBrandExport:
    """Parse an external yearly product export (CSV/XLSX) into normalized ASIN-level rows.

    The column-alias mapping is flexible enough to support different
    ASIN-keyed yearly summary formats without coupling the product flow
    to a specific external vendor.
    """
    raw_df = _read_export_dataframe(file_bytes, filename)
    raw_df = raw_df.dropna(how="all")
    if raw_df.empty:
        raise BrandAnalysisDataError("Source file is empty")

    available = [str(column) for column in raw_df.columns]
    normalized_to_column = {_normalize_header(column): column for column in raw_df.columns}
    mapping: dict[str, Optional[str]] = {}
    for canonical, aliases in NORMALIZED_ALIASES.items():
        mapping[canonical] = None
        for normalized, column in normalized_to_column.items():
            if normalized in aliases:
                mapping[canonical] = column
                break

    detected_mapping = {canonical: source for canonical, source in mapping.items() if source}
    required_found = [c for c in REQUIRED_CANONICAL_COLUMNS if mapping.get(c)]
    required_missing = [c for c in REQUIRED_CANONICAL_COLUMNS if not mapping.get(c)]
    optional_found = [c for c in OPTIONAL_CANONICAL_COLUMNS if mapping.get(c)]
    optional_missing = [c for c in OPTIONAL_CANONICAL_COLUMNS if not mapping.get(c)]
    validation = ColumnValidationReport(
        required_found=required_found,
        required_missing=required_missing,
        optional_found=optional_found,
        optional_missing=optional_missing,
        detected_mapping=detected_mapping,
        available_columns=available,
    )

    if required_missing:
        raise MissingColumnError(required_missing, available, year=year)

    normalized = pd.DataFrame()
    for canonical in COLUMN_ALIASES:
        source = mapping.get(canonical)
        if source is None:
            normalized[canonical] = None
        elif canonical in NUMERIC_COLUMNS:
            normalized[canonical] = raw_df[source].apply(_parse_number)
        else:
            normalized[canonical] = raw_df[source].apply(
                lambda value: None if pd.isna(value) else str(value).strip()
            )

    normalized["asin"] = normalized["asin"].apply(lambda value: str(value or "").strip().upper())
    normalized = normalized[normalized["asin"] != ""]
    normalized = normalized[normalized["asin"] != DAILY_TOTAL_ASIN]
    if normalized.empty:
        raise BrandAnalysisDataError("Source file does not contain any valid ASIN rows")

    rows = []
    for asin, group in normalized.groupby("asin", sort=False):
        rows.append(
            {
                "asin": asin,
                "product_name": _first_present(group["product_name"]) or asin,
                "brand": _first_present(group["brand"]),
                "category": _first_present(group["category"]),
                "subcategory": _first_present(group["subcategory"]) or "Uncategorized",
                "revenue": float(group["revenue"].fillna(0).sum()),
                "units": float(group["units"].fillna(0).sum()) if mapping.get("units") else None,
                "price": _mean_present(group["price"]) if mapping.get("price") else None,
                "rating": _mean_present(group["rating"]) if mapping.get("rating") else None,
                "reviews": _max_present(group["reviews"]) if mapping.get("reviews") else None,
                "images": _max_present(group["images"]) if mapping.get("images") else None,
                "sellers": _max_present(group["sellers"]) if mapping.get("sellers") else None,
                "offer_count": _max_present(group["offer_count"]) if mapping.get("offer_count") else None,
                "buy_box_owner": _first_present(group["buy_box_owner"]) if mapping.get("buy_box_owner") else None,
                "bsr": _max_present(group["bsr"]) if mapping.get("bsr") else None,
                "status": _first_present(group["status"]) if mapping.get("status") else None,
                "bullets": _first_present(group["bullets"]) if mapping.get("bullets") else None,
                "bullet_count": _max_present(group["bullet_count"]) if mapping.get("bullet_count") else None,
                "description": _first_present(group["description"]) if mapping.get("description") else None,
                "aplus_content": _first_present(group["aplus_content"]) if mapping.get("aplus_content") else None,
                "has_aplus_content": _first_present(group["has_aplus_content"]) if mapping.get("has_aplus_content") else None,
                "aplus_module_count": _max_present(group["aplus_module_count"]) if mapping.get("aplus_module_count") else None,
                "text_module_count": _max_present(group["text_module_count"]) if mapping.get("text_module_count") else None,
                "image_module_count": _max_present(group["image_module_count"]) if mapping.get("image_module_count") else None,
                "aplus_source": _first_present(group["aplus_source"]) if mapping.get("aplus_source") else None,
                "aplus_limitation": _first_present(group["aplus_limitation"]) if mapping.get("aplus_limitation") else None,
                "fulfillment": _first_present(group["fulfillment"]) if mapping.get("fulfillment") else None,
                "fba_fees": _mean_present(group["fba_fees"]) if mapping.get("fba_fees") else None,
                "actual_fba_fees": _mean_present(group["actual_fba_fees"]) if mapping.get("actual_fba_fees") else None,
                "estimated_fba_fees": _mean_present(group["estimated_fba_fees"]) if mapping.get("estimated_fba_fees") else None,
                "fee_source": _first_present(group["fee_source"]) if mapping.get("fee_source") else None,
                "fee_confidence": _first_present(group["fee_confidence"]) if mapping.get("fee_confidence") else None,
                "fee_limitation": _first_present(group["fee_limitation"]) if mapping.get("fee_limitation") else None,
            }
        )

    return ParsedBrandExport(
        rows=pd.DataFrame(rows),
        columns=available,
        row_count=len(raw_df),
        source_name="manual_upload",
        year=year,
        validation=validation,
    )


def build_metric_provenance(
    export_2024: ParsedBrandExport,
    export_2025: ParsedBrandExport,
) -> dict[str, dict[str, Any]]:
    """Return a per-metric provenance map (years, source columns, formula).

    Every entry documents which data the metric was derived from so the
    PPTX/PDF output is auditable and never reliant on AI inference.
    """
    mapping_2024 = (export_2024.validation.detected_mapping if export_2024.validation else {})
    mapping_2025 = (export_2025.validation.detected_mapping if export_2025.validation else {})

    def cols(canonical: str, *years_data: tuple[int, dict[str, str]]) -> list[str]:
        names = []
        for year, mapping in years_data:
            src = mapping.get(canonical)
            if src:
                names.append(f"{year}:{src}")
        return names

    yrs_2024_2025 = ((2024, mapping_2024), (2025, mapping_2025))
    yrs_2025 = ((2025, mapping_2025),)
    yrs_2024 = ((2024, mapping_2024),)

    return {
        "total_revenue_2024": {
            "source_years": [2024],
            "source_columns": cols("revenue", *yrs_2024),
            "source_name": export_2024.source_name,
            "formula": "sum(revenue) over all ASIN rows for 2024",
        },
        "total_revenue_2025": {
            "source_years": [2025],
            "source_columns": cols("revenue", *yrs_2025),
            "source_name": export_2025.source_name,
            "formula": "sum(revenue) over all ASIN rows for 2025",
        },
        "yoy_percent": {
            "source_years": [2024, 2025],
            "source_columns": cols("revenue", *yrs_2024_2025),
            "source_name": export_2025.source_name,
            "formula": "(total_revenue_2025 - total_revenue_2024) / total_revenue_2024 * 100",
        },
        "weighted_average_rating": {
            "source_years": [2025],
            "source_columns": cols("rating", *yrs_2025) + cols("units", *yrs_2025),
            "source_name": export_2025.source_name,
            "formula": "sum(rating * units) / sum(units), falling back to revenue weights, then unweighted mean",
        },
        "total_units_sold_2025": {
            "source_years": [2025],
            "source_columns": cols("units", *yrs_2025),
            "source_name": export_2025.source_name,
            "formula": "sum(units) over all ASIN rows for 2025",
        },
        "average_price_per_asin": {
            "source_years": [2025],
            "source_columns": cols("price", *yrs_2025),
            "source_name": export_2025.source_name,
            "formula": "mean(price) per ASIN; fallback total_revenue_2025 / total_units_2025",
        },
        "active_asins_2025": {
            "source_years": [2025],
            "source_columns": cols("revenue", *yrs_2025),
            "source_name": export_2025.source_name,
            "formula": "count(ASIN where revenue > 0) for 2025",
        },
        "inactive_asins_2025": {
            "source_years": [2025],
            "source_columns": cols("revenue", *yrs_2025),
            "source_name": export_2025.source_name,
            "formula": "count(ASIN where revenue = 0) for 2025",
        },
        "new_asins_yoy": {
            "source_years": [2024, 2025],
            "source_columns": cols("asin", *yrs_2024_2025),
            "source_name": export_2025.source_name,
            "formula": "len(asin_set_2025 - asin_set_2024)",
        },
        "top_5_revenue_share": {
            "source_years": [2025],
            "source_columns": cols("revenue", *yrs_2025),
            "source_name": export_2025.source_name,
            "formula": "sum(top-5 revenue) / total_revenue_2025 * 100",
        },
        "top_10_revenue_share": {
            "source_years": [2025],
            "source_columns": cols("revenue", *yrs_2025),
            "source_name": export_2025.source_name,
            "formula": "sum(top-10 revenue) / total_revenue_2025 * 100",
        },
        "average_revenue_per_active_asin": {
            "source_years": [2025],
            "source_columns": cols("revenue", *yrs_2025),
            "source_name": export_2025.source_name,
            "formula": "total_revenue_2025 / active_asins_2025",
        },
        "revenue_by_subcategory": {
            "source_years": [2024, 2025],
            "source_columns": cols("subcategory", *yrs_2024_2025) + cols("revenue", *yrs_2024_2025),
            "source_name": export_2025.source_name,
            "formula": "group by subcategory and sum(revenue) per year",
        },
        "reseller_buy_box_distribution": {
            "source_years": [2025],
            "source_columns": cols("buy_box_owner", *yrs_2025) + cols("revenue", *yrs_2025),
            "source_name": export_2025.source_name,
            "formula": "group by buy_box_owner; count ASINs and sum(revenue) for 2025",
        },
        "market_share_2025": {
            "source_years": [2025],
            "source_columns": cols("brand", *yrs_2025) + cols("revenue", *yrs_2025),
            "source_name": export_2025.source_name,
            "formula": "brand revenue / total market revenue * 100; returned as N/A unless source rows contain a broad competitor market with brand revenue",
        },
        "content_health": {
            "source_years": [2025],
            "source_columns": (
                cols("product_name", *yrs_2025)
                + cols("bullets", *yrs_2025)
                + cols("bullet_count", *yrs_2025)
                + cols("description", *yrs_2025)
                + cols("aplus_content", *yrs_2025)
                + cols("images", *yrs_2025)
                + cols("reviews", *yrs_2025)
                + cols("rating", *yrs_2025)
            ),
            "source_name": export_2025.source_name,
            "formula": "deterministic counts of missing/thin listing content when those fields are present; unavailable fields remain N/A",
        },
    }


def build_metric_source_registry(metrics: dict[str, Any], source_name: str) -> dict[str, dict[str, Any]]:
    """Define source priority and quality for Brand Analysis metrics."""

    market = metrics.get("market_analysis") or {}
    fee = metrics.get("fee_summary") or {}
    content = metrics.get("content_health") or {}
    seller = metrics.get("seller_buy_box_summary") or {}
    internal_source = source_name == "internal"

    registry: dict[str, dict[str, Any]] = {
        "total_revenue_2024": {
            "preferred_source": "internal sales_data",
            "fallback_sources": ["Data Kiosk / Sales & Traffic sync", "external_upload"],
            "required_fields": ["asin", "revenue"],
            "quality": "exact" if metrics.get("total_revenue_2024") is not None else "unavailable",
            "formula": "sum(revenue) for source rows scoped to the brand/ASIN universe",
        },
        "total_revenue_2025": {
            "preferred_source": "internal sales_data",
            "fallback_sources": ["Data Kiosk / Sales & Traffic sync", "external_upload"],
            "required_fields": ["asin", "revenue"],
            "quality": "exact" if metrics.get("total_revenue_2025") is not None else "unavailable",
            "formula": "sum(revenue) for source rows scoped to the brand/ASIN universe",
        },
        "market_revenue_share": {
            "preferred_source": "trusted competitor revenue export or licensed internal competitor revenue",
            "fallback_sources": [],
            "required_fields": ["competitor brand", "competitor revenue"],
            "quality": "exact" if market.get("market_size_status") == "revenue_calculated" else "unavailable",
            "formula": "brand revenue / trusted total market revenue",
            "limitation": market.get("limitation"),
        },
        "search_purchase_share": {
            "preferred_source": "Brand Analytics Search Query/Catalog Performance",
            "fallback_sources": [],
            "required_fields": ["search purchase share"],
            "quality": "proxy" if market.get("search_purchase_share") is not None else "unavailable",
            "formula": "Brand Analytics search purchase share; never used as revenue share",
            "limitation": market.get("search_share_limitation"),
        },
        "search_click_share": {
            "preferred_source": "Brand Analytics Search Query/Catalog Performance",
            "fallback_sources": [],
            "required_fields": ["search click share"],
            "quality": "proxy" if market.get("search_click_share") is not None else "unavailable",
            "formula": "Brand Analytics click share",
            "limitation": market.get("search_share_limitation"),
        },
        "search_cart_add_share": {
            "preferred_source": "Brand Analytics Search Query/Catalog Performance",
            "fallback_sources": [],
            "required_fields": ["search cart add share"],
            "quality": "proxy" if market.get("search_cart_add_share") is not None else "unavailable",
            "formula": "Brand Analytics cart-add share",
            "limitation": market.get("search_share_limitation"),
        },
        "fba_fees": {
            "preferred_source": "settlement/finance actual fees",
            "fallback_sources": ["FBA fee reports", "Product Fees API estimate"],
            "required_fields": ["actual_fba_fees or estimated_fba_fees"],
            "quality": fee.get("fee_confidence", "unavailable"),
            "formula": "average actual FBA fees when available; otherwise average Product Fees API estimates",
            "limitation": fee.get("fee_limitation"),
        },
        "buy_box_percentage": {
            "preferred_source": "Sales & Traffic by ASIN buyBoxPercentage",
            "fallback_sources": [],
            "required_fields": ["buyBoxPercentage"],
            "quality": "unavailable",
            "formula": "N/A until buyBoxPercentage is stored from Sales & Traffic",
            "limitation": seller.get("buy_box_percentage_limitation"),
        },
        "buy_box_owner_history": {
            "preferred_source": "stored asin_offer_snapshots history",
            "fallback_sources": [],
            "required_fields": ["observed_at", "buy_box_owner_name"],
            "quality": "unavailable" if not seller.get("buy_box_owner_history_available") else "partial",
            "formula": "owner trend over actually observed snapshots only",
            "limitation": seller.get("buy_box_owner_history_limitation"),
        },
        "seller_count_current": {
            "preferred_source": "Product Pricing offers snapshot",
            "fallback_sources": ["stored market research snapshot"],
            "required_fields": ["seller_count", "offer_count"],
            "quality": "exact" if seller.get("seller_count_available") else "unavailable",
            "formula": "current Product Pricing offer count per ASIN",
        },
        "aplus_content": {
            "preferred_source": "A+ Content API",
            "fallback_sources": ["internal cached A+ detection"],
            "required_fields": ["content document / ASIN relation"],
            "quality": "exact" if content.get("aplus_content_available") else "unavailable",
            "formula": "count ASINs with A+ content only when A+ API/cache source confirms it",
            "limitation": content.get("aplus_limitation"),
        },
    }

    shared_exact_keys = [
        "yoy_percent",
        "weighted_average_rating",
        "total_units_sold_2025",
        "average_price_per_asin",
        "total_asins_2024",
        "total_asins_2025",
        "active_asins_2024",
        "active_asins_2025",
        "inactive_asins_2025",
        "new_asins_yoy",
        "top_5_asins",
        "top_5_revenue_share",
        "top_10_revenue_share",
        "average_revenue_per_active_asin",
        "revenue_by_subcategory",
        "percentage_inactive_asins",
        "declining_asins_count",
        "percentage_declining_asins_among_active",
        "content_health",
        "review_rating_weaknesses",
        "seller_buy_box_summary",
        "fulfillment_breakdown",
        "growth_projection_scenarios",
    ]
    for key in shared_exact_keys:
        registry.setdefault(
            key,
            {
                "preferred_source": "deterministic calculation from selected source rows",
                "fallback_sources": ["external_upload"] if internal_source else [],
                "required_fields": ["asin", "revenue"],
                "quality": "exact" if metrics.get(key) is not None else "unavailable",
                "formula": "deterministic calculation; see metric_provenance formula",
            },
        )
    return registry


def enrich_metric_provenance(
    provenance: dict[str, dict[str, Any]],
    metrics: dict[str, Any],
    source_name: str,
) -> dict[str, dict[str, Any]]:
    registry = build_metric_source_registry(metrics, source_name)
    enriched = dict(provenance)
    for metric_key, source_spec in registry.items():
        existing = dict(enriched.get(metric_key, {}))
        existing.update({key: value for key, value in source_spec.items() if value is not None})
        enriched[metric_key] = existing
    return enriched


DECK_NUMERIC_PROVENANCE_KEYS = {
    "total_revenue_2024",
    "total_revenue_2025",
    "yoy_percent",
    "weighted_average_rating",
    "total_units_sold_2025",
    "average_price_per_asin",
    "total_asins_2024",
    "total_asins_2025",
    "active_asins_2025",
    "inactive_asins_2025",
    "new_asins_yoy",
    "top_5_asins",
    "top_5_revenue_share",
    "top_10_revenue_share",
    "average_revenue_per_active_asin",
    "revenue_by_subcategory",
    "percentage_inactive_asins",
    "percentage_declining_asins_among_active",
    "seller_buy_box_summary",
    "content_health",
    "review_rating_weaknesses",
    "market_revenue_share",
    "search_purchase_share",
    "search_click_share",
    "search_cart_add_share",
    "fba_fees",
    "buy_box_percentage",
    "buy_box_owner_history",
    "seller_count_current",
    "aplus_content",
    "growth_projection_scenarios",
}


def validate_metric_provenance_for_deck(metrics: dict[str, Any], provenance: dict[str, Any]) -> None:
    """Fail before PPTX if any deck numeric family lacks provenance."""
    missing = []
    for key in DECK_NUMERIC_PROVENANCE_KEYS:
        if key not in provenance:
            missing.append(key)
            continue
        entry = provenance.get(key) or {}
        if not entry.get("formula") and not entry.get("limitation"):
            missing.append(key)
    if missing:
        raise BrandAnalysisDataError(
            "Cannot generate Brand Analysis PPTX because metric provenance is missing for: "
            + ", ".join(sorted(missing))
        )


def build_limitation_summary(metrics: dict[str, Any], capability_matrix: Optional[dict] = None, data_coverage: Optional[dict] = None) -> dict[str, Any]:
    limitations: list[dict[str, str]] = []

    market = metrics.get("market_analysis") or {}
    if market.get("limitation"):
        limitations.append({"area": "market_revenue_share", "message": str(market["limitation"])})
    if market.get("search_share_limitation"):
        limitations.append({"area": "search_share_proxy", "message": str(market["search_share_limitation"])})

    fee = metrics.get("fee_summary") or {}
    if fee.get("fee_limitation"):
        limitations.append({"area": "fba_fees", "message": str(fee["fee_limitation"])})

    content = metrics.get("content_health") or {}
    if content.get("aplus_limitation"):
        limitations.append({"area": "aplus_content", "message": str(content["aplus_limitation"])})

    seller = metrics.get("seller_buy_box_summary") or {}
    for key in ("buy_box_owner_history_limitation", "buy_box_percentage_limitation"):
        if seller.get(key):
            limitations.append({"area": key.replace("_limitation", ""), "message": str(seller[key])})

    for year, report in ((data_coverage or {}).get("years") or {}).items():
        if report.get("classification") != "complete":
            for message in report.get("limitations") or ["Internal sales coverage is not complete."]:
                limitations.append({"area": f"sales_data_{year}", "message": str(message)})

    missing_roles = (capability_matrix or {}).get("missing_roles") or []
    for reason in missing_roles:
        limitations.append({"area": "missing_permission", "message": str(reason)})

    return {
        "has_limitations": bool(limitations),
        "items": limitations,
    }


def assess_data_completeness(
    export_2024: ParsedBrandExport,
    export_2025: ParsedBrandExport,
) -> dict[str, Any]:
    """Summarize missing optional fields without blocking deterministic metrics.

    Revenue and ASIN are required. Catalog fields such as reviews, images,
    sellers and Buy Box are useful but not guaranteed by SP-API or uploaded
    exports. The pipeline leaves those values as N/A and surfaces this
    completeness report to the UI/deck metadata.
    """

    def missing_optional(df: pd.DataFrame) -> list[str]:
        missing = []
        for field in COMPLETENESS_OPTIONAL_FIELDS:
            if field not in df.columns or not df[field].notna().any():
                missing.append(field)
        return missing

    missing_2024 = missing_optional(export_2024.rows)
    missing_2025 = missing_optional(export_2025.rows)
    return {
        "required_years_present": [2024, 2025],
        "missing_optional_fields_2024": missing_2024,
        "missing_optional_fields_2025": missing_2025,
        "optional_fields_complete": not missing_2024 and not missing_2025,
        "source_names": {
            "2024": export_2024.source_name,
            "2025": export_2025.source_name,
        },
    }


def get_recoverable_sales_window(
    account: AmazonAccount,
    marketplace_id: str,
    report_type: str,
    now: datetime,
) -> dict[str, Any]:
    """Return the configured SP-API sales recovery window for preflight.

    Amazon can still reject a specific historical range; when that happens the
    worker stores the API error against the exact requested period.
    """
    recovery_days = settings.BRAND_ANALYSIS_SALES_TRAFFIC_RECOVERY_DAYS
    end = now.date()
    start = end - timedelta(days=recovery_days)
    return {
        "account_id": str(account.id),
        "marketplace_id": marketplace_id,
        "report_type": report_type,
        "start_date": start.isoformat(),
        "end_date": end.isoformat(),
        "configured_days": recovery_days,
        "source": "settings.BRAND_ANALYSIS_SALES_TRAFFIC_RECOVERY_DAYS",
    }


def _month_key(value: date) -> str:
    return f"{value.year}-{value.month:02d}"


def _quarter_key(value: date) -> str:
    return f"{value.year}-Q{((value.month - 1) // 3) + 1}"


def _year_months(year: int) -> list[str]:
    return [f"{year}-{month:02d}" for month in range(1, 13)]


def _year_quarters(year: int) -> list[str]:
    return [f"{year}-Q{quarter}" for quarter in range(1, 5)]


def _month_bounds(month_key: str) -> tuple[date, date]:
    year_s, month_s = month_key.split("-")
    year = int(year_s)
    month = int(month_s)
    start = date(year, month, 1)
    if month == 12:
        end = date(year, 12, 31)
    else:
        end = date(year, month + 1, 1) - timedelta(days=1)
    return start, end


def classify_sales_year_coverage(
    *,
    year: int,
    row_count: int,
    asin_count: int,
    dates: list[date],
    recoverable_start: date,
    recoverable_end: date,
) -> dict[str, Any]:
    """Classify internal sales_data coverage for one analysis year."""
    year_start = date(year, 1, 1)
    year_end = date(year, 12, 31)
    unique_dates = sorted({value for value in dates if value})
    months_present = sorted({_month_key(value) for value in unique_dates})
    quarters_present = sorted({_quarter_key(value) for value in unique_dates})
    first_date = unique_dates[0] if unique_dates else None
    last_date = unique_dates[-1] if unique_dates else None

    annual_aggregate_only = bool(
        row_count > 0
        and unique_dates
        and set(unique_dates) == {year_end}
        and asin_count > 0
    )
    missing_months = [] if annual_aggregate_only else [
        month for month in _year_months(year) if month not in months_present
    ]
    missing_quarters = [] if annual_aggregate_only else [
        quarter for quarter in _year_quarters(year) if quarter not in quarters_present
    ]

    missing_periods: list[dict[str, Any]] = []
    for month in missing_months:
        start, end = _month_bounds(month)
        recoverable = end >= recoverable_start and start <= recoverable_end
        missing_periods.append(
            {
                "label": month,
                "start_date": start.isoformat(),
                "end_date": end.isoformat(),
                "recoverable": recoverable,
            }
        )

    has_recoverable_gap = any(period["recoverable"] for period in missing_periods)
    if row_count == 0:
        if year_end >= recoverable_start and year_start <= recoverable_end:
            classification = "recoverable_gap"
        else:
            classification = "unavailable"
    elif not missing_months and not annual_aggregate_only:
        classification = "complete"
    elif has_recoverable_gap:
        classification = "recoverable_gap"
    elif len(months_present) >= settings.BRAND_ANALYSIS_PARTIAL_USABLE_MONTHS or annual_aggregate_only:
        classification = "partial_but_usable"
    else:
        classification = "unrecoverable_gap"

    sync_start = max(year_start, recoverable_start)
    sync_end = min(year_end, recoverable_end)
    recoverable_window = (
        {"start_date": sync_start.isoformat(), "end_date": sync_end.isoformat()}
        if sync_start <= sync_end and classification == "recoverable_gap"
        else None
    )

    limitations: list[str] = []
    if annual_aggregate_only:
        limitations.append(
            "ASIN-level Sales & Traffic rows are stored on the report end date; daily month coverage is not available."
        )
    if classification in {"partial_but_usable", "recoverable_gap", "unrecoverable_gap"} and missing_months:
        limitations.append("Internal sales_data coverage is incomplete for this calendar year.")
    if row_count == 0 and classification == "unavailable":
        limitations.append("No internal sales_data rows exist and the requested year is outside the configured recovery window.")

    return {
        "year": year,
        "classification": classification,
        "first_date": first_date.isoformat() if first_date else None,
        "last_date": last_date.isoformat() if last_date else None,
        "row_count": int(row_count),
        "asin_count": int(asin_count),
        "missing_months": missing_months,
        "missing_quarters": missing_quarters,
        "missing_periods": missing_periods,
        "recoverable_window": recoverable_window,
        "annual_aggregate_only": annual_aggregate_only,
        "limitations": limitations,
    }


async def inspect_internal_sales_data_coverage(
    db: AsyncSession,
    account: AmazonAccount,
    *,
    years: tuple[int, ...] = BRAND_ANALYSIS_YEARS,
    now: Optional[datetime] = None,
) -> dict[str, Any]:
    """Inspect internal sales_data readiness before metrics generation."""
    current = now or datetime.utcnow()
    recovery = get_recoverable_sales_window(
        account,
        account.marketplace_id,
        SALES_AND_TRAFFIC_REPORT_TYPE,
        current,
    )
    recoverable_start = date.fromisoformat(recovery["start_date"])
    recoverable_end = date.fromisoformat(recovery["end_date"])
    years_report: dict[str, Any] = {}

    for year in years:
        start = date(year, 1, 1)
        end = date(year, 12, 31)
        filters = and_(
            SalesData.account_id == account.id,
            SalesData.date >= start,
            SalesData.date <= end,
            SalesData.asin != DAILY_TOTAL_ASIN,
        )
        summary = await db.execute(
            select(
                func.count(SalesData.id),
                func.count(func.distinct(SalesData.asin)),
            ).where(filters)
        )
        row_count, asin_count = summary.one()
        date_result = await db.execute(select(SalesData.date).where(filters).distinct())
        dates = [row[0] for row in date_result.all() if row and row[0]]
        years_report[str(year)] = classify_sales_year_coverage(
            year=year,
            row_count=int(row_count or 0),
            asin_count=int(asin_count or 0),
            dates=dates,
            recoverable_start=recoverable_start,
            recoverable_end=recoverable_end,
        )

    sync_windows = {
        year: report["recoverable_window"]
        for year, report in years_report.items()
        if report.get("classification") == "recoverable_gap" and report.get("recoverable_window")
    }
    classifications = {year: report["classification"] for year, report in years_report.items()}
    return {
        "checked_at": current.isoformat(),
        "account_id": str(account.id),
        "marketplace_id": account.marketplace_id,
        "recoverable_sales_window": recovery,
        "years": years_report,
        "classifications": classifications,
        "needs_sync": bool(sync_windows),
        "sync_windows": sync_windows,
        "has_usable_data": any(
            report["classification"] in {"complete", "partial_but_usable", "recoverable_gap"}
            and report["row_count"] > 0
            for report in years_report.values()
        ),
    }


def calculate_brand_metrics(
    export_2024: ParsedBrandExport,
    export_2025: ParsedBrandExport,
    *,
    brand_name: str,
) -> dict[str, Any]:
    """Compute all requested Brand Analysis metrics before any narrative step."""
    market_df24 = export_2024.rows.copy()
    market_df25 = export_2025.rows.copy()
    for df in (market_df24, market_df25):
        for column in COLUMN_ALIASES:
            if column not in df.columns:
                df[column] = None
        if "product_name" in df.columns:
            df["product_name"] = df["product_name"].fillna(df["asin"])
        if "subcategory" in df.columns:
            df["subcategory"] = df["subcategory"].fillna("Uncategorized")

    def _has_any(df: pd.DataFrame, column: str) -> bool:
        return bool(column in df.columns and df[column].notna().any())

    def _has_text(value: Any) -> bool:
        if value is None:
            return False
        if isinstance(value, (list, tuple, set)):
            return any(_has_text(item) for item in value)
        try:
            if pd.isna(value):
                return False
        except (TypeError, ValueError):
            pass
        return bool(str(value).strip())

    def _text_present_any(df: pd.DataFrame, column: str) -> bool:
        return bool(column in df.columns and df[column].apply(_has_text).any())

    def _truthy_content(value: Any) -> bool:
        if isinstance(value, bool):
            return value
        if value is None:
            return False
        text = str(value).strip().lower()
        return text in {"1", "true", "yes", "y", "present", "available", "has_aplus", "has a+ content"}

    def _brand_mask(df: pd.DataFrame) -> Optional[pd.Series]:
        if not _text_present_any(df, "brand"):
            return None
        mask = df["brand"].apply(lambda value: brand_matches(value, brand_name))
        return mask if bool(mask.any()) else None

    # Generic external yearly market exports often contain the target brand
    # plus competitors. In that case brand metrics are calculated on the
    # target-brand rows and market size/share on the full export. Internal
    # Inthezon source rows are already scoped to the job's brand/ASIN set.
    mask24 = _brand_mask(market_df24)
    mask25 = _brand_mask(market_df25)
    can_filter_external_market = (
        (export_2024.source_name == "manual_upload" or export_2025.source_name == "manual_upload")
        and (mask24 is not None or mask25 is not None)
    )
    df24 = market_df24[mask24].copy() if can_filter_external_market and mask24 is not None else market_df24.copy()
    df25 = market_df25[mask25].copy() if can_filter_external_market and mask25 is not None else market_df25.copy()

    previous_by_asin = {
        row.asin: float(row.revenue or 0)
        for row in df24.itertuples(index=False)
    }

    total_revenue_2024 = float(df24["revenue"].fillna(0).sum())
    total_revenue_2025 = float(df25["revenue"].fillna(0).sum())
    total_units_2025 = float(df25["units"].fillna(0).sum()) if _has_any(df25, "units") else None
    total_asins_2024 = int(df24["asin"].nunique())
    total_asins_2025 = int(df25["asin"].nunique())
    active_asins_2024 = int((df24["revenue"].fillna(0) > 0).sum())
    active_asins_2025 = int((df25["revenue"].fillna(0) > 0).sum())
    inactive_asins_2025 = int((df25["revenue"].fillna(0) == 0).sum())
    asin_set_2024 = set(df24["asin"])
    asin_set_2025 = set(df25["asin"])

    rating_rows = df25[df25["rating"].notna()].copy() if _has_any(df25, "rating") else pd.DataFrame()
    weighted_average_rating = None
    if not rating_rows.empty:
        weights = rating_rows["units"].fillna(0)
        if float(weights.sum()) <= 0:
            weights = rating_rows["revenue"].fillna(0)
        if float(weights.sum()) > 0:
            weighted_average_rating = float((rating_rows["rating"] * weights).sum() / weights.sum())
        else:
            weighted_average_rating = float(rating_rows["rating"].mean())

    if _has_any(df25, "price"):
        average_price_per_asin = float(df25["price"].dropna().mean())
    elif total_units_2025:
        average_price_per_asin = total_revenue_2025 / total_units_2025
    else:
        average_price_per_asin = None

    df25["revenue_2024"] = df25["asin"].map(previous_by_asin).fillna(0)
    df25["yoy_percent"] = df25.apply(lambda row: yoy_percent(row["revenue"], row["revenue_2024"]), axis=1)
    top_2025 = df25.sort_values("revenue", ascending=False).head(10)
    top_5_asins = [
        {
            "asin": row.asin,
            "product_name": row.product_name,
            "revenue_2025": round(float(row.revenue or 0), 2),
            "revenue_2024": round(float(row.revenue_2024 or 0), 2),
            "yoy_percent": row.yoy_percent,
        }
        for row in top_2025.head(5).itertuples(index=False)
    ]

    comparable_active = df25[(df25["revenue"] > 0) & (df25["revenue_2024"] > 0)]
    declining_asins = comparable_active[comparable_active["revenue"] < comparable_active["revenue_2024"]]

    sub24 = df24.groupby("subcategory", dropna=False)["revenue"].sum().to_dict()
    sub25 = df25.groupby("subcategory", dropna=False)["revenue"].sum().to_dict()
    subcategories = sorted(set(sub24) | set(sub25), key=lambda name: float(sub25.get(name, 0)), reverse=True)
    revenue_by_subcategory = []
    for subcategory in subcategories:
        revenue_2025 = float(sub25.get(subcategory, 0))
        revenue_2024 = float(sub24.get(subcategory, 0))
        revenue_by_subcategory.append(
            {
                "subcategory": str(subcategory or "Uncategorized"),
                "revenue_2025": round(revenue_2025, 2),
                "revenue_2024": round(revenue_2024, 2),
                "yoy_percent": yoy_percent(revenue_2025, revenue_2024),
            }
        )

    declining_subcategories = [
        item for item in revenue_by_subcategory
        if item["yoy_percent"] is not None and item["yoy_percent"] < 0
    ]
    largest_decline = min(declining_subcategories, key=lambda item: item["yoy_percent"], default=None)

    sellers_gt_one = None
    average_seller_count = None
    if _has_any(df25, "sellers"):
        sellers_gt_one = int((df25["sellers"].fillna(0) > 1).sum())
        average_seller_count = float(df25["sellers"].dropna().mean())

    average_offer_count = None
    if _has_any(df25, "offer_count"):
        average_offer_count = float(df25["offer_count"].dropna().mean())

    avg_images = None
    asins_with_fewer_than_5_images = None
    if _has_any(df25, "images"):
        avg_images = float(df25["images"].dropna().mean())
        asins_with_fewer_than_5_images = int((df25["images"].fillna(0) < 5).sum())

    asins_with_fewer_than_15_reviews = None
    if _has_any(df25, "reviews"):
        asins_with_fewer_than_15_reviews = int((df25["reviews"].fillna(0) < 15).sum())

    asins_with_rating_below_4 = None
    if _has_any(df25, "rating"):
        asins_with_rating_below_4 = int((df25["rating"].dropna() < 4.0).sum())

    reseller_distribution = []
    buy_box_owner_is_historical = export_2025.source_name != "internal"
    if _text_present_any(df25, "buy_box_owner") and buy_box_owner_is_historical:
        grouped = (
            df25.assign(buy_box_owner=df25["buy_box_owner"].fillna("Unknown"))
            .groupby("buy_box_owner", dropna=False)
            .agg(asin_count=("asin", "count"), revenue=("revenue", "sum"))
            .reset_index()
            .sort_values("revenue", ascending=False)
        )
        reseller_distribution = [
            {
                "reseller": str(row.buy_box_owner),
                "asin_count": int(row.asin_count),
                "revenue": round(float(row.revenue or 0), 2),
                "share_percent": _pct(float(row.revenue or 0), total_revenue_2025),
            }
            for row in grouped.itertuples(index=False)
        ]

    current_buy_box_snapshot_distribution = []
    if _text_present_any(df25, "buy_box_owner"):
        grouped_current = (
            df25.assign(buy_box_owner=df25["buy_box_owner"].fillna("Unknown"))
            .groupby("buy_box_owner", dropna=False)
            .agg(asin_count=("asin", "count"))
            .reset_index()
            .sort_values("asin_count", ascending=False)
        )
        current_buy_box_snapshot_distribution = [
            {
                "reseller": str(row.buy_box_owner),
                "asin_count": int(row.asin_count),
                "source": "current_product_pricing_snapshot" if export_2025.source_name == "internal" else export_2025.source_name,
            }
            for row in grouped_current.itertuples(index=False)
        ]

    buy_box_missing_asins = None
    if _text_present_any(df25, "buy_box_owner"):
        buy_box_missing_asins = int(total_asins_2025 - df25["buy_box_owner"].apply(_has_text).sum())

    bullet_counts = None
    if _has_any(df25, "bullet_count"):
        bullet_counts = df25["bullet_count"]
    elif _text_present_any(df25, "bullets"):
        bullet_counts = df25["bullets"].apply(
            lambda value: len([part for part in re.split(r"\n|;|\|", str(value)) if part.strip()])
            if _has_text(value) else None
        )

    asins_missing_bullets = None
    asins_with_fewer_than_3_bullets = None
    if bullet_counts is not None and pd.Series(bullet_counts).notna().any():
        asins_missing_bullets = int(pd.Series(bullet_counts).fillna(0).eq(0).sum())
        asins_with_fewer_than_3_bullets = int(pd.Series(bullet_counts).fillna(0).lt(3).sum())

    def _missing_text_count(column: str) -> Optional[int]:
        if not _text_present_any(df25, column):
            return None
        return int(total_asins_2025 - df25[column].apply(_has_text).sum())

    asins_missing_description = _missing_text_count("description")
    if _has_any(df25, "has_aplus_content"):
        has_aplus_series = df25["has_aplus_content"].apply(_truthy_content)
        asins_missing_aplus_content = int(total_asins_2025 - has_aplus_series.sum())
        aplus_content_available = True
    else:
        asins_missing_aplus_content = _missing_text_count("aplus_content")
        aplus_content_available = _text_present_any(df25, "aplus_content")
    title_missing_count = int(
        sum(
            1
            for row in df25.itertuples(index=False)
            if not _has_text(getattr(row, "product_name", None))
            or str(getattr(row, "product_name", "")).strip().upper() == str(row.asin).strip().upper()
        )
    )
    short_title_count = int(
        sum(
            1
            for row in df25.itertuples(index=False)
            if _has_text(getattr(row, "product_name", None))
            and str(getattr(row, "product_name")).strip().upper() != str(row.asin).strip().upper()
            and len(str(getattr(row, "product_name")).strip()) < 60
        )
    )

    content_gap_asins = []
    for row in top_2025.head(10).itertuples(index=False):
        issues = []
        title = str(getattr(row, "product_name", "") or "").strip()
        if not title or title.upper() == str(row.asin).upper():
            issues.append("missing_title")
        elif len(title) < 60:
            issues.append("short_title")
        row_bullets = getattr(row, "bullet_count", None)
        if row_bullets is None and _has_text(getattr(row, "bullets", None)):
            row_bullets = len([part for part in re.split(r"\n|;|\|", str(getattr(row, "bullets"))) if part.strip()])
        if row_bullets is not None and not pd.isna(row_bullets) and float(row_bullets) < 3:
            issues.append("thin_bullets")
        if asins_missing_description is not None and not _has_text(getattr(row, "description", None)):
            issues.append("missing_description")
        if getattr(row, "images", None) is not None and not pd.isna(getattr(row, "images")) and float(getattr(row, "images")) < 5:
            issues.append("low_image_count")
        if getattr(row, "reviews", None) is not None and not pd.isna(getattr(row, "reviews")) and float(getattr(row, "reviews")) < 15:
            issues.append("low_reviews")
        if getattr(row, "rating", None) is not None and not pd.isna(getattr(row, "rating")) and float(getattr(row, "rating")) < 4.0:
            issues.append("low_rating")
        if issues:
            content_gap_asins.append(
                {
                    "asin": row.asin,
                    "product_name": title or row.asin,
                    "issues": issues,
                    "revenue_2025": round(float(row.revenue or 0), 2),
                }
            )

    review_weakness_asins = []
    for row in top_2025.head(10).itertuples(index=False):
        issues = []
        if getattr(row, "reviews", None) is not None and not pd.isna(getattr(row, "reviews")) and float(getattr(row, "reviews")) < 15:
            issues.append("low_reviews")
        if getattr(row, "rating", None) is not None and not pd.isna(getattr(row, "rating")) and float(getattr(row, "rating")) < 4.0:
            issues.append("low_rating")
        if issues:
            review_weakness_asins.append(
                {
                    "asin": row.asin,
                    "product_name": getattr(row, "product_name", row.asin),
                    "reviews": _safe_round(getattr(row, "reviews", None), 0),
                    "rating": _safe_round(getattr(row, "rating", None), 2),
                    "issues": issues,
                }
            )

    fulfillment_breakdown = []
    if _text_present_any(df25, "fulfillment"):
        fulfillment_grouped = (
            df25.assign(fulfillment=df25["fulfillment"].fillna("Unknown"))
            .groupby("fulfillment", dropna=False)
            .agg(asin_count=("asin", "count"), revenue=("revenue", "sum"))
            .reset_index()
            .sort_values("revenue", ascending=False)
        )
        fulfillment_breakdown = [
            {
                "fulfillment": str(row.fulfillment),
                "asin_count": int(row.asin_count),
                "revenue": round(float(row.revenue or 0), 2),
                "share_percent": _pct(float(row.revenue or 0), total_revenue_2025),
            }
            for row in fulfillment_grouped.itertuples(index=False)
        ]

    actual_fba_fee_series = None
    if _has_any(df25, "actual_fba_fees"):
        actual_fba_fee_series = df25["actual_fba_fees"]
    elif _has_any(df25, "fba_fees"):
        actual_fba_fee_series = df25["fba_fees"]
    estimated_fba_fee_series = df25["estimated_fba_fees"] if _has_any(df25, "estimated_fba_fees") else None
    average_actual_fba_fees = float(actual_fba_fee_series.dropna().mean()) if actual_fba_fee_series is not None else None
    average_estimated_fba_fees = (
        float(estimated_fba_fee_series.dropna().mean())
        if estimated_fba_fee_series is not None
        else None
    )
    average_fba_fees = average_actual_fba_fees if average_actual_fba_fees is not None else average_estimated_fba_fees
    fee_source = "actual_settlement_or_upload" if average_actual_fba_fees is not None else (
        "product_fees_api_estimate" if average_estimated_fba_fees is not None else "unavailable"
    )
    fee_confidence = "actual" if average_actual_fba_fees is not None else (
        "estimated" if average_estimated_fba_fees is not None else "unavailable"
    )
    fee_limitation = None
    if fee_confidence == "estimated":
        fee_limitation = "Product Fees API values are estimates based on current price and are not actual settlement fees."
    elif fee_confidence == "unavailable":
        fee_limitation = "FBA fees require settlement/finance reports, FBA fee reports, or Product Fees API estimates."

    market_has_brand_column = _text_present_any(market_df25, "brand")
    market_brand_values = (
        {normalize_brand_text(value) for value in market_df25["brand"] if _has_text(value)}
        if market_has_brand_column else set()
    )
    broad_market_available = (
        can_filter_external_market
        and market_has_brand_column
        and bool(market_brand_values)
        and (len(market_brand_values) > 1 or len(market_df25) > len(df25))
    )
    market_size_2024 = float(market_df24["revenue"].fillna(0).sum()) if broad_market_available else None
    market_size_2025 = float(market_df25["revenue"].fillna(0).sum()) if broad_market_available else None
    market_share_2024 = _pct(total_revenue_2024, market_size_2024) if market_size_2024 else None
    market_share_2025 = _pct(total_revenue_2025, market_size_2025) if market_size_2025 else None
    market_share_status = "calculated_from_external_market_export" if broad_market_available else "not_available"
    market_size_status = "revenue_calculated" if broad_market_available else (
        "partial_internal_only" if export_2025.source_name == "internal" else "not_available"
    )
    market_share_limitation = None if broad_market_available else (
        "Market share requires a broad ASIN-level market dataset with competitor revenue. "
        "Internal Amazon data is brand/account-scoped and SP-API catalog search does not expose reliable revenue."
    )

    competitive_brand_distribution = []
    if broad_market_available:
        brand_grouped = (
            market_df25.assign(brand=market_df25["brand"].fillna("Unknown"))
            .groupby("brand", dropna=False)
            .agg(asin_count=("asin", "nunique"), revenue=("revenue", "sum"))
            .reset_index()
            .sort_values("revenue", ascending=False)
        )
        competitive_brand_distribution = [
            {
                "brand": str(row.brand),
                "asin_count": int(row.asin_count),
                "revenue": round(float(row.revenue or 0), 2),
                "market_share_percent": _pct(float(row.revenue or 0), market_size_2025 or 0),
            }
            for row in brand_grouped.head(10).itertuples(index=False)
        ]

    top_5_revenue = float(top_2025.head(5)["revenue"].fillna(0).sum())
    top_10_revenue = float(top_2025.head(10)["revenue"].fillna(0).sum())
    projection_ranges = {
        "conservative": {
            "growth_low": 10,
            "growth_high": 15,
            "revenue_low": round(total_revenue_2025 * 1.10, 2),
            "revenue_high": round(total_revenue_2025 * 1.15, 2),
        },
        "realistic": {
            "growth_low": 25,
            "growth_high": 35,
            "revenue_low": round(total_revenue_2025 * 1.25, 2),
            "revenue_high": round(total_revenue_2025 * 1.35, 2),
        },
        "optimistic": {
            "growth_low": 40,
            "growth_high": 55,
            "revenue_low": round(total_revenue_2025 * 1.40, 2),
            "revenue_high": round(total_revenue_2025 * 1.55, 2),
        },
    }

    metrics = {
        "brand_name": brand_name,
        "currency": "EUR",
        "source_years": [2024, 2025],
        "total_revenue_2024": round(total_revenue_2024, 2),
        "total_revenue_2025": round(total_revenue_2025, 2),
        "yoy_percent": yoy_percent(total_revenue_2025, total_revenue_2024),
        "weighted_average_rating": _safe_round(weighted_average_rating, 2),
        "total_units_sold_2025": int(total_units_2025) if total_units_2025 is not None else None,
        "average_price_per_asin": _safe_round(average_price_per_asin, 2),
        "total_asins_2024": total_asins_2024,
        "total_asins_2025": total_asins_2025,
        "active_asins_2024": active_asins_2024,
        "active_asins_2025": active_asins_2025,
        "new_asins_yoy": len(asin_set_2025 - asin_set_2024),
        "inactive_asins_2025": inactive_asins_2025,
        "top_5_asins": top_5_asins,
        "average_images_per_asin": _safe_round(avg_images, 1),
        "asins_with_fewer_than_5_images": asins_with_fewer_than_5_images,
        "asins_with_fewer_than_15_reviews": asins_with_fewer_than_15_reviews,
        "revenue_by_subcategory": revenue_by_subcategory,
        "percentage_inactive_asins": _pct(inactive_asins_2025, total_asins_2025),
        "declining_asins_count": int(len(declining_asins)),
        "comparable_active_asins_count": int(len(comparable_active)),
        "percentage_declining_asins_among_active": _pct(len(declining_asins), len(comparable_active)),
        "asins_with_more_than_1_seller": sellers_gt_one,
        "average_seller_count": _safe_round(average_seller_count, 1),
        "average_offer_count": _safe_round(average_offer_count, 1),
        "asins_missing_buy_box_owner": buy_box_missing_asins,
        "subcategory_with_largest_decline": largest_decline,
        "top_5_revenue_share": _pct(top_5_revenue, total_revenue_2025),
        "top_10_revenue_share": _pct(top_10_revenue, total_revenue_2025),
        "average_revenue_per_active_asin": _safe_round(
            total_revenue_2025 / active_asins_2025 if active_asins_2025 else None,
            2,
        ),
        "reseller_buy_box_distribution": reseller_distribution,
        "seller_buy_box_summary": {
            "seller_count_available": _has_any(df25, "sellers"),
            "offer_count_available": _has_any(df25, "offer_count"),
            "buy_box_owner_available": _text_present_any(df25, "buy_box_owner"),
            "buy_box_owner_history_available": buy_box_owner_is_historical and _text_present_any(df25, "buy_box_owner"),
            "buy_box_owner_history_limitation": None if buy_box_owner_is_historical else (
                "Historical Buy Box owner is N/A because current Product Pricing snapshots cannot be used as historical ownership."
            ),
            "current_snapshot_available": _text_present_any(df25, "buy_box_owner") or _has_any(df25, "sellers") or _has_any(df25, "offer_count"),
            "average_seller_count": _safe_round(average_seller_count, 1),
            "average_offer_count": _safe_round(average_offer_count, 1),
            "asins_with_more_than_1_seller": sellers_gt_one,
            "asins_missing_buy_box_owner": buy_box_missing_asins,
            "current_buy_box_snapshot_distribution": current_buy_box_snapshot_distribution,
            "buy_box_percentage": None,
            "buy_box_percentage_limitation": "Sales & Traffic buyBoxPercentage is not stored in the current sales_data model.",
        },
        "content_health": {
            "title_missing_count": title_missing_count,
            "short_title_count": short_title_count,
            "asins_missing_bullets": asins_missing_bullets,
            "asins_with_fewer_than_3_bullets": asins_with_fewer_than_3_bullets,
            "asins_missing_description": asins_missing_description,
            "asins_missing_aplus_content": asins_missing_aplus_content,
            "aplus_content_available": aplus_content_available,
            "aplus_source": _first_present(df25["aplus_source"]) if _text_present_any(df25, "aplus_source") else "unavailable",
            "aplus_limitation": _first_present(df25["aplus_limitation"]) if _text_present_any(df25, "aplus_limitation") else (
                None if aplus_content_available else "A+ Content API data is unavailable or no A+ document was found."
            ),
            "content_gap_asins": content_gap_asins,
        },
        "review_rating_weaknesses": {
            "asins_with_fewer_than_15_reviews": asins_with_fewer_than_15_reviews,
            "asins_with_rating_below_4": asins_with_rating_below_4,
            "weak_asins": review_weakness_asins,
        },
        "fulfillment_breakdown": fulfillment_breakdown,
        "average_fba_fees": _safe_round(average_fba_fees, 2),
        "average_actual_fba_fees": _safe_round(average_actual_fba_fees, 2),
        "average_estimated_fba_fees": _safe_round(average_estimated_fba_fees, 2),
        "fee_summary": {
            "actual_fba_fees": _safe_round(average_actual_fba_fees, 2),
            "estimated_fba_fees": _safe_round(average_estimated_fba_fees, 2),
            "fee_source": fee_source,
            "fee_confidence": fee_confidence,
            "fee_limitation": fee_limitation,
        },
        "market_size_2024": _safe_round(market_size_2024, 2),
        "market_size_2025": _safe_round(market_size_2025, 2),
        "market_share_2024": market_share_2024,
        "market_share_2025": market_share_2025,
        "market_analysis": {
            "status": market_share_status,
            "market_size_status": market_size_status,
            "market_size_2024": _safe_round(market_size_2024, 2),
            "market_size_2025": _safe_round(market_size_2025, 2),
            "market_share_2024": market_share_2024,
            "market_share_2025": market_share_2025,
            "limitation": market_share_limitation,
            "competitive_brand_distribution": competitive_brand_distribution,
            "search_purchase_share": None,
            "search_click_share": None,
            "search_cart_add_share": None,
            "search_share_limitation": "Brand Analytics search share is unavailable unless the Brand Analytics capability is present; it is never used as revenue market share.",
        },
        "growth_projection_scenarios": projection_ranges,
        "rules": {
            "can_mention_vine": total_revenue_2025 >= 100000,
        },
        "automation_metadata": {
            "narrative_template_version": NARRATIVE_TEMPLATE_VERSION,
            "pptx_template_version": PPTX_TEMPLATE_VERSION,
            "numbers_are_deterministic": True,
            "ai_role": "narrative_only",
        },
        "source_metadata": {
            "rows_2024": export_2024.row_count,
            "rows_2025": export_2025.row_count,
            "brand_rows_2024": int(len(df24)),
            "brand_rows_2025": int(len(df25)),
            "market_rows_2024": int(len(market_df24)),
            "market_rows_2025": int(len(market_df25)),
            "columns_2024": export_2024.columns,
            "columns_2025": export_2025.columns,
        },
    }
    return _json_safe(metrics)


def _json_safe(value: Any) -> Any:
    if isinstance(value, dict):
        return {key: _json_safe(item) for key, item in value.items()}
    if isinstance(value, list):
        return [_json_safe(item) for item in value]
    if value is None:
        return None
    if isinstance(value, float):
        if math.isnan(value) or math.isinf(value):
            return None
        return value
    if pd.isna(value):
        return None
    return value


def _contains_vine(value: Any) -> bool:
    return isinstance(value, str) and "vine" in value.lower()


def _remove_vine_mentions(value: Any) -> Any:
    if isinstance(value, dict):
        return {key: _remove_vine_mentions(item) for key, item in value.items() if not _contains_vine(item)}
    if isinstance(value, list):
        return [_remove_vine_mentions(item) for item in value if not _contains_vine(item)]
    return "" if _contains_vine(value) else value


def build_fallback_narrative(metrics: dict[str, Any], language: str = "en") -> dict[str, Any]:
    """Build deterministic strategic text when AI is unavailable."""
    brand = metrics["brand_name"]
    top_category = (metrics.get("revenue_by_subcategory") or [{}])[0].get("subcategory", "core categories")
    decline = metrics.get("subcategory_with_largest_decline") or {}
    active = metrics.get("active_asins_2025") or 0
    total = metrics.get("total_asins_2025") or 0
    inactive_pct = metrics.get("percentage_inactive_asins") or 0

    narrative = {
        "overview": (
            f"{brand} generated {format_currency(metrics.get('total_revenue_2025'))} in 2025 "
            f"with {active} active ASINs out of {total}. The main opportunity is to recover inactive catalog coverage "
            f"and focus content investment on {top_category}."
        ),
        "strengths": [
            f"Revenue is concentrated in identifiable hero ASINs, making prioritization practical.",
            f"Top category momentum is visible in {top_category}.",
            f"Average rating is {format_number(metrics.get('weighted_average_rating'), 2)} where ratings are present.",
        ],
        "weaknesses": [
            f"{format_percent(inactive_pct)} of the 2025 catalog is inactive.",
            f"{metrics.get('declining_asins_count', 0)} comparable active ASINs declined year over year.",
            "Listing content and reseller control require targeted cleanup where source columns show gaps.",
        ],
        "approach_pillars": [
            {
                "title": "Cleanup & Advanced SEO",
                "body": f"Prioritize the top revenue ASINs and fix content gaps across the active {brand} catalog.",
            },
            {
                "title": "Brand Protection",
                "body": "Monitor seller fragmentation, Buy Box ownership and pricing pressure where seller data is present.",
            },
            {
                "title": "Integrated ADV Strategy",
                "body": f"Concentrate campaigns on {top_category} and defend branded demand before scaling broader coverage.",
            },
        ],
        "roadmap": [
            {
                "phase": "01",
                "title": "Audit & Quick Wins - Months 1-3",
                "body": "Catalog audit, SEO cleanup, image fixes on priority ASINs, reseller review and initial campaigns.",
            },
            {
                "phase": "02",
                "title": "Scaling & Optimization - Months 4-8",
                "body": "A+ content rollout, selective inactive ASIN reactivation, campaign expansion and Buy Box monitoring.",
            },
            {
                "phase": "03",
                "title": "Consolidation & Growth - Months 9-12",
                "body": "Brand Store optimization, bundles or variants where supported, monthly KPI reporting and ROI tuning.",
            },
        ],
        "conclusions": {
            "current_situation": [
                f"Revenue 2025: {format_currency(metrics.get('total_revenue_2025'))}",
                f"Active ASINs: {active} out of {total}",
            ],
            "strengths": [
                f"Top category: {top_category}",
                f"Top 5 ASINs drive {format_share(metrics.get('top_5_revenue_share'))} of revenue",
            ],
            "plan": [
                "Use a 3-phase 12-month operating plan",
                f"Illustrative scenario range: {format_currency((metrics.get('growth_projection_scenarios') or {}).get('realistic', {}).get('revenue_low'))} - {format_currency((metrics.get('growth_projection_scenarios') or {}).get('realistic', {}).get('revenue_high'))}",
            ],
            "urgency": [
                f"Inactive catalog: {format_percent(inactive_pct)}",
                f"Largest declining subcategory: {decline.get('subcategory', 'N/A')}",
            ],
        },
    }
    if not metrics.get("rules", {}).get("can_mention_vine", False):
        narrative = _remove_vine_mentions(narrative)
    return narrative


def build_priority_actions(metrics: dict[str, Any], language: str = "en") -> list[str]:
    """Brand-specific priority actions derived from the deck's real metrics.

    Replaces the old hard-coded per-scenario playbook (identical for every brand)
    with actions grounded in this brand's actual gaps. Only actions whose
    underlying metric is present and non-zero are emitted, so the list is never
    boilerplate; the largest gaps come first.
    """
    it = language == "it"
    content = metrics.get("content_health") or {}
    candidates: list[tuple[int, str]] = []

    def add(count: Any, en_text: str, it_text: str) -> None:
        n = int(count or 0)
        if n > 0:
            candidates.append((n, (it_text if it else en_text).format(n=n)))

    add(metrics.get("inactive_asins_2025"),
        "Reactivate {n} inactive ASINs", "Riattivare {n} ASIN inattivi")
    add(metrics.get("asins_with_fewer_than_5_images"),
        "Add images to {n} ASINs (under 5)", "Aggiungere immagini a {n} ASIN (meno di 5)")
    add(content.get("asins_missing_aplus_content"),
        "Add A+ content to {n} ASINs", "Aggiungere contenuti A+ a {n} ASIN")
    add(content.get("short_title_count"),
        "Improve {n} short titles", "Migliorare {n} titoli troppo corti")
    add(content.get("asins_missing_description"),
        "Add descriptions to {n} ASINs", "Aggiungere descrizioni a {n} ASIN")
    add(metrics.get("declining_asins_count"),
        "Review {n} declining ASINs", "Rivedere {n} ASIN in calo")
    add(metrics.get("asins_with_fewer_than_15_reviews"),
        "Build reviews on {n} ASINs (under 15)", "Aumentare le recensioni su {n} ASIN (meno di 15)")

    candidates.sort(key=lambda item: item[0], reverse=True)
    actions = [text for _, text in candidates[:6]]
    if not actions:
        actions = [
            "Mantenere copertura di catalogo e qualità dei contenuti"
            if it
            else "Maintain catalog coverage and content quality"
        ]
    return actions


class BrandAnalysisNarrativeService:
    """Generate strategic narrative with Anthropic, with deterministic fallback."""

    def __init__(self, api_key: Optional[str] = None) -> None:
        self.api_key = api_key or settings.ANTHROPIC_API_KEY

    def generate(
        self,
        metrics: dict[str, Any],
        language: str = "en",
        *,
        provenance: Optional[dict[str, Any]] = None,
        limitations: Optional[dict[str, Any]] = None,
    ) -> dict[str, Any]:
        if not self.api_key:
            return build_fallback_narrative(metrics, language)

        try:
            import anthropic

            client = anthropic.Anthropic(api_key=self.api_key)
            lang_instruction = "Respond entirely in Italian." if language == "it" else "Respond entirely in English."
            vine_rule = (
                "You may mention Amazon Vine only if can_mention_vine is true. "
                "If false, do not mention Amazon Vine anywhere."
            )
            prompt = f"""You are an Amazon marketplace strategy analyst. Create narrative text for a Brand Analysis deck using internal template {NARRATIVE_TEMPLATE_VERSION}.

{lang_instruction}

Metrics are already calculated. Do not calculate, infer, invent, round, or change numbers.
Use only the deterministic metrics in this JSON:
{json.dumps(metrics, ensure_ascii=False, indent=2)}

Metric provenance:
{json.dumps(provenance or {}, ensure_ascii=False, indent=2)}

Limitations and N/A reasons:
{json.dumps(limitations or {}, ensure_ascii=False, indent=2)}

Never convert search proxy metrics into revenue market share. Never infer historical Buy Box owner from a current snapshot. Never invent competitor revenue or market size.

{vine_rule}

Return ONLY valid JSON with exactly this structure:
{{
  "overview": "2 sentence current-state overview",
  "strengths": ["strength 1", "strength 2", "strength 3"],
  "weaknesses": ["weakness 1", "weakness 2", "weakness 3"],
  "approach_pillars": [
    {{"title": "Cleanup & Advanced SEO", "body": "specific body"}},
    {{"title": "Brand Protection", "body": "specific body"}},
    {{"title": "Integrated ADV Strategy", "body": "specific body"}}
  ],
  "roadmap": [
    {{"phase": "01", "title": "Audit & Quick Wins - Months 1-3", "body": "specific actions"}},
    {{"phase": "02", "title": "Scaling & Optimization - Months 4-8", "body": "specific actions"}},
    {{"phase": "03", "title": "Consolidation & Growth - Months 9-12", "body": "specific actions"}}
  ],
  "conclusions": {{
    "current_situation": ["bullet", "bullet"],
    "strengths": ["bullet", "bullet"],
    "plan": ["bullet", "bullet"],
    "urgency": ["bullet", "bullet"]
  }}
}}"""
            message = client.messages.create(
                model="claude-sonnet-4-6",
                max_tokens=2200,
                messages=[{"role": "user", "content": prompt}],
            )
            response_text = message.content[0].text.strip()
            if response_text.startswith("```"):
                lines = response_text.splitlines()
                lines = lines[1:-1] if lines and lines[-1].strip() == "```" else lines[1:]
                response_text = "\n".join(lines)
            parsed = json.loads(response_text)
            narrative = self._validate(parsed)
        except Exception:
            logger.exception("Brand analysis AI narrative failed; using fallback")
            narrative = build_fallback_narrative(metrics, language)

        if not metrics.get("rules", {}).get("can_mention_vine", False):
            narrative = _remove_vine_mentions(narrative)
        return narrative

    @staticmethod
    def _validate(data: dict[str, Any]) -> dict[str, Any]:
        required = {"overview", "strengths", "weaknesses", "approach_pillars", "roadmap", "conclusions"}
        missing = required - set(data)
        if missing:
            raise ValueError(f"AI response missing keys: {missing}")
        if not isinstance(data["approach_pillars"], list) or len(data["approach_pillars"]) < 3:
            raise ValueError("AI response approach_pillars must contain at least 3 items")
        if not isinstance(data["roadmap"], list) or len(data["roadmap"]) < 3:
            raise ValueError("AI response roadmap must contain at least 3 items")
        if not isinstance(data["conclusions"], dict):
            raise ValueError("AI response conclusions must be an object")
        return data


def format_currency(value: Optional[float], digits: int = 0) -> str:
    if value is None:
        return "N/A"
    return f"EUR {float(value):,.{digits}f}"


def format_percent(value: Optional[float]) -> str:
    if value is None:
        return "New"
    return f"{float(value):+.1f}%"


def format_share(value: Optional[float]) -> str:
    if value is None:
        return "N/A"
    return f"{float(value):.1f}%"


def format_number(value: Optional[float], digits: int = 0) -> str:
    if value is None:
        return "N/A"
    return f"{float(value):,.{digits}f}"


PPTX_STATIC_STRINGS: dict[str, dict[str, str]] = {
    "en": {
        # Cover
        "cover_on_amazon": "ON AMAZON",
        "cover_subtitle": "Market analysis and growth strategy",
        # As-is
        "as_is_title": "Current Amazon Performance",
        "as_is_subtitle": "Overview of current Amazon marketplace presence",
        "kpi_revenue_2025": "Revenue 2025",
        "kpi_revenue_2024": "Revenue 2024",
        "kpi_yoy_change": "YoY Change",
        "kpi_average_rating": "Average Rating",
        "kpi_units_sold_2025": "Units Sold 2025",
        "kpi_avg_price_per_asin": "Avg Price / ASIN",
        # Revenue YoY
        "revenue_yoy_title": "Revenue 2024 vs 2025",
        "revenue_yoy_subtitle": "Year-over-year performance calculated from source rows",
        "kpi_yoy": "YoY",
        "revenue_yoy_footnote": "All revenue numbers are deterministic sums from the selected internal or uploaded yearly source data.",
        # Catalog health
        "catalog_health_title": "Catalog Health",
        "catalog_health_subtitle": "ASIN universe, discovery coverage and enrichment quality",
        "kpi_asins_2025": "ASINs 2025",
        "kpi_active_2025": "Active 2025",
        "kpi_inactive_2025": "Inactive 2025",
        "kpi_new_yoy": "New YoY",
        "readiness_discovered_asins": "Discovered ASINs",
        "readiness_lookups_attempted": "Catalog lookups attempted",
        "readiness_failed_asins": "Catalog failed ASINs",
        "readiness_partial_enrichment": "Partial enrichment",
        "value_yes": "Yes",
        "value_no": "No",
        "value_none": "None",
        "table_readiness_signal": "Readiness signal",
        "table_value": "Value",
        "missing_optional_fields_2025": "Optional fields missing in 2025",
        "source_limitations": "Source limitations",
        # Active / inactive
        "active_inactive_title": "Active / Inactive ASINs",
        "active_inactive_subtitle": "Zero-revenue ASINs are treated as inactive",
        "catalog_split_2025": "2025 catalog split",
        "kpi_active": "Active",
        "kpi_inactive": "Inactive",
        "kpi_pct_inactive": "% Inactive",
        "active_inactive_footnote": "ASINs discovered through catalog search or local catalog enrichment are included with EUR 0 revenue when the account has yearly history but the scoped ASIN has no sales.",
        # Top performers
        "top_performers_title": "Top Performing ASINs",
        "top_performers_subtitle": "Revenue leaders and YoY movement",
        "table_asin": "ASIN",
        "table_product": "Product",
        "table_rev_2025": "Rev. 2025",
        "table_rev_2024": "Rev. 2024",
        "table_yoy_pct": "YoY %",
        # Catalog audit (alternate layout)
        "catalog_audit_title": "Catalog Audit",
        "catalog_audit_subtitle": "Catalog composition and performance on Amazon",
        "kpi_asins_2024": "ASINs 2024",
        "kpi_inactive_asins": "Inactive ASINs",
        # Content audit
        "content_audit_title": "SEO & Content Audit",
        "content_audit_subtitle": "Content quality and product listing optimization",
        "kpi_avg_images_per_asin": "Avg Images / ASIN",
        "kpi_missing_bullets": "Missing Bullets",
        "kpi_missing_description": "Missing Description",
        "kpi_short_titles": "Short Titles",
        "table_detected_gaps": "Detected gaps",
        "content_no_gaps": "No source-backed content gaps",
        # Image / review weaknesses
        "review_image_title": "Image / Review Weaknesses",
        "review_image_subtitle": "Listing trust signals where source data is available",
        "kpi_asins_few_images": "ASINs <5 Images",
        "kpi_asins_few_reviews": "ASINs <15 Reviews",
        "kpi_rating_below_4": "Rating <4.0",
        "table_reviews": "Reviews",
        "table_rating": "Rating",
        "table_issue": "Issue",
        "review_no_weakness": "No review/rating weakness in source-backed fields",
        # Subcategory performance
        "subcategory_title": "Subcategory Performance",
        "subcategory_subtitle": "Revenue contribution and YoY by subcategory",
        "table_subcategory": "Subcategory",
        # Operational gap
        "operational_gap_title": "Operational Performance Gap",
        "operational_gap_subtitle": "Key issues limiting growth potential",
        "kpi_pct_inactive_asins": "% Inactive ASINs",
        "kpi_pct_declining_asins": "% ASINs Declining YoY",
        "kpi_asins_multi_seller": "ASINs with >1 Seller",
        "kpi_largest_subcat_decline": "Largest Subcategory Decline",
        "revenue_concentration": "Revenue Concentration",
        "kpi_top_5_asins": "Top 5 ASINs",
        "kpi_top_10_asins": "Top 10 ASINs",
        "kpi_avg_rev_per_active_asin": "Avg Rev / Active ASIN",
        # Channel gap
        "channel_gap_title": "Channel Performance Gap",
        "channel_gap_subtitle": "Current reseller snapshot and Buy Box availability",
        "kpi_avg_sellers": "Avg Sellers",
        "kpi_avg_offers": "Avg Offers",
        "kpi_missing_buy_box": "Missing Buy Box",
        "table_reseller_buy_box": "Reseller / Buy Box",
        "table_asins": "ASINs",
        "table_revenue": "Revenue",
        "table_pct_impact": "% Impact",
        "current_snapshot": "Current snapshot",
        "channel_not_available": "Not available in source data",
        # Concentration risk
        "concentration_risk_title": "Concentration Risk",
        "concentration_risk_subtitle": "How much revenue depends on the top ASINs",
        "kpi_top_5_revenue_share": "Top 5 Revenue Share",
        "kpi_top_10_revenue_share": "Top 10 Revenue Share",
        "table_yoy": "YoY",
        # Market share
        "market_share_title": "Market Share & Competition",
        "market_share_subtitle": "Calculated only when a reliable competitor revenue base exists",
        "kpi_market_size_2025": "Market Size 2025",
        "kpi_revenue_share_2025": "Revenue Share 2025",
        "kpi_revenue_share_2024": "Revenue Share 2024",
        "table_brand": "Brand",
        "table_share": "Share",
        "market_share_unavailable": "Revenue market share: N/A [UNAVAILABLE]",
        "market_share_no_base": "No reliable competitor revenue base was provided.",
        "market_share_no_invent": "The deck does not invent market size. Upload a broad yearly market export with competitor brand revenue to calculate this section.",
        # Approach
        "approach_title": "Our Approach",
        "approach_subtitle": "Three pillars for professional Amazon channel management",
        "approach_visibility": "VISIBILITY",
        "approach_visibility_sub": "SEO - ADV - Keywords",
        "approach_conversion": "CONVERSION",
        "approach_conversion_sub": "Content - Images - A+",
        "approach_loyalty": "LOYALTY",
        "approach_loyalty_sub": "Brand Store - Reviews",
        # Roadmap
        "roadmap_title": "Operational Roadmap",
        "roadmap_subtitle": "3-phase action plan for the first 12 months",
        # Projection
        "projection_title": "Growth Scenarios",
        "projection_subtitle": "Illustrative 12-month revenue ranges, not a forecast",
        "current_situation": "Current Situation",
        "projection_active_asins": "active ASINs out of",
        "scenario_conservative": "CONSERVATIVE",
        "scenario_realistic": "REALISTIC",
        "scenario_optimistic": "OPTIMISTIC",
        "projection_actions_title": "Priority actions for this brand",
        "projection_disclaimer": "Illustrative scenarios based on uniform growth assumptions, not a forecast. Adjust with brand-specific targets.",
        # Conclusions
        "conclusions_title": "Conclusions",
        "conclusions_subtitle": "Summary and next steps",
        "conclusions_current_situation": "Current Situation",
        "conclusions_strengths": "Strengths",
        "conclusions_plan": "3-Phase Plan - 12 Months",
        "conclusions_urgency": "Urgency",
    },
    "it": {
        # Cover
        "cover_on_amazon": "SU AMAZON",
        "cover_subtitle": "Analisi di mercato e strategia di crescita",
        # As-is
        "as_is_title": "Performance attuale su Amazon",
        "as_is_subtitle": "Panoramica della presenza attuale sul marketplace Amazon",
        "kpi_revenue_2025": "Fatturato 2025",
        "kpi_revenue_2024": "Fatturato 2024",
        "kpi_yoy_change": "Variazione YoY",
        "kpi_average_rating": "Valutazione media",
        "kpi_units_sold_2025": "Unità vendute 2025",
        "kpi_avg_price_per_asin": "Prezzo medio / ASIN",
        # Revenue YoY
        "revenue_yoy_title": "Fatturato 2024 vs 2025",
        "revenue_yoy_subtitle": "Performance anno su anno calcolata dai dati di origine",
        "kpi_yoy": "YoY",
        "revenue_yoy_footnote": "Tutti i valori di fatturato sono somme deterministiche dei dati annuali di origine interni o caricati selezionati.",
        # Catalog health
        "catalog_health_title": "Salute del catalogo",
        "catalog_health_subtitle": "Universo ASIN, copertura della discovery e qualità dell'arricchimento",
        "kpi_asins_2025": "ASIN 2025",
        "kpi_active_2025": "Attivi 2025",
        "kpi_inactive_2025": "Inattivi 2025",
        "kpi_new_yoy": "Nuovi YoY",
        "readiness_discovered_asins": "ASIN individuati",
        "readiness_lookups_attempted": "Ricerche catalogo tentate",
        "readiness_failed_asins": "ASIN catalogo non riusciti",
        "readiness_partial_enrichment": "Arricchimento parziale",
        "value_yes": "Sì",
        "value_no": "No",
        "value_none": "Nessuno",
        "table_readiness_signal": "Indicatore di completezza",
        "table_value": "Valore",
        "missing_optional_fields_2025": "Campi opzionali mancanti nel 2025",
        "source_limitations": "Limiti dei dati di origine",
        # Active / inactive
        "active_inactive_title": "ASIN attivi / inattivi",
        "active_inactive_subtitle": "Gli ASIN con fatturato zero sono trattati come inattivi",
        "catalog_split_2025": "Suddivisione catalogo 2025",
        "kpi_active": "Attivi",
        "kpi_inactive": "Inattivi",
        "kpi_pct_inactive": "% Inattivi",
        "active_inactive_footnote": "Gli ASIN individuati tramite ricerca a catalogo o arricchimento del catalogo locale sono inclusi con fatturato pari a EUR 0 quando l'account ha uno storico annuale ma l'ASIN considerato non ha vendite.",
        # Top performers
        "top_performers_title": "ASIN più performanti",
        "top_performers_subtitle": "Leader di fatturato e andamento YoY",
        "table_asin": "ASIN",
        "table_product": "Prodotto",
        "table_rev_2025": "Fatt. 2025",
        "table_rev_2024": "Fatt. 2024",
        "table_yoy_pct": "YoY %",
        # Catalog audit (alternate layout)
        "catalog_audit_title": "Audit del catalogo",
        "catalog_audit_subtitle": "Composizione del catalogo e performance su Amazon",
        "kpi_asins_2024": "ASIN 2024",
        "kpi_inactive_asins": "ASIN inattivi",
        # Content audit
        "content_audit_title": "Audit SEO e contenuti",
        "content_audit_subtitle": "Qualità dei contenuti e ottimizzazione delle schede prodotto",
        "kpi_avg_images_per_asin": "Immagini medie / ASIN",
        "kpi_missing_bullets": "Bullet mancanti",
        "kpi_missing_description": "Descrizione mancante",
        "kpi_short_titles": "Titoli brevi",
        "table_detected_gaps": "Lacune rilevate",
        "content_no_gaps": "Nessuna lacuna di contenuto supportata dai dati",
        # Image / review weaknesses
        "review_image_title": "Debolezze immagini / recensioni",
        "review_image_subtitle": "Segnali di affidabilità delle schede dove i dati di origine sono disponibili",
        "kpi_asins_few_images": "ASIN <5 immagini",
        "kpi_asins_few_reviews": "ASIN <15 recensioni",
        "kpi_rating_below_4": "Valutazione <4.0",
        "table_reviews": "Recensioni",
        "table_rating": "Valutazione",
        "table_issue": "Problema",
        "review_no_weakness": "Nessuna debolezza di recensioni/valutazioni nei campi supportati dai dati",
        # Subcategory performance
        "subcategory_title": "Performance per sottocategoria",
        "subcategory_subtitle": "Contributo al fatturato e YoY per sottocategoria",
        "table_subcategory": "Sottocategoria",
        # Operational gap
        "operational_gap_title": "Gap di performance operativa",
        "operational_gap_subtitle": "Principali criticità che limitano il potenziale di crescita",
        "kpi_pct_inactive_asins": "% ASIN inattivi",
        "kpi_pct_declining_asins": "% ASIN in calo YoY",
        "kpi_asins_multi_seller": "ASIN con >1 venditore",
        "kpi_largest_subcat_decline": "Maggior calo per sottocategoria",
        "revenue_concentration": "Concentrazione del fatturato",
        "kpi_top_5_asins": "Top 5 ASIN",
        "kpi_top_10_asins": "Top 10 ASIN",
        "kpi_avg_rev_per_active_asin": "Fatt. medio / ASIN attivo",
        # Channel gap
        "channel_gap_title": "Gap di performance di canale",
        "channel_gap_subtitle": "Snapshot attuale dei rivenditori e disponibilità della Buy Box",
        "kpi_avg_sellers": "Venditori medi",
        "kpi_avg_offers": "Offerte medie",
        "kpi_missing_buy_box": "Buy Box mancante",
        "table_reseller_buy_box": "Rivenditore / Buy Box",
        "table_asins": "ASIN",
        "table_revenue": "Fatturato",
        "table_pct_impact": "% Impatto",
        "current_snapshot": "Snapshot attuale",
        "channel_not_available": "Non disponibile nei dati di origine",
        # Concentration risk
        "concentration_risk_title": "Rischio di concentrazione",
        "concentration_risk_subtitle": "Quanto fatturato dipende dagli ASIN principali",
        "kpi_top_5_revenue_share": "Quota fatturato Top 5",
        "kpi_top_10_revenue_share": "Quota fatturato Top 10",
        "table_yoy": "YoY",
        # Market share
        "market_share_title": "Quota di mercato e concorrenza",
        "market_share_subtitle": "Calcolata solo quando esiste una base affidabile di fatturato dei concorrenti",
        "kpi_market_size_2025": "Dimensione mercato 2025",
        "kpi_revenue_share_2025": "Quota fatturato 2025",
        "kpi_revenue_share_2024": "Quota fatturato 2024",
        "table_brand": "Brand",
        "table_share": "Quota",
        "market_share_unavailable": "Quota di mercato a fatturato: N/A [NON DISPONIBILE]",
        "market_share_no_base": "Non è stata fornita una base affidabile di fatturato dei concorrenti.",
        "market_share_no_invent": "Il deck non inventa la dimensione del mercato. Carica un export di mercato annuale completo con il fatturato dei brand concorrenti per calcolare questa sezione.",
        # Approach
        "approach_title": "Il nostro approccio",
        "approach_subtitle": "Tre pilastri per una gestione professionale del canale Amazon",
        "approach_visibility": "VISIBILITÀ",
        "approach_visibility_sub": "SEO - ADV - Keyword",
        "approach_conversion": "CONVERSIONE",
        "approach_conversion_sub": "Contenuti - Immagini - A+",
        "approach_loyalty": "FIDELIZZAZIONE",
        "approach_loyalty_sub": "Brand Store - Recensioni",
        # Roadmap
        "roadmap_title": "Roadmap operativa",
        "roadmap_subtitle": "Piano d'azione in 3 fasi per i primi 12 mesi",
        # Projection
        "projection_title": "Scenari di crescita",
        "projection_subtitle": "Intervalli di fatturato illustrativi a 12 mesi, non una previsione",
        "current_situation": "Situazione attuale",
        "projection_active_asins": "ASIN attivi su",
        "scenario_conservative": "CONSERVATIVO",
        "scenario_realistic": "REALISTICO",
        "scenario_optimistic": "OTTIMISTICO",
        "projection_actions_title": "Azioni prioritarie per questa marca",
        "projection_disclaimer": "Scenari illustrativi basati su ipotesi di crescita uniforme, non una previsione. Da adattare con obiettivi specifici della marca.",
        # Conclusions
        "conclusions_title": "Conclusioni",
        "conclusions_subtitle": "Sintesi e prossimi passi",
        "conclusions_current_situation": "Situazione attuale",
        "conclusions_strengths": "Punti di forza",
        "conclusions_plan": "Piano in 3 fasi - 12 mesi",
        "conclusions_urgency": "Urgenza",
    },
}


class BrandAnalysisPptxBuilder:
    """Generate a PowerPoint deck from deterministic metrics."""

    def __init__(
        self,
        metrics: dict[str, Any],
        narrative: dict[str, Any],
        language: str = "en",
    ) -> None:
        self.metrics = metrics
        self.narrative = narrative
        self.brand = str(metrics.get("brand_name") or "Brand").upper()
        self.language = "it" if str(language or "").lower().startswith("it") else "en"

    def _t(self, key: str) -> str:
        strings = PPTX_STATIC_STRINGS.get(self.language) or PPTX_STATIC_STRINGS["en"]
        return strings.get(key) or PPTX_STATIC_STRINGS["en"].get(key, key)

    def _badge(self, metric_key: str) -> str:
        registry = self.metrics.get("metric_source_registry") or {}
        quality = (registry.get(metric_key) or {}).get("quality")
        if not quality:
            return ""
        return f" [{str(quality).upper()}]"

    def _has_market_share(self) -> bool:
        return (self.metrics.get("market_analysis") or {}).get("status") == "calculated_from_external_market_export"

    def _has_channel_data(self) -> bool:
        bb = self.metrics.get("seller_buy_box_summary") or {}
        return any(bb.get(k) for k in ("seller_count_available", "offer_count_available", "buy_box_owner_available", "current_snapshot_available"))

    def build(self) -> bytes:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.util import Inches, Pt

        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)

        self._RGBColor = RGBColor
        self._PP_ALIGN = PP_ALIGN
        self._MSO_ANCHOR = MSO_ANCHOR
        self._Inches = Inches
        self._Pt = Pt

        self._slide_cover(prs)
        slides = [
            (self._slide_as_is, True),
            (self._slide_revenue_yoy, True),
            (self._slide_catalog_health, True),
            (self._slide_active_inactive, True),
            (self._slide_top_performers, True),
            (self._slide_content_audit, True),
            (self._slide_review_image_weaknesses, True),
            (self._slide_subcategory_performance, True),
            (self._slide_operational_gap, True),
            (self._slide_channel_gap, self._has_channel_data()),
            (self._slide_concentration_risk, True),
            (self._slide_market_share, self._has_market_share()),
            (self._slide_projection, True),
            (self._slide_roadmap, True),
            (self._slide_conclusions, True),
        ]
        page = 2
        for method, include in slides:
            if not include:
                continue
            method(prs, page)
            page += 1

        output = io.BytesIO()
        prs.save(output)
        return output.getvalue()

    def _blank(self, prs):
        return prs.slides.add_slide(prs.slide_layouts[6])

    def _add_header(self, slide, page: int) -> None:
        red = self._RGBColor(212, 39, 45)
        self._rect(slide, 0, 0, 10, 0.34, red)
        self._text(slide, 0.35, 0.09, 2.8, 0.18, self.brand, size=8, bold=True, color=(255, 255, 255))
        self._text(slide, 9.38, 0.08, 0.25, 0.18, str(page), size=8, bold=True, color=(255, 255, 255), align="right")
        self._footer(slide)

    def _footer(self, slide) -> None:
        colors = [(29, 78, 216), (234, 88, 12), (22, 163, 74), (220, 38, 38), (245, 158, 11), (14, 165, 233), (100, 116, 139)]
        x = 0
        for r, g, b in colors:
            self._rect(slide, x, 5.49, 10 / len(colors), 0.08, self._RGBColor(r, g, b))
            x += 10 / len(colors)

    def _slide_cover(self, prs) -> None:
        slide = self._blank(prs)
        self._rect(slide, 0, 0, 10, 5.625, self._RGBColor(212, 39, 45))
        self._text(slide, 1.2, 2.25, 7.6, 0.55, f"{self.brand} {self._t('cover_on_amazon')}", size=30, bold=True, color=(255, 255, 255), align="center")
        self._text(slide, 2.7, 2.92, 4.6, 0.3, self._t("cover_subtitle"), size=13, color=(255, 255, 255), align="center")
        self._footer(slide)

    def _slide_as_is(self, prs, page: int) -> None:
        slide = self._blank(prs)
        self._add_header(slide, page)
        self._title(slide, self._t("as_is_title"), self._t("as_is_subtitle"))
        kpis = [
            (self._t("kpi_revenue_2025") + self._badge("total_revenue_2025"), format_currency(self.metrics.get("total_revenue_2025")), 0.65, 1.4, 2.3, 0.95),
            (self._t("kpi_revenue_2024") + self._badge("total_revenue_2024"), format_currency(self.metrics.get("total_revenue_2024")), 3.05, 1.4, 2.3, 0.95),
            (self._t("kpi_yoy_change") + self._badge("yoy_percent"), format_percent(self.metrics.get("yoy_percent")), 5.45, 1.4, 1.65, 0.95),
            (self._t("kpi_average_rating") + self._badge("weighted_average_rating"), format_number(self.metrics.get("weighted_average_rating"), 2), 7.25, 1.4, 1.65, 0.95),
            (self._t("kpi_units_sold_2025") + self._badge("total_units_sold_2025"), format_number(self.metrics.get("total_units_sold_2025"), 0), 5.45, 2.55, 1.65, 0.95),
            (self._t("kpi_avg_price_per_asin") + self._badge("average_price_per_asin"), format_currency(self.metrics.get("average_price_per_asin"), 2), 7.25, 2.55, 1.65, 0.95),
        ]
        for label, value, x, y, w, h in kpis:
            self._kpi(slide, x, y, w, h, label, value)
        self._body_box(slide, 0.65, 3.75, 8.4, 0.75, self.narrative.get("overview", ""))

    def _slide_revenue_yoy(self, prs, page: int) -> None:
        slide = self._blank(prs)
        self._add_header(slide, page)
        self._title(slide, self._t("revenue_yoy_title"), self._t("revenue_yoy_subtitle"))
        bars = [
            ("2024", self.metrics.get("total_revenue_2024"), (100, 116, 139)),
            ("2025", self.metrics.get("total_revenue_2025"), (212, 39, 45)),
        ]
        max_value = max(float(value or 0) for _, value, _ in bars) or 1
        for idx, (label, value, color) in enumerate(bars):
            x = 1.0 + idx * 4.25
            bar_h = 2.4 * (float(value or 0) / max_value)
            self._text(slide, x, 1.25, 2.7, 0.24, label, size=12, bold=True, align="center")
            self._rect(slide, x + 0.35, 4.1 - bar_h, 2.0, bar_h, self._RGBColor(*color))
            self._text(slide, x, 4.25, 2.7, 0.28, format_currency(value), size=16, bold=True, align="center", color=color)
        self._kpi(slide, 3.95, 1.35, 1.9, 0.9, self._t("kpi_yoy"), format_percent(self.metrics.get("yoy_percent")))
        self._body_box(
            slide,
            0.8,
            4.82,
            8.4,
            0.42,
            self._t("revenue_yoy_footnote"),
        )

    def _slide_catalog_health(self, prs, page: int) -> None:
        slide = self._blank(prs)
        self._add_header(slide, page)
        self._title(slide, self._t("catalog_health_title"), self._t("catalog_health_subtitle"))
        kpis = [
            (self._t("kpi_asins_2025"), str(self.metrics.get("total_asins_2025", 0))),
            (self._t("kpi_active_2025"), str(self.metrics.get("active_asins_2025", 0))),
            (self._t("kpi_inactive_2025"), str(self.metrics.get("inactive_asins_2025", 0))),
            (self._t("kpi_new_yoy"), str(self.metrics.get("new_asins_yoy", 0))),
        ]
        for idx, (label, value) in enumerate(kpis):
            self._kpi(slide, 0.65 + idx * 2.2, 1.25, 1.85, 0.78, label, value)
        readiness = self.metrics.get("data_readiness") or {}
        catalog = readiness.get("catalog_enrichment") or {}
        rows = [
            [self._t("readiness_discovered_asins"), str(readiness.get("discovered_asins_count", "N/A"))],
            [self._t("readiness_lookups_attempted"), str(catalog.get("attempted", "N/A"))],
            [self._t("readiness_failed_asins"), str(len(catalog.get("failed_asins") or []))],
            [self._t("readiness_partial_enrichment"), self._t("value_yes") if catalog.get("partial") else self._t("value_no")],
        ]
        self._table(slide, 0.8, 2.35, 4.0, 2.25, [self._t("table_readiness_signal"), self._t("table_value")], rows, [2.6, 1.1])
        completeness = self.metrics.get("data_completeness") or {}
        missing = completeness.get("missing_optional_fields_2025") or []
        limitations = (self.metrics.get("limitations") or {}).get("items") or []
        none_bullet = f"- {self._t('value_none')}"
        self._body_box(
            slide,
            5.15,
            2.35,
            3.95,
            1.45,
            f"{self._t('missing_optional_fields_2025')}\n" + ("\n".join(f"- {item}" for item in missing[:8]) if missing else none_bullet),
            title_bold=True,
        )
        self._body_box(
            slide,
            5.15,
            3.95,
            3.95,
            0.85,
            f"{self._t('source_limitations')}\n" + ("\n".join(f"- {item.get('area')}" for item in limitations[:3]) if limitations else none_bullet),
            title_bold=True,
        )

    def _slide_active_inactive(self, prs, page: int) -> None:
        slide = self._blank(prs)
        self._add_header(slide, page)
        self._title(slide, self._t("active_inactive_title"), self._t("active_inactive_subtitle"))
        active = int(self.metrics.get("active_asins_2025") or 0)
        inactive = int(self.metrics.get("inactive_asins_2025") or 0)
        total = max(active + inactive, 1)
        active_w = 7.8 * active / total
        inactive_w = 7.8 - active_w
        self._rect(slide, 1.1, 1.7, active_w, 0.72, self._RGBColor(22, 163, 74))
        self._rect(slide, 1.1 + active_w, 1.7, inactive_w, 0.72, self._RGBColor(212, 39, 45))
        self._text(slide, 1.1, 1.35, 2.5, 0.22, self._t("catalog_split_2025"), size=12, bold=True)
        self._kpi(slide, 1.1, 2.75, 2.4, 0.86, self._t("kpi_active"), str(active))
        self._kpi(slide, 3.8, 2.75, 2.4, 0.86, self._t("kpi_inactive"), str(inactive))
        self._kpi(slide, 6.5, 2.75, 2.4, 0.86, self._t("kpi_pct_inactive"), format_share(self.metrics.get("percentage_inactive_asins")))
        self._body_box(
            slide,
            1.1,
            4.05,
            7.8,
            0.72,
            self._t("active_inactive_footnote"),
        )

    def _slide_top_performers(self, prs, page: int) -> None:
        slide = self._blank(prs)
        self._add_header(slide, page)
        self._title(slide, self._t("top_performers_title"), self._t("top_performers_subtitle"))
        rows = [
            [item["asin"], item["product_name"], format_currency(item["revenue_2025"]), format_currency(item["revenue_2024"]), format_percent(item["yoy_percent"])]
            for item in self.metrics.get("top_5_asins", [])
        ]
        self._table(
            slide, 0.65, 1.25, 8.7, 3.65,
            [self._t("table_asin"), self._t("table_product"), self._t("table_rev_2025"), self._t("table_rev_2024"), self._t("table_yoy_pct")], rows, [1.25, 3.45, 1.35, 1.35, 0.95],
            yoy_columns=[4],
        )

    def _slide_catalog_audit(self, prs, page: int) -> None:
        slide = self._blank(prs)
        self._add_header(slide, page)
        self._title(slide, self._t("catalog_audit_title"), self._t("catalog_audit_subtitle"))
        kpis = [
            (self._t("kpi_asins_2025"), str(self.metrics.get("total_asins_2025", 0))),
            (self._t("kpi_asins_2024"), str(self.metrics.get("total_asins_2024", 0))),
            (self._t("kpi_new_yoy"), str(self.metrics.get("new_asins_yoy", 0))),
            (self._t("kpi_inactive_asins"), str(self.metrics.get("inactive_asins_2025", 0))),
        ]
        for idx, (label, value) in enumerate(kpis):
            self._kpi(slide, 0.65 + idx * 2.2, 1.25, 1.85, 0.78, label, value)
        rows = [
            [item["asin"], item["product_name"], format_currency(item["revenue_2025"]), format_percent(item["yoy_percent"])]
            for item in self.metrics.get("top_5_asins", [])
        ]
        self._table(
            slide, 0.65, 2.35, 8.7, 2.55,
            [self._t("table_asin"), self._t("table_product"), self._t("table_rev_2025"), self._t("table_yoy_pct")], rows, [1.35, 4.25, 1.55, 1.0],
            yoy_columns=[3],
        )

    def _slide_content_audit(self, prs, page: int) -> None:
        slide = self._blank(prs)
        self._add_header(slide, page)
        self._title(slide, self._t("content_audit_title"), self._t("content_audit_subtitle"))
        content = self.metrics.get("content_health") or {}
        kpis = [
            (self._t("kpi_avg_images_per_asin"), format_number(self.metrics.get("average_images_per_asin"), 1)),
            (self._t("kpi_missing_bullets"), format_number(content.get("asins_missing_bullets"), 0)),
            (self._t("kpi_missing_description"), format_number(content.get("asins_missing_description"), 0)),
            (self._t("kpi_short_titles"), format_number(content.get("short_title_count"), 0)),
        ]
        for idx, (label, value) in enumerate(kpis):
            self._kpi(slide, 0.65 + idx * 2.25, 1.25, 2.0, 0.78, label, value)
        rows = [
            [item.get("asin"), item.get("product_name"), ", ".join(item.get("issues") or []), format_currency(item.get("revenue_2025"))]
            for item in (content.get("content_gap_asins") or [])[:6]
        ]
        if not rows:
            rows = [["N/A", self._t("content_no_gaps"), "N/A", "N/A"]]
        self._table(
            slide, 0.65, 2.25, 8.7, 2.75,
            [self._t("table_asin"), self._t("table_product"), self._t("table_detected_gaps"), self._t("table_rev_2025")], rows, [1.2, 3.0, 2.8, 1.2],
        )

    def _slide_review_image_weaknesses(self, prs, page: int) -> None:
        slide = self._blank(prs)
        self._add_header(slide, page)
        self._title(slide, self._t("review_image_title"), self._t("review_image_subtitle"))
        weaknesses = self.metrics.get("review_rating_weaknesses") or {}
        kpis = [
            (self._t("kpi_asins_few_images"), format_number(self.metrics.get("asins_with_fewer_than_5_images"), 0)),
            (self._t("kpi_asins_few_reviews"), format_number(weaknesses.get("asins_with_fewer_than_15_reviews"), 0)),
            (self._t("kpi_rating_below_4"), format_number(weaknesses.get("asins_with_rating_below_4"), 0)),
        ]
        for idx, (label, value) in enumerate(kpis):
            self._kpi(slide, 0.85 + idx * 2.65, 1.25, 2.25, 0.82, label, value)
        rows = [
            [item.get("asin"), item.get("product_name"), format_number(item.get("reviews"), 0), format_number(item.get("rating"), 2), ", ".join(item.get("issues") or [])]
            for item in (weaknesses.get("weak_asins") or [])[:7]
        ]
        if not rows:
            rows = [["N/A", self._t("review_no_weakness"), "N/A", "N/A", "N/A"]]
        self._table(
            slide, 0.65, 2.35, 8.7, 2.55,
            [self._t("table_asin"), self._t("table_product"), self._t("table_reviews"), self._t("table_rating"), self._t("table_issue")], rows, [1.15, 3.1, 0.95, 0.85, 2.0],
        )

    def _slide_subcategory_performance(self, prs, page: int) -> None:
        slide = self._blank(prs)
        self._add_header(slide, page)
        self._title(slide, self._t("subcategory_title"), self._t("subcategory_subtitle"))
        rows = [
            [item["subcategory"], format_currency(item["revenue_2025"]), format_currency(item["revenue_2024"]), format_percent(item["yoy_percent"])]
            for item in (self.metrics.get("revenue_by_subcategory") or [])[:8]
        ]
        self._table(
            slide, 0.65, 2.25, 8.7, 2.75,
            [self._t("table_subcategory"), self._t("table_rev_2025"), self._t("table_rev_2024"), self._t("table_yoy_pct")], rows, [3.6, 1.65, 1.65, 1.1],
            yoy_columns=[3],
        )

    def _slide_operational_gap(self, prs, page: int) -> None:
        slide = self._blank(prs)
        self._add_header(slide, page)
        self._title(slide, self._t("operational_gap_title"), self._t("operational_gap_subtitle"))
        decline = self.metrics.get("subcategory_with_largest_decline") or {}
        kpis = [
            (self._t("kpi_pct_inactive_asins"), format_share(self.metrics.get("percentage_inactive_asins"))),
            (self._t("kpi_pct_declining_asins"), format_share(self.metrics.get("percentage_declining_asins_among_active"))),
            (self._t("kpi_asins_multi_seller"), format_number(self.metrics.get("asins_with_more_than_1_seller"), 0)),
            (self._t("kpi_largest_subcat_decline"), f"{decline.get('subcategory', 'N/A')} {format_percent(decline.get('yoy_percent')) if decline else ''}".strip()),
        ]
        for idx, (label, value) in enumerate(kpis):
            x = 0.65 + (idx % 2) * 4.3
            y = 1.25 + (idx // 2) * 1.05
            self._kpi(slide, x, y, 3.8, 0.78, label, value)
        self._text(slide, 0.85, 3.6, 2.8, 0.24, self._t("revenue_concentration"), size=13, bold=True, color=(212, 39, 45))
        concentration = [
            (self._t("kpi_top_5_asins"), format_share(self.metrics.get("top_5_revenue_share"))),
            (self._t("kpi_top_10_asins"), format_share(self.metrics.get("top_10_revenue_share"))),
            (self._t("kpi_avg_rev_per_active_asin"), format_currency(self.metrics.get("average_revenue_per_active_asin"))),
        ]
        for idx, (label, value) in enumerate(concentration):
            self._kpi(slide, 0.85 + idx * 2.75, 3.95, 2.25, 0.72, label, value)

    def _slide_channel_gap(self, prs, page: int) -> None:
        slide = self._blank(prs)
        self._add_header(slide, page)
        self._title(slide, self._t("channel_gap_title"), self._t("channel_gap_subtitle"))
        summary = self.metrics.get("seller_buy_box_summary") or {}
        self._kpi(slide, 0.65, 1.15, 2.0, 0.72, self._t("kpi_avg_sellers"), format_number(summary.get("average_seller_count"), 1))
        self._kpi(slide, 2.9, 1.15, 2.0, 0.72, self._t("kpi_avg_offers"), format_number(summary.get("average_offer_count"), 1))
        self._kpi(slide, 5.15, 1.15, 2.0, 0.72, self._t("kpi_missing_buy_box"), format_number(summary.get("asins_missing_buy_box_owner"), 0))
        rows = [
            [item["reseller"], str(item["asin_count"]), format_currency(item["revenue"]), format_share(item["share_percent"])]
            for item in (self.metrics.get("reseller_buy_box_distribution") or [])[:8]
        ]
        if not rows and summary.get("current_buy_box_snapshot_distribution"):
            rows = [
                [item.get("reseller"), str(item.get("asin_count")), "N/A", self._t("current_snapshot")]
                for item in summary.get("current_buy_box_snapshot_distribution", [])[:8]
            ]
        if not rows:
            rows = [[self._t("channel_not_available"), "N/A", "N/A", "N/A"]]
        self._table(slide, 0.65, 2.15, 8.7, 2.8, [self._t("table_reseller_buy_box"), self._t("table_asins"), self._t("table_revenue"), self._t("table_pct_impact")], rows, [4.2, 1.0, 1.8, 1.2])

    def _slide_concentration_risk(self, prs, page: int) -> None:
        slide = self._blank(prs)
        self._add_header(slide, page)
        self._title(slide, self._t("concentration_risk_title"), self._t("concentration_risk_subtitle"))
        self._kpi(slide, 0.85, 1.25, 2.3, 0.9, self._t("kpi_top_5_revenue_share"), format_share(self.metrics.get("top_5_revenue_share")))
        self._kpi(slide, 3.45, 1.25, 2.3, 0.9, self._t("kpi_top_10_revenue_share"), format_share(self.metrics.get("top_10_revenue_share")))
        self._kpi(slide, 6.05, 1.25, 2.7, 0.9, self._t("kpi_avg_rev_per_active_asin"), format_currency(self.metrics.get("average_revenue_per_active_asin")))
        rows = [
            [item["asin"], item["product_name"], format_currency(item["revenue_2025"]), format_percent(item["yoy_percent"])]
            for item in self.metrics.get("top_5_asins", [])
        ]
        self._table(
            slide, 0.65, 2.55, 8.7, 2.35,
            [self._t("table_asin"), self._t("table_product"), self._t("table_revenue"), self._t("table_yoy")], rows, [1.25, 4.35, 1.45, 0.95],
            yoy_columns=[3],
        )

    def _slide_market_share(self, prs, page: int) -> None:
        slide = self._blank(prs)
        self._add_header(slide, page)
        self._title(slide, self._t("market_share_title"), self._t("market_share_subtitle"))
        market = self.metrics.get("market_analysis") or {}
        if market.get("status") == "calculated_from_external_market_export":
            kpis = [
                (self._t("kpi_market_size_2025") + self._badge("market_revenue_share"), format_currency(market.get("market_size_2025"))),
                (self._t("kpi_revenue_share_2025") + self._badge("market_revenue_share"), format_share(market.get("market_share_2025"))),
                (self._t("kpi_revenue_share_2024") + self._badge("market_revenue_share"), format_share(market.get("market_share_2024"))),
            ]
            for idx, (label, value) in enumerate(kpis):
                self._kpi(slide, 0.75 + idx * 2.8, 1.25, 2.35, 0.85, label, value)
            rows = [
                [item.get("brand"), str(item.get("asin_count")), format_currency(item.get("revenue")), format_share(item.get("market_share_percent"))]
                for item in (market.get("competitive_brand_distribution") or [])[:7]
            ]
            self._table(slide, 0.65, 2.45, 8.7, 2.45, [self._t("table_brand"), self._t("table_asins"), self._t("table_revenue"), self._t("table_share")], rows, [3.4, 1.0, 1.8, 1.0])
        else:
            self._body_box(
                slide,
                0.85,
                1.45,
                8.3,
                2.2,
                f"{self._t('market_share_unavailable')}\n"
                + str(market.get("limitation") or self._t("market_share_no_base")),
                title_bold=True,
            )
            self._body_box(
                slide,
                0.85,
                3.95,
                8.3,
                0.72,
                self._t("market_share_no_invent"),
            )

    def _slide_approach(self, prs, page: int) -> None:
        slide = self._blank(prs)
        self._add_header(slide, page)
        self._title(slide, self._t("approach_title"), self._t("approach_subtitle"))
        pillars = (self.narrative.get("approach_pillars") or [])[:3]
        for idx, pillar in enumerate(pillars):
            x = 0.65 + idx * 3.1
            self._rect(slide, x, 1.35, 2.75, 2.1, self._RGBColor(25, 35, 50))
            self._rect(slide, x, 1.35, 2.75, 0.12, self._RGBColor(212, 39, 45))
            self._text(slide, x + 0.18, 1.65, 2.35, 0.35, pillar.get("title", ""), size=13, bold=True, color=(255, 255, 255))
            self._text(slide, x + 0.18, 2.12, 2.35, 1.0, pillar.get("body", ""), size=8.5, color=(255, 255, 255))
        banners = [
            (self._t("approach_visibility"), self._t("approach_visibility_sub"), (29, 78, 216)),
            (self._t("approach_conversion"), self._t("approach_conversion_sub"), (234, 88, 12)),
            (self._t("approach_loyalty"), self._t("approach_loyalty_sub"), (22, 163, 74)),
        ]
        for idx, (title, subtitle, color) in enumerate(banners):
            x = 0.65 + idx * 3.1
            self._rect(slide, x, 3.85, 2.75, 0.7, self._RGBColor(*color))
            self._text(slide, x, 3.98, 2.75, 0.18, title, size=10, bold=True, color=(255, 255, 255), align="center")
            self._text(slide, x, 4.25, 2.75, 0.16, subtitle, size=7.5, color=(255, 255, 255), align="center")

    def _slide_roadmap(self, prs, page: int) -> None:
        slide = self._blank(prs)
        self._add_header(slide, page)
        self._title(slide, self._t("roadmap_title"), self._t("roadmap_subtitle"))
        colors = [(212, 39, 45), (234, 88, 12), (22, 163, 74)]
        roadmap = (self.narrative.get("roadmap") or [])[:3]
        for idx, item in enumerate(roadmap):
            y = 1.25 + idx * 1.15
            self._rect(slide, 0.78, y, 0.55, 0.55, self._RGBColor(*colors[idx]))
            self._text(slide, 0.78, y + 0.17, 0.55, 0.14, item.get("phase", f"0{idx+1}"), size=10, bold=True, color=(255, 255, 255), align="center")
            self._text(slide, 1.55, y, 7.5, 0.24, item.get("title", ""), size=13, bold=True, color=(20, 20, 20))
            self._text(slide, 1.55, y + 0.38, 7.5, 0.45, item.get("body", ""), size=9.5, color=(70, 70, 70))

    def _slide_projection(self, prs, page: int) -> None:
        slide = self._blank(prs)
        self._add_header(slide, page)
        self._title(slide, self._t("projection_title"), self._t("projection_subtitle"))
        self._rect(slide, 0.65, 1.18, 8.7, 0.65, self._RGBColor(245, 245, 245))
        self._text(slide, 0.85, 1.27, 2.0, 0.18, self._t("current_situation"), size=9, bold=True, color=(80, 80, 80))
        self._text(slide, 2.55, 1.22, 2.2, 0.35, format_currency(self.metrics.get("total_revenue_2025")), size=22, bold=True, color=(212, 39, 45))
        self._text(
            slide,
            4.75,
            1.34,
            4.0,
            0.2,
            f"{format_percent(self.metrics.get('yoy_percent'))} YoY | {self.metrics.get('active_asins_2025', 0)} {self._t('projection_active_asins')} {self.metrics.get('total_asins_2025', 0)}",
            size=8.5,
            color=(80, 80, 80),
        )
        scenarios = [
            (self._t("scenario_conservative"), "conservative", (100, 116, 139)),
            (self._t("scenario_realistic"), "realistic", (22, 163, 74)),
            (self._t("scenario_optimistic"), "optimistic", (212, 39, 45)),
        ]
        projections = self.metrics.get("growth_projection_scenarios") or {}
        for idx, (label, key, color) in enumerate(scenarios):
            x = 0.65 + idx * 3.0
            scenario = projections.get(key, {})
            self._rect(slide, x, 2.05, 2.7, 1.45, self._RGBColor(250, 250, 250))
            self._rect(slide, x, 2.05, 2.7, 0.18, self._RGBColor(*color))
            growth = f"+{scenario.get('growth_low', 0)}-{scenario.get('growth_high', 0)}%"
            self._text(slide, x + 0.12, 2.40, 2.45, 0.2, label, size=9, bold=True, color=color)
            self._text(slide, x + 0.12, 2.70, 2.45, 0.3, growth, size=18, bold=True, color=color)
            self._text(slide, x + 0.12, 3.10, 2.45, 0.25, f"{format_currency(scenario.get('revenue_low'))} - {format_currency(scenario.get('revenue_high'))}", size=8.5, bold=True)
        # Brand-specific priority actions derived from real metrics (not boilerplate).
        actions = build_priority_actions(self.metrics, self.language)
        self._text(slide, 0.65, 3.72, 8.7, 0.2, self._t("projection_actions_title"), size=11, bold=True, color=(40, 40, 40))
        self._text(slide, 0.65, 4.00, 4.25, 1.1, "\n".join(f"-> {a}" for a in actions[0:3]), size=9, color=(70, 70, 70))
        if actions[3:6]:
            self._text(slide, 5.10, 4.00, 4.25, 1.1, "\n".join(f"-> {a}" for a in actions[3:6]), size=9, color=(70, 70, 70))
        # Honest framing: illustrative ranges, not a forecast.
        self._text(slide, 0.65, 5.18, 8.7, 0.22, self._t("projection_disclaimer"), size=7, color=(150, 150, 150))

    def _slide_conclusions(self, prs, page: int) -> None:
        slide = self._blank(prs)
        self._add_header(slide, page)
        self._title(slide, self._t("conclusions_title"), self._t("conclusions_subtitle"))
        sections = [
            (self._t("conclusions_current_situation"), self.narrative.get("conclusions", {}).get("current_situation", []), (212, 39, 45)),
            (self._t("conclusions_strengths"), self.narrative.get("conclusions", {}).get("strengths", []), (22, 163, 74)),
            (self._t("conclusions_plan"), self.narrative.get("conclusions", {}).get("plan", []), (29, 78, 216)),
            (self._t("conclusions_urgency"), self.narrative.get("conclusions", {}).get("urgency", []), (234, 88, 12)),
        ]
        for idx, (title, bullets, bar_color) in enumerate(sections):
            x = 0.65 + (idx % 2) * 4.35
            y = 1.25 + (idx // 2) * 1.65
            # Left coloured marker bar
            self._rect(slide, x, y, 0.12, 1.25, self._RGBColor(*bar_color))
            self._body_box(
                slide, x + 0.12, y, 3.83, 1.25,
                f"{title}\n" + "\n".join(f"- {bullet}" for bullet in bullets[:4]),
                title_bold=True,
            )

    def _title(self, slide, title: str, subtitle: str) -> None:
        self._text(slide, 0.65, 0.62, 4.5, 0.32, title, size=22, bold=True, color=(20, 20, 20))
        self._text(slide, 0.67, 1.0, 5.5, 0.18, subtitle, size=9, color=(90, 90, 90))

    def _kpi(self, slide, x: float, y: float, w: float, h: float, label: str, value: str) -> None:
        self._rect(slide, x, y, w, h, self._RGBColor(248, 248, 248), line=(230, 230, 230))
        self._text(slide, x + 0.12, y + 0.12, w - 0.24, 0.18, label, size=7.5, bold=True, color=(90, 90, 90))
        self._text(slide, x + 0.12, y + 0.38, w - 0.24, 0.28, value, size=15, bold=True, color=(30, 30, 30))

    def _body_box(self, slide, x: float, y: float, w: float, h: float, text: str, title_bold: bool = False) -> None:
        self._rect(slide, x, y, w, h, self._RGBColor(250, 250, 250), line=(230, 230, 230))
        shape = self._text(slide, x + 0.18, y + 0.14, w - 0.36, h - 0.2, text, size=8.8, color=(55, 55, 55))
        if title_bold and shape.text_frame.paragraphs:
            shape.text_frame.paragraphs[0].runs[0].font.bold = True
            shape.text_frame.paragraphs[0].runs[0].font.size = self._Pt(11)

    def _table(
        self,
        slide,
        x: float,
        y: float,
        w: float,
        h: float,
        headers: list[str],
        rows: list[list[str]],
        widths: list[float],
        *,
        yoy_columns: Optional[list[int]] = None,
    ) -> None:
        """Render a table. If ``yoy_columns`` is provided, cells in those columns
        are coloured green/red based on the sign of their numeric value (the
        ``+``/``-`` prefix or a leading minus sign is the cue).
        """
        table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), self._Inches(x), self._Inches(y), self._Inches(w), self._Inches(h))
        table = table_shape.table
        for idx, width in enumerate(widths):
            table.columns[idx].width = self._Inches(width)
        for col, header in enumerate(headers):
            cell = table.cell(0, col)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = self._RGBColor(245, 245, 245)
            self._cell_font(cell, bold=True, size=7.5)
        for row_idx, row in enumerate(rows, start=1):
            for col_idx, value in enumerate(row):
                cell = table.cell(row_idx, col_idx)
                cell.text = str(value)
                colour = None
                if yoy_columns and col_idx in yoy_columns:
                    colour = _yoy_cell_color(value)
                self._cell_font(cell, size=7.2, color=colour)

    def _cell_font(
        self,
        cell,
        *,
        bold: bool = False,
        size: float = 8,
        color: Optional[tuple[int, int, int]] = None,
    ) -> None:
        cell.margin_left = self._Inches(0.04)
        cell.margin_right = self._Inches(0.04)
        cell.margin_top = self._Inches(0.02)
        cell.margin_bottom = self._Inches(0.02)
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = "Nunito"
                run.font.size = self._Pt(size)
                run.font.bold = bold
                if color is not None:
                    run.font.color.rgb = self._RGBColor(*color)

    def _rect(self, slide, x: float, y: float, w: float, h: float, fill, line: Optional[tuple[int, int, int]] = None) -> None:
        from pptx.enum.shapes import MSO_SHAPE

        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, self._Inches(x), self._Inches(y), self._Inches(w), self._Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        if line is None:
            shape.line.fill.background()
        else:
            shape.line.color.rgb = self._RGBColor(*line)

    def _text(
        self,
        slide,
        x: float,
        y: float,
        w: float,
        h: float,
        text: str,
        *,
        size: float,
        bold: bool = False,
        color: tuple[int, int, int] = (20, 20, 20),
        align: str = "left",
    ):
        shape = slide.shapes.add_textbox(self._Inches(x), self._Inches(y), self._Inches(w), self._Inches(h))
        frame = shape.text_frame
        frame.clear()
        frame.margin_left = 0
        frame.margin_right = 0
        frame.margin_top = 0
        frame.margin_bottom = 0
        frame.word_wrap = True
        paragraph = frame.paragraphs[0]
        paragraph.alignment = {"center": self._PP_ALIGN.CENTER, "right": self._PP_ALIGN.RIGHT}.get(align, self._PP_ALIGN.LEFT)
        run = paragraph.add_run()
        run.text = str(text or "")
        run.font.name = "Nunito"
        run.font.size = self._Pt(size)
        run.font.bold = bold
        run.font.color.rgb = self._RGBColor(*color)
        return shape


class BrandAnalysisService:
    """Create, update and process Brand Analysis jobs."""

    def __init__(self, db: AsyncSession):
        self.db = db

    async def create_job(self, *, org_id: UUID, user_id: UUID, data) -> BrandAnalysisJob:
        account_id = UUID(data.account_id) if data.account_id else None
        if account_id is not None:
            result = await self.db.execute(
                select(AmazonAccount.id).where(
                    AmazonAccount.id == account_id,
                    AmazonAccount.organization_id == org_id,
                )
            )
            if not result.scalar_one_or_none():
                raise ValueError("Account not found or does not belong to organization")

        if data.market_type == "asin" and not data.asin_list:
            raise ValueError("asin_list is required when market_type is asin")
        if data.market_type == "brand" and not (data.market_query or data.brand_name):
            raise ValueError("market_query or brand_name is required when market_type is brand")

        job = BrandAnalysisJob(
            organization_id=org_id,
            created_by_id=user_id,
            account_id=account_id,
            brand_name=data.brand_name.strip(),
            language=data.language,
            mode=_canonical_mode(data.mode),
            market_type=data.market_type,
            market_query=(data.market_query or data.brand_name).strip(),
            asin_list=data.asin_list or None,
            status="pending",
            progress_step="Waiting for source data",
            progress_pct=0,
        )
        self.db.add(job)
        await self.db.flush()
        await self.db.refresh(job)
        return job

    async def get_job(self, job_id: UUID, org_id: UUID) -> Optional[BrandAnalysisJob]:
        result = await self.db.execute(
            select(BrandAnalysisJob)
            .options(selectinload(BrandAnalysisJob.source_files))
            .where(
                BrandAnalysisJob.id == job_id,
                BrandAnalysisJob.organization_id == org_id,
            )
        )
        return result.scalar_one_or_none()

    async def list_jobs(self, org_id: UUID, limit: int = 50, offset: int = 0) -> list[BrandAnalysisJob]:
        result = await self.db.execute(
            select(BrandAnalysisJob)
            .options(selectinload(BrandAnalysisJob.source_files))
            .where(BrandAnalysisJob.organization_id == org_id)
            .order_by(BrandAnalysisJob.created_at.desc())
            .offset(offset)
            .limit(limit)
        )
        return list(result.scalars().all())

    async def request_cancel(self, job_id: UUID, org_id: UUID) -> Optional[BrandAnalysisJob]:
        """Cooperatively cancel a job.

        A job that has not started running yet (``pending`` or already in a
        non-running, non-terminal state) is moved straight to ``cancelled``.
        A running job is flagged ``cancel_requested`` and parked in
        ``cancelling`` so the processor aborts at its next phase boundary; the
        Celery task is also revoked best-effort. Terminal jobs are rejected by
        the caller (409) before reaching here.
        """
        job = await self.get_job(job_id, org_id)
        if not job:
            return None
        if job.status in TERMINAL_STATUSES:
            return job

        job.cancel_requested = True
        if job.celery_task_id:
            _revoke_celery_task(job.celery_task_id)

        if job.status in RUNNING_STATUSES:
            job.status = "cancelling"
            job.progress_step = "Cancellation requested"
            job.progress_pct = STATUS_PROGRESS["cancelling"]
        else:
            job.status = "cancelled"
            job.progress_step = "Cancelled by user"
            job.progress_pct = 100
            job.error_message = None
            job.completed_at = datetime.utcnow()
        job.updated_at = datetime.utcnow()
        await self.db.flush()
        return job

    async def delete_job(self, job_id: UUID, org_id: UUID) -> bool:
        job = await self.get_job(job_id, org_id)
        if not job:
            return False
        # A running job must be cancelled first: deleting the row out from under
        # the worker leaves orphaned artifacts and races the final commit.
        if job.status in RUNNING_STATUSES:
            raise BrandAnalysisJobRunningError(
                "Cancel the running job before deleting it."
            )
        await self.db.delete(job)
        await self.db.flush()
        return True

    async def save_source_file(
        self,
        *,
        job: BrandAnalysisJob,
        year: int,
        filename: str,
        content_type: Optional[str],
        data: bytes,
        uploaded_by_id: Optional[UUID],
    ) -> BrandAnalysisSourceFile:
        if year not in {2024, 2025}:
            raise ValueError("Only 2024 and 2025 exports are supported")
        parsed = parse_brand_export(data, filename, year=year)
        await self.db.execute(
            delete(BrandAnalysisSourceFile).where(
                BrandAnalysisSourceFile.job_id == job.id,
                BrandAnalysisSourceFile.year == year,
            )
        )
        source = BrandAnalysisSourceFile(
            job_id=job.id,
            organization_id=job.organization_id,
            uploaded_by_id=uploaded_by_id,
            year=year,
            filename=filename,
            content_type=content_type,
            file_size=len(data),
            file_data=data,
            row_count=parsed.row_count,
            columns=parsed.columns,
            column_validation=parsed.validation.to_dict() if parsed.validation else None,
        )
        self.db.add(source)
        job.status = "pending"
        job.progress_step = "Waiting for processing"
        job.progress_pct = 0
        job.error_message = None
        job.error_code = None
        await self.db.flush()
        await self.db.refresh(source)
        return source


def _safe_filename_part(value: str) -> str:
    safe = re.sub(r"[^A-Za-z0-9_-]+", "_", value.strip())[:60].strip("_")
    return safe or "brand"


def _yoy_cell_color(value: Any) -> Optional[tuple[int, int, int]]:
    """Return a green/red RGB tuple for a YoY-style cell value, or None for neutral.

    Accepts strings like ``"+12.3%"``, ``"-5.0%"``, ``"–0.2%"``, ``"N/A"`` and
    raw numbers. Returns ``None`` when the sign cannot be determined.
    """
    if value is None:
        return None
    text = str(value).strip()
    if not text or text.upper() in {"N/A", "NAN", "-", "--"}:
        return None
    # Normalize en-dash / minus-sign characters to ASCII minus.
    normalized = text.replace("–", "-").replace("−", "-")
    if normalized.startswith("+"):
        return (22, 163, 74)
    if normalized.startswith("-") and len(normalized) > 1 and normalized[1].isdigit():
        return (212, 39, 45)
    # Try to parse numeric content.
    cleaned = re.sub(r"[^0-9.\-]+", "", normalized)
    try:
        number = float(cleaned)
    except ValueError:
        return None
    if number > 0:
        return (22, 163, 74)
    if number < 0:
        return (212, 39, 45)
    return None


def build_brand_analysis_pptx(metrics: dict[str, Any], narrative: dict[str, Any], language: str = "en") -> bytes:
    return BrandAnalysisPptxBuilder(metrics, narrative, language).build()


def validate_pptx_bytes(pptx_bytes: bytes) -> dict[str, Any]:
    """Open a generated PPTX and return a structural fingerprint for tests.

    Raises a :class:`ValueError` if the deck doesn't open or doesn't meet
    the expected 12-16 slide structure (N/A slides are skipped, so the count
    varies with the available data).
    """
    from pptx import Presentation

    prs = Presentation(io.BytesIO(pptx_bytes))
    slide_count = len(prs.slides)
    if slide_count < 12 or slide_count > 16:
        raise ValueError(f"Expected 12-16 slides, got {slide_count}")

    def slide_text(slide) -> str:
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
            elif shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        parts.append(cell.text_frame.text)
        return "\n".join(parts)

    texts = [slide_text(slide) for slide in prs.slides]
    return {
        "slide_count": slide_count,
        "slide_texts": texts,
    }


# Alert types surfaced to the NotificationBell on terminal brand-analysis
# transitions. The alerts API exposes ``alert_type`` from the owning rule, so
# each org gets one auto-created rule per type the first time a job finishes.
BRAND_ANALYSIS_READY_ALERT_TYPE = "brand_analysis_ready"
BRAND_ANALYSIS_FAILED_ALERT_TYPE = "brand_analysis_failed"


async def _ensure_brand_analysis_alert_rule(db: "AsyncSession", organization_id: UUID, alert_type: str):
    """Return (creating one if needed) the org's auto rule for a job alert type.

    Job-completion alerts are not user-configured; they hang off a hidden
    per-org rule so they reuse the existing alerts join/serialization and
    render in the NotificationBell without bespoke API work.
    """
    from app.models.alert import AlertRule

    result = await db.execute(
        select(AlertRule).where(
            AlertRule.organization_id == organization_id,
            AlertRule.alert_type == alert_type,
        )
    )
    rule = result.scalars().first()
    if rule:
        return rule

    rule = AlertRule(
        organization_id=organization_id,
        name="Brand Analysis",
        alert_type=alert_type,
        conditions={"auto_created": True},
        applies_to_accounts=None,
        applies_to_asins=None,
        notification_channels=[],
        notification_emails=None,
        webhook_url=None,
        is_enabled=True,
    )
    db.add(rule)
    await db.flush()
    return rule


def _emit_terminal_notification_factory(session_factory):
    """Build a coroutine that records a terminal-transition Alert.

    Bound to a session factory so it can run on the worker's private engine
    after the job's own commit. Best-effort: notification failures never fail
    the job.
    """

    async def _emit(*, organization_id, account_id, brand_name: str, status: str) -> None:
        from app.models.alert import Alert

        if status in {"completed", "completed_with_limitations"}:
            alert_type = BRAND_ANALYSIS_READY_ALERT_TYPE
            severity = "info"
            message = f"Brand analysis for {brand_name} is ready."
        elif status == "failed":
            alert_type = BRAND_ANALYSIS_FAILED_ALERT_TYPE
            severity = "critical"
            message = f"Brand analysis for {brand_name} failed."
        else:
            return

        dedup_key = f"{alert_type}:{organization_id}:{brand_name}:{status}"
        now = datetime.utcnow()
        try:
            async with session_factory() as ndb:
                rule = await _ensure_brand_analysis_alert_rule(ndb, organization_id, alert_type)
                existing = await ndb.execute(
                    select(Alert).where(
                        Alert.rule_id == rule.id,
                        Alert.dedup_key == dedup_key,
                        Alert.resolved_at.is_(None),
                    )
                )
                alert = existing.scalars().first()
                if alert is not None:
                    alert.message = message
                    alert.severity = severity
                    alert.is_read = False
                    alert.last_seen_at = now
                else:
                    ndb.add(
                        Alert(
                            rule_id=rule.id,
                            organization_id=organization_id,
                            account_id=account_id,
                            asin=None,
                            event_kind=alert_type,
                            dedup_key=dedup_key,
                            message=message,
                            details={"brand_name": brand_name, "status": status},
                            severity=severity,
                            is_read=False,
                            triggered_at=now,
                            last_seen_at=now,
                            notification_status="pending",
                        )
                    )
                rule.last_triggered_at = now
                await ndb.commit()
        except Exception as exc:
            logger.warning("Brand analysis notification emit failed for %s: %s", organization_id, exc)

    return _emit


def process_brand_analysis_job(job_id: str) -> None:
    """Process a brand analysis job in a background worker or thread."""
    from app.db.session import db_url as _db_url
    from app.services.brand_analysis_sources import (
        AmazonAccountDataSource,
        ManualUploadDataSource,
    )
    from app.services.brand_analysis_capabilities import detect_brand_analysis_capabilities
    from app.services.brand_analysis_storage import BrandAnalysisStorage

    _local_engine = create_async_engine(
        _db_url,
        echo=settings.APP_DEBUG,
        pool_size=2,
        max_overflow=1,
    )
    _LocalSession = async_sessionmaker(
        bind=_local_engine,
        class_=AsyncSession,
        expire_on_commit=False,
        autoflush=False,
    )

    _emit_terminal_notification = _emit_terminal_notification_factory(_LocalSession)

    async def _set_status(status_: str, step: str, pct: Optional[int] = None) -> None:
        try:
            async with _LocalSession() as pdb:
                await pdb.execute(
                    sa_text(
                        "UPDATE brand_analysis_jobs "
                        "SET status = :status, progress_step = :step, progress_pct = :pct, "
                        "updated_at = :updated_at, heartbeat_at = :updated_at "
                        "WHERE id = :rid"
                    ),
                    {
                        "status": status_,
                        "step": step,
                        "pct": STATUS_PROGRESS.get(status_, 0) if pct is None else pct,
                        "updated_at": datetime.utcnow(),
                        "rid": job_id,
                    },
                )
                await pdb.commit()
        except Exception as exc:
            logger.warning("Brand analysis progress update failed for %s: %s", job_id, exc)

    async def _job_state() -> Optional[tuple[bool, bool]]:
        """Return (still_exists, cancel_requested) for the job, on a fresh session.

        Used at phase boundaries so a delete or a cancel that lands mid-run is
        observed cooperatively instead of being silently overwritten by the
        worker's final commit.
        """
        try:
            async with _LocalSession() as pdb:
                row = (
                    await pdb.execute(
                        sa_text("SELECT cancel_requested FROM brand_analysis_jobs WHERE id = :rid"),
                        {"rid": job_id},
                    )
                ).first()
            if row is None:
                return None
            return True, bool(row[0])
        except Exception as exc:
            logger.warning("Brand analysis state probe failed for %s: %s", job_id, exc)
            return True, False

    async def _finalize_cancel() -> None:
        try:
            async with _LocalSession() as pdb:
                await pdb.execute(
                    sa_text(
                        "UPDATE brand_analysis_jobs "
                        "SET status = 'cancelled', progress_step = 'Cancelled by user', "
                        "progress_pct = 100, error_message = NULL, completed_at = :now, "
                        "updated_at = :now, heartbeat_at = :now "
                        "WHERE id = :rid AND status NOT IN ('completed', 'completed_with_limitations', 'failed', 'cancelled')"
                    ),
                    {"now": datetime.utcnow(), "rid": job_id},
                )
                await pdb.commit()
        except Exception as exc:
            logger.warning("Brand analysis cancel finalize failed for %s: %s", job_id, exc)

    async def _cancel_or_deleted() -> bool:
        """True if the job has been deleted or asked to cancel.

        On cancel, finalizes the row to ``cancelled`` so the caller can abort
        cleanly without committing a half-built artifact.
        """
        state = await _job_state()
        if state is None:
            logger.info("Brand analysis job %s deleted mid-run; aborting", job_id)
            return True
        _, cancel_requested = state
        if cancel_requested:
            logger.info("Brand analysis job %s cancellation requested; aborting", job_id)
            await _finalize_cancel()
            return True
        return False

    def _resolve_adapter(job: BrandAnalysisJob, source_by_year: dict) -> Any:
        """Pick the right adapter for the job.

        Brand Analysis is autonomous: ``internal`` (the canonical
        Inthezon SP-API + Market Research path) is preferred. ``manual``
        uses uploaded yearly product exports. Deprecated external-provider
        modes always fall back to manual upload; they are no longer offered
        in the UI but old DB rows still resolve cleanly.
        """
        mode = _canonical_mode(job.mode)

        has_complete_manual_fallback = {2024, 2025}.issubset(source_by_year)
        if mode == "manual" or has_complete_manual_fallback:
            return ManualUploadDataSource(
                source_files={
                    year: (source.file_data, source.filename)
                    for year, source in source_by_year.items()
                }
            )

        if mode == "internal":
            if not job.account_id:
                raise BrandAnalysisDataError(
                    "Internal data mode requires a connected Amazon account. "
                    "Pick an account or switch to upload of external yearly exports."
                )
            return AmazonAccountDataSource(
                db=None,  # set per request, see fetch loop
                account_id=job.account_id,
                organization_id=job.organization_id,
                brand_filter=(job.market_query or job.brand_name) if job.market_type == "brand" else None,
                asin_list=job.asin_list,
            )

        # Legacy / unknown modes route through manual upload, the only
        # fallback data source. The adapter raises InsufficientDataError
        # per year when a file is missing so the processor can surface a
        # clean ``waiting_for_user_action`` state.
        return ManualUploadDataSource(
            source_files={
                year: (source.file_data, source.filename)
                for year, source in source_by_year.items()
            }
        )

    async def _process() -> None:
        async with _LocalSession() as db:
            result = await db.execute(
                select(BrandAnalysisJob)
                .options(selectinload(BrandAnalysisJob.source_files))
                .where(BrandAnalysisJob.id == UUID(job_id))
            )
            job = result.scalar_one_or_none()
            if not job:
                logger.error("Brand analysis job %s not found", job_id)
                return

            try:
                source_by_year = {source.year: source for source in job.source_files}
                capability_matrix: dict[str, Any] = job.capability_matrix or {}
                data_coverage: dict[str, Any] = job.data_coverage or {}

                if _canonical_mode(job.mode) == "internal":
                    account_result = await db.execute(
                        select(AmazonAccount).where(
                            AmazonAccount.id == job.account_id,
                            AmazonAccount.organization_id == job.organization_id,
                        )
                    )
                    account = account_result.scalar_one_or_none()
                    if not account:
                        raise BrandAnalysisDataError("Internal data mode requires a connected Amazon account.")

                    await _set_status("capability_checking", "Checking Amazon/SP-API capabilities", 8)
                    capability_result = await detect_brand_analysis_capabilities(
                        db,
                        account,
                        force_refresh=True,
                    )
                    capability_matrix = capability_result.to_dict()
                    job.capability_matrix = capability_matrix
                    await db.commit()

                    await _set_status("preflight_checking", "Checking internal sales coverage", 14)
                    data_coverage = await inspect_internal_sales_data_coverage(db, account)
                    job.data_coverage = data_coverage
                    await db.commit()

                    if data_coverage.get("needs_sync") and job.sync_attempt_count < settings.BRAND_ANALYSIS_MAX_SYNC_ATTEMPTS:
                        sync_windows = data_coverage.get("sync_windows") or {}
                        idempotency_key = f"{job.id}:{json.dumps(sync_windows, sort_keys=True)}"
                        if job.sync_idempotency_key != idempotency_key:
                            job.sync_attempt_count = int(job.sync_attempt_count or 0) + 1
                            job.sync_idempotency_key = idempotency_key
                            job.status = "internal_sync_requested"
                            job.progress_step = "Internal sales sync requested"
                            job.progress_pct = STATUS_PROGRESS["internal_sync_requested"]
                            job.updated_at = datetime.utcnow()
                            await db.commit()

                            await _set_status("syncing_internal_data", "Syncing recoverable Amazon sales windows", 28)
                            try:
                                from app.models.amazon_account import AccountType
                                from app.services.data_extraction import DataExtractionService

                                extraction = DataExtractionService(db)
                                organization = await extraction._load_organization(account)
                                synced_records = 0
                                api_limitations: list[dict[str, str]] = []
                                for year, window in sync_windows.items():
                                    if not window:
                                        continue
                                    start_date = date.fromisoformat(window["start_date"])
                                    end_date = date.fromisoformat(window["end_date"])
                                    try:
                                        if account.account_type == AccountType.VENDOR:
                                            synced_records += await extraction.sync_vendor_sales_data(
                                                account,
                                                organization,
                                                start_date,
                                                end_date,
                                            )
                                        else:
                                            synced_records += await extraction.sync_sales_data(
                                                account,
                                                organization,
                                                start_date,
                                                end_date,
                                            )
                                    except Exception as sync_window_exc:
                                        api_limitations.append(
                                            {
                                                "year": str(year),
                                                "start_date": window["start_date"],
                                                "end_date": window["end_date"],
                                                "error": str(sync_window_exc)[:500],
                                            }
                                        )
                                await db.flush()
                                data_coverage = await inspect_internal_sales_data_coverage(db, account)
                                data_coverage["sync_result"] = {
                                    "attempt_count": job.sync_attempt_count,
                                    "records": synced_records,
                                    "api_limitations": api_limitations,
                                }
                                job.data_coverage = data_coverage
                                job.last_sync_error = None if not api_limitations else json.dumps(api_limitations)
                                job.status = "internal_sync_completed" if not api_limitations else "internal_sync_failed"
                                job.progress_step = "Internal sync completed" if not api_limitations else "Internal sync completed with API limitations"
                                job.progress_pct = STATUS_PROGRESS[job.status]
                                job.updated_at = datetime.utcnow()
                                await db.commit()
                            except Exception as sync_exc:
                                logger.exception("Automatic Brand Analysis sales sync failed for %s", job_id)
                                await db.rollback()
                                fail_result = await db.execute(select(BrandAnalysisJob).where(BrandAnalysisJob.id == UUID(job_id)))
                                job = fail_result.scalar_one()
                                job.status = "internal_sync_failed"
                                job.progress_step = "Internal sync failed; continuing with available data"
                                job.progress_pct = STATUS_PROGRESS["internal_sync_failed"]
                                job.last_sync_error = str(sync_exc)[:1000]
                                job.next_retry_at = datetime.utcnow() + timedelta(hours=1)
                                data_coverage = dict(data_coverage)
                                data_coverage["sync_result"] = {
                                    "attempt_count": job.sync_attempt_count,
                                    "error": job.last_sync_error,
                                }
                                job.data_coverage = data_coverage
                                await db.commit()

                if await _cancel_or_deleted():
                    return

                adapter = _resolve_adapter(job, source_by_year)
                if isinstance(adapter, AmazonAccountDataSource):
                    adapter.db = db
                source_name = getattr(adapter, "source_name", "unknown")

                await _set_status("collecting_source_data", "Resolving ASINs and loading 2024 performance", 25)
                try:
                    parsed_2024 = await adapter.fetch_year(2024)
                    await _set_status("collecting_source_data", "Loading 2025 yearly performance", 40)
                    parsed_2025 = await adapter.fetch_year(2025)
                except InsufficientDataError as recovery_error:
                    logger.info(
                        "Brand analysis job %s needs manual fallback data: %s",
                        job_id, recovery_error,
                    )
                    error_code = "manual_upload_required"
                    if isinstance(recovery_error, InsufficientDataError):
                        if recovery_error.year == 2024:
                            error_code = "missing_2024_data"
                        elif recovery_error.year == 2025:
                            error_code = "missing_2025_data"
                        elif source_name == "internal":
                            error_code = "internal_data_missing"
                        else:
                            error_code = "insufficient_yearly_data"
                    job.status = "waiting_for_user_action"
                    job.progress_step = "Waiting for external yearly export upload"
                    job.progress_pct = STATUS_PROGRESS["waiting_for_user_action"]
                    job.error_message = str(recovery_error)[:1000]
                    job.error_code = error_code
                    job.data_source_name = source_name
                    job.updated_at = datetime.utcnow()
                    await db.commit()
                    return

                # Persist column validation reports for manual sources.
                if parsed_2024.validation and 2024 in source_by_year:
                    source_by_year[2024].column_validation = parsed_2024.validation.to_dict()
                if parsed_2025.validation and 2025 in source_by_year:
                    source_by_year[2025].column_validation = parsed_2025.validation.to_dict()

                if parsed_2024.source_name != "manual_upload":
                    await _set_status("enriching_catalog", "Enriching catalog from Market Research")

                # Track partial catalog enrichment so the UI can flag a
                # completed analysis as having missing optional fields.
                partial_enrichment_code: Optional[str] = None
                if isinstance(adapter, AmazonAccountDataSource) and adapter.enrichment_partial:
                    partial_enrichment_code = "catalog_enrichment_partial"

                if await _cancel_or_deleted():
                    return

                await _set_status("generating_metrics", "Calculating deterministic metrics")
                metrics = calculate_brand_metrics(parsed_2024, parsed_2025, brand_name=job.brand_name)
                completeness = assess_data_completeness(parsed_2024, parsed_2025)
                metrics["data_completeness"] = completeness
                metrics["capability_matrix"] = capability_matrix
                metrics["data_coverage"] = data_coverage
                if isinstance(adapter, AmazonAccountDataSource):
                    readiness = adapter.describe_readiness()
                    metrics["data_readiness"] = readiness
                    incomplete_years = [
                        int(year)
                        for year, report in (readiness.get("years") or {}).items()
                        if report and not report.get("complete_year", False)
                    ]
                    metrics["data_completeness"]["history_incomplete_years"] = incomplete_years
                completion_code = partial_enrichment_code
                if completion_code is None and not completeness["optional_fields_complete"]:
                    completion_code = "analysis_completed_with_missing_optional_fields"
                limitations = build_limitation_summary(metrics, capability_matrix, data_coverage)
                if not completeness["optional_fields_complete"]:
                    limitations["items"].append(
                        {
                            "area": "optional_catalog_fields",
                            "message": "Some optional catalog fields are unavailable and appear as N/A.",
                        }
                    )
                    limitations["has_limitations"] = True
                metrics["limitations"] = limitations
                metrics["metric_source_registry"] = build_metric_source_registry(metrics, source_name)
                provenance = enrich_metric_provenance(
                    build_metric_provenance(parsed_2024, parsed_2025),
                    metrics,
                    source_name,
                )
                validate_metric_provenance_for_deck(metrics, provenance)

                if await _cancel_or_deleted():
                    return

                await _set_status("generating_narrative", "Generating strategic narrative", 82)
                narrative = BrandAnalysisNarrativeService().generate(
                    metrics,
                    job.language,
                    provenance=provenance,
                    limitations=limitations,
                )

                if await _cancel_or_deleted():
                    return

                await _set_status("generating_pptx", "Generating PowerPoint deck")
                pptx_bytes = build_brand_analysis_pptx(metrics, narrative, language=job.language)
                # Open the deck back up to confirm the file is structurally
                # valid (16 slides, correct OOXML). If it isn't, surface as a
                # job failure rather than letting the user download a broken
                # artifact.
                try:
                    pptx_fingerprint = validate_pptx_bytes(pptx_bytes)
                except Exception as pptx_exc:
                    raise BrandAnalysisDataError(
                        f"Generated Brand Analysis deck failed validation: {pptx_exc}"
                    ) from pptx_exc
                metrics["pptx_fingerprint"] = {
                    "slide_count": pptx_fingerprint.get("slide_count"),
                }
                filename = f"brand_analysis_{_safe_filename_part(job.brand_name)}_{datetime.utcnow().date().isoformat()}.pptx"

                storage = BrandAnalysisStorage()
                artifact_ref = storage.save_artifact(
                    organization_id=job.organization_id,
                    job_id=job.id,
                    filename=filename,
                    content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    data=pptx_bytes,
                )

                # Re-select on the work session before the final commit so a
                # delete or cancel that landed mid-run cannot be resurrected by
                # writing the completed artifact back onto a vanished/cancelled row.
                guard = await db.execute(
                    select(BrandAnalysisJob).where(BrandAnalysisJob.id == UUID(job_id)).with_for_update()
                )
                fresh_job = guard.scalar_one_or_none()
                if fresh_job is None:
                    logger.info("Brand analysis job %s deleted before completion; discarding artifact", job_id)
                    await db.rollback()
                    return
                if fresh_job.cancel_requested or fresh_job.status in TERMINAL_STATUSES:
                    logger.info("Brand analysis job %s cancelled before completion; finalizing", job_id)
                    await db.rollback()
                    await _finalize_cancel()
                    return
                job = fresh_job

                job.metrics = metrics
                job.metric_provenance = provenance
                job.capability_matrix = capability_matrix
                job.data_coverage = data_coverage
                job.limitations = limitations
                job.data_source_name = source_name
                job.narrative = narrative
                job.artifact_data = pptx_bytes
                job.artifact_filename = filename
                job.artifact_content_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                job.storage_ref = artifact_ref.to_dict()
                has_limitations = limitations.get("has_limitations") or completion_code
                job.status = "completed_with_limitations" if has_limitations else "completed"
                job.progress_step = "Complete with limitations" if has_limitations else "Complete"
                job.progress_pct = 100
                job.error_message = None
                job.error_code = completion_code
                job.completed_at = datetime.utcnow()
                job.updated_at = datetime.utcnow()
                await db.commit()

                await _emit_terminal_notification(
                    organization_id=job.organization_id,
                    account_id=job.account_id,
                    brand_name=job.brand_name,
                    status=job.status,
                )
            except Exception as exc:
                logger.exception("Brand analysis job %s failed", job_id)
                await db.rollback()
                notify_args: Optional[dict] = None
                async with _LocalSession() as fdb:
                    fail_result = await fdb.execute(select(BrandAnalysisJob).where(BrandAnalysisJob.id == UUID(job_id)))
                    failed_job = fail_result.scalar_one_or_none()
                    if failed_job:
                        # A delete or cancel that landed mid-run wins over the failure.
                        if failed_job.cancel_requested:
                            failed_job.status = "cancelled"
                            failed_job.error_message = None
                            failed_job.progress_step = "Cancelled by user"
                        else:
                            failed_job.status = "failed"
                            failed_job.error_message = str(exc)[:1000]
                            failed_job.progress_step = "Failed"
                            notify_args = {
                                "organization_id": failed_job.organization_id,
                                "account_id": failed_job.account_id,
                                "brand_name": failed_job.brand_name,
                                "status": "failed",
                            }
                        failed_job.progress_pct = 100
                        failed_job.completed_at = datetime.utcnow()
                        failed_job.updated_at = datetime.utcnow()
                        failed_job.heartbeat_at = datetime.utcnow()
                        await fdb.commit()
                if notify_args is not None:
                    await _emit_terminal_notification(**notify_args)

    import asyncio

    loop = asyncio.new_event_loop()
    asyncio.set_event_loop(loop)
    try:
        loop.run_until_complete(_process())
    finally:
        loop.run_until_complete(_local_engine.dispose())
        loop.close()
