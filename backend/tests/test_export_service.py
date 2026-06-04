from datetime import date
from importlib.util import module_from_spec, spec_from_file_location
from pathlib import Path
from types import SimpleNamespace
import sys
import types
from uuid import uuid4

import pytest


MODULE_PATH = Path(__file__).resolve().parents[1] / "app" / "services" / "export_service.py"
services_pkg = types.ModuleType("app.services")
services_pkg.__path__ = [str(MODULE_PATH.parent)]
sys.modules.setdefault("app.services", services_pkg)

data_extraction_stub = types.ModuleType("app.services.data_extraction")
data_extraction_stub.DAILY_TOTAL_ASIN = "__DAILY_TOTAL__"
sys.modules["app.services.data_extraction"] = data_extraction_stub

SPEC = spec_from_file_location("export_service_under_test", MODULE_PATH)
MODULE = module_from_spec(SPEC)
assert SPEC is not None and SPEC.loader is not None
SPEC.loader.exec_module(MODULE)

ExportService = MODULE.ExportService
DAILY_TOTAL_ASIN = data_extraction_stub.DAILY_TOTAL_ASIN
_previous_period = MODULE._previous_period
_trend_period_buckets = MODULE._trend_period_buckets


def test_previous_period_uses_same_inclusive_length():
    prev_start, prev_end = _previous_period(date(2026, 3, 1), date(2026, 3, 31))

    assert prev_start == date(2026, 1, 29)
    assert prev_end == date(2026, 2, 28)


def test_csv_bytes_uses_bom_and_localized_headers():
    service = ExportService(None)  # type: ignore[arg-type]
    rows = [
        {
            "section": "Metadata",
            "metric": "Report Type",
            "current_value": "Sales Report",
            "previous_value": "",
            "change_percent": "",
            "unit": "",
            "notes": "",
        }
    ]

    payload = service._csv_bytes(
        rows,
        ["section", "metric", "current_value", "previous_value", "change_percent", "unit", "notes"],
        "it",
    )

    text = payload.decode("utf-8-sig")
    assert text.startswith("Sezione,Metrica,Valore Attuale")
    assert "Sales Report" in text


def test_inventory_metrics_respects_low_stock_filter():
    service = ExportService(None)  # type: ignore[arg-type]
    rows = [
        {"total_available": 0, "inbound_total": 5, "reserved_quantity": 1},
        {"total_available": 8, "inbound_total": 3, "reserved_quantity": 0},
        {"total_available": 40, "inbound_total": 2, "reserved_quantity": 4},
    ]

    metrics = service._inventory_metrics(rows, True)

    assert metrics["total_skus"] == 3
    assert metrics["exported_rows"] == 2
    assert metrics["low_stock_skus"] == 2
    assert metrics["total_available"] == 48


# --------------------------------------------------------------------------- #
# Revenue-trend zero-fill + sentinel fallback
# --------------------------------------------------------------------------- #
class _FakeResult:
    def __init__(self, rows):
        self._rows = rows

    def all(self):
        return self._rows


class _FakeAsyncSession:
    """Returns queued query results in order, like the analytics-service tests."""

    def __init__(self, responses):
        self._responses = list(responses)

    async def execute(self, query):
        return _FakeResult(self._responses.pop(0) if self._responses else [])


def _trend_row(report_date, revenue, units=0, orders=0, currency="EUR"):
    return SimpleNamespace(
        report_date=report_date,
        revenue=revenue,
        units=units,
        orders=orders,
        currency=currency,
    )


def test_trend_period_buckets_fills_every_month():
    buckets = _trend_period_buckets(date(2026, 1, 15), date(2026, 6, 10), "month")

    assert buckets == [
        date(2026, 1, 1),
        date(2026, 2, 1),
        date(2026, 3, 1),
        date(2026, 4, 1),
        date(2026, 5, 1),
        date(2026, 6, 1),
    ]


@pytest.mark.asyncio
async def test_sales_trend_rows_multi_month_returns_more_than_one_row():
    session = _FakeAsyncSession(
        [
            [
                _trend_row(date(2026, 1, 1), 1000.0, 10, 8),
                _trend_row(date(2026, 3, 1), 3000.0, 30, 24),
            ]
        ]
    )
    service = ExportService(session)  # type: ignore[arg-type]

    rows = await service._sales_trend_rows(
        [uuid4()], date(2026, 1, 1), date(2026, 3, 31), "month"
    )

    assert len(rows) == 3
    assert [r["report_date"] for r in rows] == [
        date(2026, 1, 1),
        date(2026, 2, 1),
        date(2026, 3, 1),
    ]


@pytest.mark.asyncio
async def test_sales_trend_rows_zero_fills_sparse_data():
    # Only February carries data; the chart must still show 6 categories.
    session = _FakeAsyncSession([[_trend_row(date(2026, 2, 1), 500.0, 5, 4)]])
    service = ExportService(session)  # type: ignore[arg-type]

    rows = await service._sales_trend_rows(
        [uuid4()], date(2026, 1, 1), date(2026, 6, 30), "month"
    )

    assert len(rows) == 6
    revenues = {r["report_date"]: r["revenue"] for r in rows}
    assert revenues[date(2026, 2, 1)] == 500.0
    assert revenues[date(2026, 1, 1)] == 0.0
    assert revenues[date(2026, 6, 1)] == 0.0


@pytest.mark.asyncio
async def test_sales_trend_rows_falls_back_to_per_asin_rows():
    # First query (sentinel rows) returns empty -> second query (per-ASIN) is used.
    session = _FakeAsyncSession(
        [
            [],
            [
                _trend_row(date(2026, 1, 1), 700.0, 7, 5),
                _trend_row(date(2026, 2, 1), 900.0, 9, 6),
            ],
        ]
    )
    service = ExportService(session)  # type: ignore[arg-type]

    rows = await service._sales_trend_rows(
        [uuid4()], date(2026, 1, 1), date(2026, 2, 28), "month"
    )

    assert len(rows) == 2
    assert rows[0]["revenue"] == 700.0
    assert rows[1]["revenue"] == 900.0


# --------------------------------------------------------------------------- #
# POST /exports/powerpoint — live deck assembly
# --------------------------------------------------------------------------- #
def _load_exports_module():
    """Load the live exports endpoint module against the real app package."""
    sys.modules.pop("app.services", None)
    sys.modules.pop("app.services.data_extraction", None)
    from app.api.v1 import exports  # noqa: WPS433

    return exports


class _FakeExportService:
    """Deterministic stand-in for ExportService used by the endpoint."""

    def __init__(self, db):
        self.db = db
        self.account = SimpleNamespace(id=uuid4(), account_name="Acme IT")

    async def _get_accounts(self, org_id, account_ids):
        return [self.account]

    async def _sales_summary(self, account_ids, start_date, end_date):
        return {
            "revenue": 12345.67,
            "units": 321,
            "orders": 210,
            "average_order_value": 58.79,
            "average_selling_price": 38.46,
            "units_per_order": 1.5,
            "active_asins": 12,
        }

    async def _sales_trend_rows(self, account_ids, start_date, end_date, group_by):
        return [
            {"report_date": bucket, "revenue": float(i * 100), "currency": "EUR"}
            for i, bucket in enumerate(_trend_period_buckets(start_date, end_date, group_by))
        ]

    async def _sales_product_rows(self, account_ids, start_date, end_date):
        return [
            {"asin": f"B0{i:08d}", "title": f"Product {i}", "units": i, "revenue": float(i * 10)}
            for i in range(1, 16)
        ]

    async def _latest_snapshot_date(self, account_ids):
        return None

    async def _inventory_snapshot_rows(self, account_ids, snapshot_date, lang):
        return []

    async def _advertising_rollup_rows(self, account_ids, start_date, end_date):
        return []


def _run_powerpoint_endpoint(exports, rec_service_factory):
    import asyncio
    import io
    from pptx import Presentation

    org = SimpleNamespace(id=uuid4())
    user = SimpleNamespace(id=uuid4())
    start_date = date(2026, 1, 1)
    end_date = date(2026, 6, 30)

    orig_service = exports.ExportService
    orig_rec = exports.StrategicRecommendationsService
    orig_latest = exports._latest_forecasts
    exports.ExportService = _FakeExportService
    exports.StrategicRecommendationsService = rec_service_factory

    async def _no_forecasts(db, account_ids):
        return []

    exports._latest_forecasts = _no_forecasts
    try:
        response = asyncio.run(
            exports.export_to_powerpoint(
                current_user=user,
                organization=org,
                db=None,
                start_date=start_date,
                end_date=end_date,
                account_ids=None,
                group_by=None,
                template="default",
                language="it",
            )
        )
        body = b"".join(asyncio.run(_drain(response.body_iterator)))
        prs = Presentation(io.BytesIO(body))
        return prs, start_date, end_date
    finally:
        exports.ExportService = orig_service
        exports.StrategicRecommendationsService = orig_rec
        exports._latest_forecasts = orig_latest


async def _drain(iterator):
    chunks = []
    async for chunk in iterator:
        chunks.append(chunk if isinstance(chunk, bytes) else chunk.encode())
    return chunks


class _OkRecService:
    def __init__(self, db):
        pass

    async def list_recommendations(self, org_id, *, limit=6):
        return [
            SimpleNamespace(category="pricing", title="Raise price on hero ASIN", expected_impact="+5% margin"),
        ]


class _BrokenRecService:
    """Simulates the Anthropic-backed service failing (e.g. credit/500 error)."""

    def __init__(self, db):
        pass

    async def list_recommendations(self, org_id, *, limit=6):
        raise RuntimeError("anthropic: insufficient credits (500)")


def test_powerpoint_endpoint_builds_valid_deck():
    exports = _load_exports_module()
    prs, start_date, end_date = _run_powerpoint_endpoint(exports, _OkRecService)

    # cover, exec summary, kpi, trend, 2x products, inventory, ads, forecast,
    # recommendations, agency -> well above the old 6-slide deck.
    assert len(prs.slides._sldIdLst) >= 10

    expected_periods = len(_trend_period_buckets(start_date, end_date, "week"))
    chart_categories = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_chart:
                chart_categories = list(shape.chart.plots[0].categories)
    assert len(chart_categories) == expected_periods
    assert expected_periods > 1


def test_powerpoint_endpoint_degrades_when_recommendations_raise():
    exports = _load_exports_module()
    # Must not raise even though the recommendation service blows up.
    prs, _, _ = _run_powerpoint_endpoint(exports, _BrokenRecService)
    assert len(prs.slides._sldIdLst) >= 10
